# add_gradient_shape.py
# -*- coding: utf-8 -*-
'''
add_gradient_shape: Add a shape with gradient fill to a PowerPoint slide.
    Required:
    - slide: Slide object to add the shape to
    - left/top/width/height: Position and size (inches or Inches/Cm/Pt objects)
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - color_start/color_end: Gradient colors in "#RRGGBB" format (default: "#FFD54F"/"#9E9E9E")
    - angle_deg: Gradient angle, 0=left→right, 90=bottom→top (default: 0)
    - method: "ooxml" (native, recommended) or "picture" (compatible) (default: "ooxml")
    - alpha_start/alpha_end: Opacity 0~1 (ooxml mode only)
    Returns: The added shape object

add_gradient_background: Apply a gradient background to an entire slide.
    Required:
    - slide: Slide object to modify
    Optional:
    - color_start/color_end: Gradient colors in "#RRGGBB" format (default: "#FFD54F"/"#9E9E9E")
    - angle_deg: Gradient angle, 0=left→right, 90=bottom→top (default: 90)
    - method: "ooxml" (native) or "picture" (compatible) (default: "ooxml")
    - alpha_start/alpha_end: Opacity 0~1 (ooxml mode only)
    - stops: Custom gradient stops [(pos,rgb,alpha),...] (overrides colors if set)
    Effect: Modifies slide background directly, no return value

add_solid_shape: Add a shape with solid color fill to a slide.
    Required:
    - slide: Slide object to add the shape to
    - left/top/width/height: Position and size (inches or Length objects)
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - fill_color: Fill color in "#RRGGBB" format (default: "#FFD54F")
    - fill_alpha: Opacity 0~1 (default: 1.0)
    - outline_color: Border color "#RRGGBB" or None for no border
    - outline_width_pt: Border width in points (None or 0 for no border)
    Returns: The added shape object

add_image_filled_shape: Add a shape filled with an image to a slide.
    Required:
    - slide: Slide object to add the shape to
    - image_path: Path to the image file
    - left/top/width/height: Position and size
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - mode: "stretch" (default) or "tile"
    - crop: Cropping (left,top,right,bottom) in 0~1 range
    - tile_scale: Scale factor for tile mode (default: 1.0)
    - remove_line: Whether to remove shape outline (default: True)
    Returns: The added shape object

add_textbox: Add a text box with / without border that auto-fits text.
    Required:
    - slide: Slide object to add to
    - text: Text content
    - left/top/width/height: Position and size (inches)
    - font_size: Target font size in points (will auto-shrink if needed)
    Optional:
    - font_name: Font name (default: "Calibri")
    - font_color: Text color (default: Black)
    - has_base_box: Whether to draw a base box (default: False)
    - border_color: Border color (default: Black)
    - border_width: Border thickness in points (default: 2)
    - corner_radius: Corner rounding 0~1 (default: 0.1)
    - has_border: Whether to show border (default: True)
    - fill_color: Background color or None for transparent
    - inner_padding: Text padding in inches (default: 0.10)
    Returns: Tuple of (border shape, text box shape)

add_line: Add a line shape to a slide.
    Required:
    - slide: Slide object to add to
    - x1/y1: Start point position in inches (inchs)
    - x2/y2: End point position in inches (inchs)
    Optional:
    - color: Line color (default: Black)
    - width_pt: Line thickness in points (default: 2)
    - dash_style: Line dash pattern keyword ("solid", "dashed", "dash_dot", "dash_dot_dot", "long_dash", "long_dash_dot", "dot", "square_dot", "mixed")
    Returns: Connector shape object representing the line
    
'''
from __future__ import annotations

from typing import Tuple, List
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR 
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import math
import re

# =========================
# 工具：单位与颜色
# =========================
def _to_length(v):
    """接受 pptx 长度对象或数字（英寸），统一转 pptx Length 对象"""
    if hasattr(v, "_EMU"):  # 已是 Length/Inches/Cm/Pt
        return v
    if isinstance(v, (int, float)):
        return Inches(v)
    raise TypeError("left/top/width/height 需为 pptx 长度对象（Inches/Cm/Pt）或数值（单位英寸）。")

def _hex_to_rgb(hexstr: str) -> Tuple[int, int, int]:
    s = hexstr.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# =========================
# 兼容：取得/创建 <p:spPr>（形状属性）
# =========================
def _get_or_add_spPr(sp_elem):
    """
    兼容不同 python-pptx 版本：返回 <p:spPr>，若不存在则按规范顺序插入。
    sp_elem = shape._element (CT_Shape/CT_AutoShape/…)
    """
    spPr = sp_elem.find(qn("p:spPr"))
    if spPr is None:
        spPr = OxmlElement("p:spPr")
        nvSpPr = sp_elem.find(qn("p:nvSpPr"))
        if nvSpPr is not None:
            sp_elem.insert(list(sp_elem).index(nvSpPr) + 1, spPr)
        else:
            sp_elem.insert(0, spPr)
    return spPr

def _remove_fills(spPr):
    """清掉已有填充节点，避免冲突"""
    for tag in ("solidFill", "pattFill", "gradFill", "noFill", "blipFill"):
        el = spPr.find(qn(f"a:{tag}"))
        if el is not None:
            spPr.remove(el)


# =========================
# 兼容：取得/创建 <p:bg> / <p:bgPr>（背景）
# =========================
def _get_or_add_bg(slide_elem):
    """
    返回 <p:bg>，若不存在则在 <p:cSld> 内按规范顺序插入：
    应位于 <p:spTree> 之前。
    """
    bg = slide_elem.find(qn("p:bg"))
    if bg is not None:
        return bg

    cSld = slide_elem.find(qn("p:cSld"))
    if cSld is None:
        raise RuntimeError("Malformed slide: missing <p:cSld>")

    bg = OxmlElement("p:bg")

    spTree = cSld.find(qn("p:spTree"))
    if spTree is not None:
        # 插到 spTree 之前（符合 schema: bg? spTree ...）
        idx = list(cSld).index(spTree)
        cSld.insert(idx, bg)
    else:
        # 没有 spTree 时，插在最前
        cSld.insert(0, bg)

    return bg

def _get_or_add_bgPr(bg_elem):
    """
    返回/创建 <p:bgPr>；若存在 <p:bgRef> 先移除以避免冲突。
    """
    bgRef = bg_elem.find(qn("p:bgRef"))
    if bgRef is not None:
        bg_elem.remove(bgRef)

    bgPr = bg_elem.find(qn("p:bgPr"))
    if bgPr is None:
        bgPr = OxmlElement("p:bgPr")
        bg_elem.append(bgPr)
    return bgPr


# =========================
# 原生梯度：写 a:gradFill
# =========================
_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

def _add_grad_fill(spPr_or_bgPr, stops: List[Tuple[int, Tuple[int, int, int], int]], angle_deg=0):
    """
    在 <p:spPr> 或 <p:bgPr> 下写 a:gradFill。
    stops: [(pos(0..100000), (r,g,b), alpha(0..100000))]
    angle_deg: 0=左→右，90=下→上（已做坐标系修正）
    """
    # 清理已有填充
    for tag in ("bgRef", "blipFill", "gradFill", "grpFill", "noFill", "solidFill", "pattFill"):
        el = spPr_or_bgPr.find(qn(f"a:{tag}"))
        if el is not None:
            spPr_or_bgPr.remove(el)

    grad = etree.SubElement(spPr_or_bgPr, etree.QName(_NS["a"], "gradFill"))
    gsLst = etree.SubElement(grad, etree.QName(_NS["a"], "gsLst"))
    for pos, (r, g, b), alpha in stops:
        gs = etree.SubElement(gsLst, etree.QName(_NS["a"], "gs"))
        gs.set("pos", str(int(pos)))
        srgb = etree.SubElement(gs, etree.QName(_NS["a"], "srgbClr"))
        srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
        if alpha is not None:
            a = etree.SubElement(srgb, etree.QName(_NS["a"], "alpha"))
            a.set("val", str(int(alpha)))
    lin = etree.SubElement(grad, etree.QName(_NS["a"], "lin"))
    lin.set("ang", str(int(((90 - angle_deg) % 360) * 60000)))  # 0=左→右，90=下→上
    lin.set("scaled", "1")

def _set_shape_linear_gradient(shape,
                               color_start: str, color_end: str,
                               angle_deg=0, alpha_start=1.0, alpha_end=1.0):
    """给可填充 autoshape/textbox 设置原生线性渐变"""
    a1 = int(max(0, min(1, alpha_start)) * 100000)
    a2 = int(max(0, min(1, alpha_end)) * 100000)
    stops = [
        (0,      _hex_to_rgb(color_start), a1),
        (100000, _hex_to_rgb(color_end),   a2),
    ]
    spPr = _get_or_add_spPr(shape._element)
    _add_grad_fill(spPr, stops, angle_deg=angle_deg)


# =========================
# 图片填充：生成渐变 PNG + <a:blipFill>
# =========================
def _make_linear_gradient_png(px_w, px_h, c1, c2, angle_deg=0) -> str:
    """
    生成线性渐变 PNG 文件路径；angle_deg 支持 0/90/180/270 精确，其他角做近似。
    """
    im = Image.new("RGB", (px_w, px_h), c1)
    px = im.load()
    if angle_deg % 180 == 0:  # 左→右 / 右→左
        rev = (angle_deg % 360 == 180)
        for x in range(px_w):
            t = x / max(1, px_w - 1)
            if rev:
                t = 1 - t
            r = int(c1[0] * (1 - t) + c2[0] * t)
            g = int(c1[1] * (1 - t) + c2[1] * t)
            b = int(c1[2] * (1 - t) + c2[2] * t)
            for y in range(px_h):
                px[x, y] = (r, g, b)
    elif angle_deg % 180 == 90:  # 下→上 / 上→下
        rev = (angle_deg % 360 == 270)
        for y in range(px_h):
            t = y / max(1, px_h - 1)
            if rev:
                t = 1 - t
            r = int(c1[0] * (1 - t) + c2[0] * t)
            g = int(c1[1] * (1 - t) + c2[1] * t)
            b = int(c1[2] * (1 - t) + c2[2] * t)
            for x in range(px_w):
                px[x, y] = (r, g, b)
    else:
        import math
        nx = abs(math.cos(math.radians(angle_deg)))
        ny = abs(math.sin(math.radians(angle_deg)))
        for y in range(px_h):
            for x in range(px_w):
                t = (x / max(1, px_w - 1)) * nx + (y / max(1, px_h - 1)) * ny
                t = max(0.0, min(1.0, t))
                r = int(c1[0] * (1 - t) + c2[0] * t)
                g = int(c1[1] * (1 - t) + c2[1] * t)
                b = int(c1[2] * (1 - t) + c2[2] * t)
                px[x, y] = (r, g, b)
    fd, path = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    im.save(path)
    return path

def _set_shape_picture_fill(shape, image_path: str):
    """
    用图片作为形状填充（<a:blipFill>），兼容没有 FillFormat.user_picture 的版本。
    """
    image_part, rId = shape.part.get_or_add_image_part(image_path)
    spPr = _get_or_add_spPr(shape._element)
    _remove_fills(spPr)

    blipFill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), rId)
    blipFill.append(blip)

    stretch = OxmlElement("a:stretch")
    fillRect = OxmlElement("a:fillRect")
    stretch.append(fillRect)
    blipFill.append(stretch)

    spPr.append(blipFill)


# =========================
# 对外：添加渐变形状
# =========================
def add_gradient_shape(
    slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(1.5),
    color_start="#FFD54F", color_end="#9E9E9E",
    angle_deg=0,
    method="ooxml",          # "ooxml" or "picture"
    alpha_start=1.0, alpha_end=1.0,
    remove_line=True
):
    """
    在 slide 上添加一个“渐变色图形”，返回 shape。
    - shape_type: MSO_SHAPE（矩形/圆/三角/箭头/文本框等）
    - left/top/width/height: pptx 长度对象或数字（英寸）
    - color_start/color_end: "#RRGGBB"
    - angle_deg: 0=左→右，90=下→上
    - method: "ooxml"=原生梯度；"picture"=位图填充（最兼容）
    - alpha_*: 透明度 0~1
    """
    L, T, W, H = map(_to_length, (left, top, width, height))
    shp = slide.shapes.add_shape(shape_type, L, T, W, H)
    if remove_line:
        shp.line.fill.background()

    if method == "ooxml":
        _set_shape_linear_gradient(
            shp, color_start=color_start, color_end=color_end,
            angle_deg=angle_deg, alpha_start=alpha_start, alpha_end=alpha_end
        )
    elif method == "picture":
        # 以 96DPI 估算位图尺寸
        px_w = max(2, int(W.inches * 96))
        px_h = max(2, int(H.inches * 96))
        c1, c2 = _hex_to_rgb(color_start), _hex_to_rgb(color_end)
        png = _make_linear_gradient_png(px_w, px_h, c1, c2, angle_deg=angle_deg)
        _set_shape_picture_fill(shp, png)
    else:
        raise ValueError("method 必须为 'ooxml' 或 'picture'")

    return shp


# =========================
# 对外：整页背景（原生渐变 & 图片填充）
# =========================
def set_slide_background_linear_gradient(
    slide,
    color_start="#FFD54F", color_end="#9E9E9E",
    angle_deg=90, alpha_start=1.0, alpha_end=1.0,
    stops: List[Tuple[int, Tuple[int, int, int], int]] | None = None
):
    """
    整页原生线性渐变背景。stops 传则覆盖 color/alpha。
    """
    if stops is None:
        a1 = int(max(0, min(1, alpha_start)) * 100000)
        a2 = int(max(0, min(1, alpha_end)) * 100000)
        stops = [
            (0,      _hex_to_rgb(color_start), a1),
            (100000, _hex_to_rgb(color_end),   a2),
        ]
    bg = _get_or_add_bg(slide._element)
    bgPr = _get_or_add_bgPr(bg)
    _add_grad_fill(bgPr, stops=stops, angle_deg=angle_deg)

def set_slide_background_picture(slide, image_path: str):
    """
    整页背景图片填充：<p:bgPr><a:blipFill r:embed="..."><a:stretch><a:fillRect/></a:stretch></a:blipFill>
    兼容没有 slide.background.fill.user_picture 的版本。
    """
    image_part, rId = slide.part.get_or_add_image_part(image_path)
    bg = _get_or_add_bg(slide._element)
    bgPr = _get_or_add_bgPr(bg)

    # 清掉原有填充
    for tag in ("bgRef", "blipFill", "gradFill", "grpFill", "noFill", "solidFill", "pattFill"):
        el = bgPr.find(qn(f"a:{tag}"))
        if el is not None:
            bgPr.remove(el)

    blipFill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), rId)
    blipFill.append(blip)

    stretch = OxmlElement("a:stretch")
    fillRect = OxmlElement("a:fillRect")
    stretch.append(fillRect)
    blipFill.append(stretch)

    bgPr.append(blipFill)

def add_gradient_background(
    slide,
    color_start: str = "#FFD54F",
    color_end:   str = "#9E9E9E",
    angle_deg:   int = 90,          # 0=左→右，90=下→上
    method:      str = "ooxml",     # "ooxml" | "picture"
    alpha_start: float = 1.0,
    alpha_end:   float = 1.0,
    stops: List[Tuple[int, Tuple[int, int, int], int]] | None = None,
    px_width:    int = 1920,        # method="picture" 时，生成位图的尺寸
    px_height:   int = 1080
):
    """
    为整页背景设置渐变：
      - method="ooxml": 原生线性渐变（写 a:gradFill；文件更轻，推荐）
      - method="picture": 生成渐变 PNG 并作为背景图片填充（最兼容）
    """
    if method == "ooxml":
        set_slide_background_linear_gradient(
            slide,
            color_start=color_start,
            color_end=color_end,
            angle_deg=angle_deg,
            alpha_start=alpha_start,
            alpha_end=alpha_end,
            stops=stops
        )
    elif method == "picture":
        if stops is not None and len(stops) >= 2:
            c1 = stops[0][1]
            c2 = stops[-1][1]
            a = angle_deg
        else:
            c1 = _hex_to_rgb(color_start)
            c2 = _hex_to_rgb(color_end)
            a = angle_deg
        png = _make_linear_gradient_png(px_width, px_height, c1, c2, angle_deg=a)
        set_slide_background_picture(slide, png)
    else:
        raise ValueError("method 必须为 'ooxml' 或 'picture'")

def add_solid_shape(
    slide,
    shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    *,
    left=Inches(1),    # 数字也可以，按英寸解释
    top=Inches(1),
    width=Inches(3),
    height=Inches(1.5),
    fill_color: str = "#FFD54F",    # 纯色填充，#RRGGBB 或 #RGB
    fill_alpha: float = 1.0,        # 不透明度：0.0(全透明) ~ 1.0(不透明)
    outline_color: str | None = None,     # None=无边框；或 "#RRGGBB"
    outline_width_pt: float | None = None # 线宽（pt），None=不改；0 或 None 都可视为无边框
):
    # 统一位置/尺寸为 pptx 的 Length 对象
    L = _to_length(left)
    T = _to_length(top)
    W = _to_length(width)
    H = _to_length(height)

    # 新建形状
    shp = slide.shapes.add_shape(shape_type, L, T, W, H)

    # 纯色填充
    shp.fill.solid()
    r, g, b = 255, 255, 255
    if type(fill_color) == str:
        r, g, b = _hex_to_rgb(fill_color)
        shp.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        shp.fill.fore_color.rgb = fill_color  # 直接用传入的颜色对象
    # python-pptx: transparency=0(不透明)~1(全透明)
    shp.fill.transparency = max(0.0, min(1.0, 1.0 - float(fill_alpha)))

    # 处理边框
    if outline_color is None or (isinstance(outline_width_pt, (int, float)) and outline_width_pt == 0):
        # 无边框
        shp.line.fill.background()
    else:
        shp.line.fill.solid()
        or_, og, ob = 0, 0, 0
        if type(outline_color) == str:
            or_, og, ob = _hex_to_rgb(outline_color)
            shp.line.fill.fore_color.rgb = RGBColor(or_, og, ob)
        else:
            shp.line.fill.fore_color.rgb = outline_color  # 直接用传入的颜色对象
        if isinstance(outline_width_pt, (int, float)) and outline_width_pt > 0:
            shp.line.width = Pt(outline_width_pt)

    return shp

def _pct100k(x: float) -> str:
    """将 0..1 的小数转为 OOXML 百分比单位 0..100000 的字符串。"""
    x = max(0.0, min(1.0, float(x)))
    return str(int(round(x * 100000)))

def fill_shape_with_image(
    shape,
    image_path: str,
    *,
    mode: str = "stretch",      # "stretch" | "tile"
    crop: tuple[float, float, float, float] = (0, 0, 0, 0),  # (left, top, right, bottom), 0..1
    tile_scale: float = 1.0,    # mode="tile" 时，缩放系数（1.0=100%）
    tile_offset: tuple[int, int] = (0, 0)  # (tx, ty) 平移，EMU；一般不需要改
):
    """
    将 image_path 指定的图片填充到给定 shape（原生 <a:blipFill>）。
    - mode="stretch": 拉伸铺满形状
    - mode="tile":    平铺（可设置 tile_scale / tile_offset）
    - crop=(l,t,r,b): 0..1，按比例从四边裁剪图片
    """
    # 1) 关系：把图片加入当前 part，拿到关系 id
    image_part, rId = shape.part.get_or_add_image_part(image_path)

    # 2) 取得/创建 <p:spPr> 并清理旧填充
    spPr = _get_or_add_spPr(shape._element)
    _remove_fills(spPr)

    # 3) 组装 <a:blipFill>
    blipFill = OxmlElement("a:blipFill")

    # a:blip r:embed
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), rId)
    blipFill.append(blip)

    # 可选裁剪 a:srcRect（按 0..100000 百分比）
    l, t, r, b = crop
    if any(v > 0 for v in (l, t, r, b)):
        srcRect = OxmlElement("a:srcRect")
        if l > 0: srcRect.set("l", _pct100k(l))
        if t > 0: srcRect.set("t", _pct100k(t))
        if r > 0: srcRect.set("r", _pct100k(r))
        if b > 0: srcRect.set("b", _pct100k(b))
        blipFill.append(srcRect)

    # 伸展 or 平铺
    if mode == "tile":
        tile = OxmlElement("a:tile")
        # sx/sy：以千分之一百分比为单位（即 100000 = 100%），和 a:lin 的单位一致
        scale = max(0.01, float(tile_scale))
        tile.set("sx", _pct100k(scale))   # 宽度缩放
        tile.set("sy", _pct100k(scale))   # 高度缩放
        # 可选平移（EMU）；一般不设
        tx, ty = tile_offset
        if isinstance(tx, (int, float)) and tx != 0:
            tile.set("tx", str(int(tx)))
        if isinstance(ty, (int, float)) and ty != 0:
            tile.set("ty", str(int(ty)))
        blipFill.append(tile)
    else:
        # 默认拉伸铺满形状
        stretch = OxmlElement("a:stretch")
        fillRect = OxmlElement("a:fillRect")
        stretch.append(fillRect)
        blipFill.append(stretch)

    spPr.append(blipFill)
    return shape  # 便于链式调用

def add_image_filled_shape(
    slide,
    image_path: str,
    *,
    shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(3),
    mode: str = "stretch",          # "stretch" | "tile"
    crop: tuple[float, float, float, float] = (0, 0, 0, 0),
    tile_scale: float = 1.0,
    tile_offset: tuple[int, int] = (0, 0),
    remove_line: bool = True
):
    """
    新建一个形状并用图片填充。返回新建的 shape。
    - 数字位置/尺寸按英寸解释；也可传 Inches()/Cm()/Pt()。
    """
    # 统一长度单位
    L = left if hasattr(left, "_EMU") else Inches(left)
    T = top if hasattr(top, "_EMU") else Inches(top)
    W = width if hasattr(width, "_EMU") else Inches(width)
    H = height if hasattr(height, "_EMU") else Inches(height)

    shp = slide.shapes.add_shape(shape_type, L, T, W, H)
    if remove_line:
        shp.line.fill.background()

    fill_shape_with_image(
        shp, image_path,
        mode=mode, crop=crop, tile_scale=tile_scale, tile_offset=tile_offset
    )
    return shp

EMU_PER_INCH = 914400  # 1英寸 = 914400 EMU

def add_textbox(
    slide,
    text,
    left,   # 英寸
    top,    # 英寸
    width,  # 英寸
    height, # 英寸
    font_size,          # pt
    font_color="#000000",   # 改为 #RRGGBB 或 #RGB
    font_name="Calibri",
    has_base_box=False,     # 是否绘制底框
    border_color="#000000", # #RRGGBB
    border_width=2,
    corner_radius=0.1,
    has_border=True,
    border_type="solid",    # solid | dashed | none
    fill_color=None,        # 背景填充，#RRGGBB 或 None
    inner_padding=0.10,     # 文本与边框内边距（英寸）
):
    """
    固定底框/文本框尺寸；若文本放不下，则仅缩小字体，直到放下或达到最小字号。
    has_base_box=False 时不绘制底框/边框，仅添加文本框；默认将内边距置 0。

    颜色参数均为 "#RRGGBB" 或 "#RGB" 字符串。
    """
    # 当没有底框时，内边距通常设为 0（也可由调用方显式传入非 0 值覆盖）
    effective_padding = inner_padding if has_base_box else 0.0

    # 1) 可选的底框（圆角矩形）——这里仍然把 hex 字符串直接传给 add_solid_shape
    shape = None
    if has_base_box:
        outline_enabled = has_border and (str(border_type).lower() != "none")
        shape = add_solid_shape(
            slide,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height,
            fill_color=fill_color if fill_color is not None else "#FFFFFF",
            outline_color=border_color if outline_enabled else None,
            outline_width_pt=border_width if outline_enabled else None
            # 如需传 dash 样式，且 add_solid_shape 支持，可追加参数：outline_style=border_type
        )
        # 圆角
        try:
            if hasattr(shape, "adjustments") and shape.adjustments:
                shape.adjustments[0] = max(0.0, min(1.0, float(corner_radius)))
        except Exception:
            pass

    # 2) 固定尺寸的文本框（不随文本变化）
    tb_left   = Inches(left + effective_padding)
    tb_top    = Inches(top + effective_padding)
    tb_width  = Inches(max(0.01, width  - 2 * effective_padding))
    tb_height = Inches(max(0.01, height - 2 * effective_padding))
    text_box = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    p.alignment = PP_ALIGN.LEFT
    run.font.name = font_name
    font = run.font

    r, g, b = _hex_to_rgb(font_color or "#000000")
    font.color.rgb = RGBColor(r, g, b)

    # ========= 手动收缩字号算法（仅改字号，不改形状尺寸） =========
    avail_w_pt = (width  - 2 * effective_padding) * 72.0
    avail_h_pt = (height - 2 * effective_padding) * 72.0

    CJK_RE = re.compile(r'[\u4E00-\u9FFF\u3400-\u4DBF\uF900-\uFAFF]')

    def visual_units(s: str) -> float:
        vu = 0.0
        for ch in s:
            if CJK_RE.match(ch):
                vu += 1.0
            elif ch.isspace():
                vu += 0.3
            else:
                vu += 0.5
        return max(vu, 0.0001)

    # def required_height_pt(txt: str, size_pt: float, width_pt: float) -> float:
    #     lines = 0
    #     for raw_line in (txt.splitlines() or [""]):
    #         units = visual_units(raw_line)
    #         per_line_units = max(int(width_pt / max(size_pt * 1.0, 1e-6)), 1)
    #         wraps = int(math.ceil(units / per_line_units)) if units > 0 else 1
    #         lines += max(wraps, 1)
    #     return lines * (1.15 * size_pt)
    def required_height_pt(txt: str, size_pt: float, width_pt: float) -> float:
        """
        估算在给定字号与宽度下需要的高度。

        逻辑：
        1. 先按手动换行符 splitlines() 拆成多行；
        2. 对每一行单独估算“视觉宽度”并计算需要几行 wraps；
        3. 所有 wraps 相加得到总行数，再乘以行高。
        """
        logical_lines = txt.splitlines() or [""]

        # 对给定字号和宽度，对应的“每行能容纳的视觉单位数”
        per_line_units = max(int(width_pt / max(size_pt * 1.0, 1e-6)), 1)

        total_lines = 0
        for raw_line in logical_lines:
            units = visual_units(raw_line)
            if units <= 0:
                wraps = 1   # 空行也算一行高度
            else:
                wraps = int(math.ceil(units / per_line_units))
            total_lines += max(wraps, 1)

        # 1.15 是行距系数，可按视觉效果微调
        return total_lines * (1.15 * size_pt)


    target = float(font_size)
    min_size = 8.0
    while target > min_size:
        need_h = required_height_pt(text, target, avail_w_pt)
        if need_h <= avail_h_pt:
            break
        target -= 1.0
    font.size = Pt(max(target, min_size))

    return shape, text_box



def add_line(
    slide,
    x1, y1,           # 起点坐标（英寸）
    x2, y2,           # 终点坐标（英寸）
    color="#000000",  # #RRGGBB 或 #RGB
    width_pt=2.0,     # 线宽（pt）
    dash_style="solid"  # "solid" | "dashed" | "dot" | "dash_dot" | "dash_dot_dot"
                        # "long_dash" | "long_dash_dot" | "square_dot" | "mixed"
):
    """
    在指定幻灯片上添加一条直线（连接器）。

    颜色参数为 "#RRGGBB" 或 "#RGB" 字符串。
    """
    begin_x = Inches(x1)
    begin_y = Inches(y1)
    end_x   = Inches(x2)
    end_y   = Inches(y2)

    if x1 == x2 and y1 == y2:
        raise ValueError("add_line: 起点和终点不能完全重合，否则线段长度为 0。")

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        begin_x, begin_y,
        end_x, end_y
    )

    line = connector.line
    r, g, b = _hex_to_rgb(color or "#000000")
    line.color.rgb = RGBColor(r, g, b)
    line.width = Pt(width_pt)
    dash_style = (dash_style or "solid").lower()

    if dash_style == "solid":
        line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    elif dash_style == "dashed":
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif dash_style == "dash_dot":
        line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT
    elif dash_style in ("dash_dot_dot", "dashdotdot"):
        line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT_DOT
    elif dash_style == "long_dash":
        line.dash_style = MSO_LINE_DASH_STYLE.LONG_DASH
    elif dash_style == "long_dash_dot":
        line.dash_style = MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    elif dash_style in ("dot", "round_dot"):
        line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    elif dash_style == "square_dot":
        line.dash_style = MSO_LINE_DASH_STYLE.SQUARE_DOT
    elif dash_style == "mixed":
        line.dash_style = MSO_LINE_DASH_STYLE.DASH_STYLE_MIXED
    else:
        line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    return connector

# =========================
# 示例（直接运行测试）
# =========================
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 背景：原生渐变（自下→上，灰→黄）
    add_gradient_background(
        slide, color_start="#9E9E9E", color_end="#FFD54F",
        angle_deg=90, method="ooxml"
    )

    # 形状1：原生渐变，圆角矩形，左→右
    add_gradient_shape(
        slide,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        left=1, top=1.1, width=6, height=1.2,   # 数字=英寸
        color_start="#FFD54F", color_end="#9E9E9E",
        angle_deg=0, method="ooxml"
    )

    # 形状2：图片填充，椭圆，上→下（270°）
    add_gradient_shape(
        slide,
        shape_type=MSO_SHAPE.OVAL,
        left=1, top=3, width=3, height=2,
        color_start="#9E9E9E", color_end="#FFD54F",
        angle_deg=270, method="picture"
    )

    out = "gradient_shapes_demo.pptx"
    prs.save(out)
    print("Saved:", out)
