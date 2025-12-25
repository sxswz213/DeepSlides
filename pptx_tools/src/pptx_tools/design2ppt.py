
# -*- coding: utf-8 -*-
"""
design2ppt_v3_ppttools.py
Converter for design_formatting_prompt_v3 adapted to pptx_tools.add_free_shape API,
with an optional `path` parameter to save the PPT after rendering.

Usage:
    from pptx_tools.add_free_shape import *
    from design2ppt_v3_ppttools import render_design_to_slide_v3_ppttools
    render_design_to_slide_v3_ppttools(slide, design_str, path="output.pptx")
"""

from __future__ import annotations
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple
import json
import re

from pptx.util import Inches

from pptx.enum.shapes import MSO_SHAPE  # shape types
from pptx import Presentation

# ---- Canvas ----
CANVAS_W, CANVAS_H = 13.33, 7.50
def _pt_from_in(inches: float) -> float: return float(inches) * 72.0

# ---- Data classes ----
@dataclass
class Border:
    on: bool = False
    weight: float = 0.02   # inches
    color: str = "#000000"
    dash: str = "solid"
    opacity: Optional[float] = None

@dataclass
class BackgroundSpec:
    type: str = "none"  # none|solid|gradient|shape
    color: Optional[str] = None
    details: Dict[str, Any] = field(default_factory=dict)
    opacity: float = 1.0
    shape: str = "rect"   # rect|rounded|circle
    radius: float = 0.0

@dataclass
class BBox:
    x: float; y: float; w: float; h: float
    def right(self) -> float: return self.x + self.w
    def bottom(self) -> float: return self.y + self.h
    def as_line(self) -> Tuple[float, float, float, float]:
        return (self.x, self.y, self.x + self.w, self.y + self.h)

@dataclass
class BlockNode:
    id: str
    function: str
    bbox: BBox
    background: BackgroundSpec = field(default_factory=BackgroundSpec)
    border: Border = field(default_factory=Border)
    description: str = ""
    z: int = 0
    rotation: float = 0.0
    children: List["BlockNode"] = field(default_factory=list)

@dataclass
class Separator:
    type: str
    bbox: BBox
    color: str = "#E0E0E0"
    style: Border = field(default_factory=Border)

# ---- Parsing ----
SECTION_RE = re.compile(r"^\[(Meta|Background|Layout|Content)\]\s*$", re.MULTILINE)

def _extract_balanced_json(s: str) -> str:
    s = s.strip()
    if not s or s[0] != "{": return s
    depth = 0
    for i, ch in enumerate(s):
        if ch == "{": depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0: return s[:i+1]
    return s

def parse_design_string(design_str: str) -> Dict[str, Any]:
    parts = {}
    matches = list(SECTION_RE.finditer(design_str))
    if len(matches) < 3:
        raise ValueError("Design string must contain [Meta], [Background], [Layout], [Content] sections.")
    for i, m in enumerate(matches):
        start = m.end()
        end = matches[i+1].start() if i+1 < len(matches) else len(design_str)
        parts[m.group(1)] = design_str[start:end].strip()

    meta = {}
    if parts.get("Meta"):
        meta = json.loads(_extract_balanced_json(parts["Meta"]))

    bg_text = parts.get("Background","")
    machine_json = {}
    machine_match = re.search(r"Machine\s*:\s*(\{[\s\S]*\})", bg_text)
    if machine_match:
        machine_json = json.loads(_extract_balanced_json(machine_match.group(1)))

    layout_text = parts.get("Layout","")
    struct_match = re.search(r"Structure\s*:\s*(\{[\s\S]*\})", layout_text)
    if not struct_match:
        raise ValueError("Layout section missing 'Structure: { ... }' JSON.")
    layout_struct = json.loads(_extract_balanced_json(struct_match.group(1)))

    content_json = {}
    if parts.get("Content"):
        content_json = json.loads(_extract_balanced_json(parts["Content"]))

    return {"meta": meta, "background_machine": machine_json, "layout": layout_struct, "content": content_json}

# ---- Geometry ----
def _round2(x: float) -> float: return round(float(x), 2)
def _bbox(d: Dict[str, Any]) -> BBox: return BBox(_round2(d["x"]), _round2(d["y"]), _round2(d["w"]), _round2(d["h"]))
def _within_canvas(b: BBox) -> bool:
    return 0 <= b.x <= CANVAS_W and 0 <= b.y <= CANVAS_H and 0 <= b.right() <= CANVAS_W and 0 <= b.bottom() <= CANVAS_H
def _overlap(a: BBox, b: BBox) -> bool:
    return not (a.right() <= b.x or b.right() <= a.x or a.bottom() <= b.y or b.bottom() <= a.y)
def _inside(outer: BBox, inner: BBox) -> bool:
    return outer.x <= inner.x and outer.y <= inner.y and outer.right() >= inner.right() and outer.bottom() >= inner.bottom()

# ---- Tree & validation ----
def _parse_block_tree(tree: Dict[str, Any]) -> List[BlockNode]:
    nodes: List[BlockNode] = []
    for key, val in (tree.get("blocks") or {}).items():
        bid = val.get("id", key)
        bg_raw = val.get("background") or {}
        node = BlockNode(
            id=bid,
            function=val.get("function","text"),
            bbox=_bbox(val["bbox_in"]),
            background=BackgroundSpec(
                type=bg_raw.get("type","none"),
                color=bg_raw.get("color"),
                details=bg_raw.get("details",{}),
                opacity=float(bg_raw.get("opacity",1.0)),
                shape=bg_raw.get("shape","rect"),
                radius=float(bg_raw.get("radius",0.0))
            ),
            border=Border(**(val.get("border") or {})) if val.get("border") else Border(),
            description=val.get("description",""),
            z=int(val.get("z",0)),
            rotation=float(val.get("rotation",0.0)),
            children=[]
        )
        if "children" in val and isinstance(val["children"], dict):
            node.children = _parse_block_tree({"blocks": val["children"]})
        nodes.append(node)
    return nodes

def _validate(nodes: List[BlockNode]) -> None:
    for i in range(len(nodes)):
        if not _within_canvas(nodes[i].bbox):
            raise ValueError(f"Block out of canvas: {nodes[i].id}")
        for j in range(i+1, len(nodes)):
            if _overlap(nodes[i].bbox, nodes[j].bbox):
                raise ValueError(f"Blocks overlap at same level: {nodes[i].id} vs {nodes[j].id}")
        _validate(nodes[i].children)
        for c in nodes[i].children:
            if not _inside(nodes[i].bbox, c.bbox):
                raise ValueError(f"Child {c.id} not inside parent {nodes[i].id}")

def _collect(nodes: List[BlockNode]) -> List[BlockNode]:
    acc: List[BlockNode] = []
    for n in nodes:
        acc.append(n); acc.extend(_collect(n.children))
    return acc

# ---- Policy ----
@dataclass
class Policy:
    avoid_content: bool = True
    safe_insets: Dict[str, float] = field(default_factory=lambda: {"top":0.0,"right":0.0,"bottom":0.0,"left":0.0})
    max_alpha: Optional[float] = None

def _policy_from_machine(machine: Dict[str, Any]) -> Policy:
    pol = machine.get("policy") or {}
    return Policy(
        avoid_content=bool(pol.get("avoid_content", True)),
        safe_insets={
            "top": float(pol.get("safe_insets", {}).get("top", 0.0)),
            "right": float(pol.get("safe_insets", {}).get("right", 0.0)),
            "bottom": float(pol.get("safe_insets", {}).get("bottom", 0.0)),
            "left": float(pol.get("safe_insets", {}).get("left", 0.0)),
        },
        max_alpha=float(pol["max_alpha"]) if "max_alpha" in pol else None
    )

def _apply_safe_insets(b: BBox, ins: Dict[str, float]) -> BBox:
    x = max(b.x, ins.get("left",0.0))
    y = max(b.y, ins.get("top",0.0))
    r = min(b.right(), CANVAS_W - ins.get("right",0.0))
    bt = min(b.bottom(), CANVAS_H - ins.get("bottom",0.0))
    return BBox(_round2(x), _round2(y), _round2(max(0.0, r-x)), _round2(max(0.0, bt-y)))

def _should_skip(b: BBox, content_blocks: List[BlockNode], policy: Policy) -> bool:
    if not policy.avoid_content:
        return False
    for blk in content_blocks:
        if _overlap(b, blk.bbox):
            return True
    return False

def _cap_opacity(op: Optional[float], policy: Policy) -> float:
    if op is None: return 1.0
    if policy.max_alpha is None: return float(op)
    return float(min(op, policy.max_alpha))

# ---- Utilities ----
def _shape_type(name: str):
    name = (name or "rect").lower()
    if name in ("rect","rectangle"): return MSO_SHAPE.RECTANGLE
    if name in ("rounded","round","rounded_rect","roundedrectangle"): return MSO_SHAPE.ROUNDED_RECTANGLE
    if name in ("circle","oval","ellipse"): return MSO_SHAPE.OVAL
    if name in ("triangle","tri"): return MSO_SHAPE.ISOSCELES_TRIANGLE
    return MSO_SHAPE.RECTANGLE

def _dash_map(name: str) -> str:
    m = (name or "solid").lower()
    table = {
        "solid":"solid", "dot":"dot", "dotted":"dot", "dashed":"dashed",
        "dash":"dashed", "dash_dot":"dash_dot", "dashdot":"dash_dot",
        "dash_dot_dot":"dash_dot_dot", "long_dash":"long_dash", "long_dash_dot":"long_dash_dot",
        "square_dot":"square_dot", "mixed":"mixed"
    }
    return table.get(m, "solid")

# ---- Background renderers (ppt_tools) ----
def _bg_fill(slide, layer: Dict[str, Any]):
    from pptx_tools.add_free_shape import add_solid_shape, add_image_filled_shape, add_gradient_background
    fill = layer.get("fill") or {}
    mode = fill.get("mode","solid")
    if mode == "solid":
        add_solid_shape(slide, left=0.0, top=0.0, width=CANVAS_W, height=CANVAS_H,
                        shape_type=MSO_SHAPE.RECTANGLE, fill_color=fill.get("color","#FFFFFF"),
                        fill_alpha=1.0, outline_color=None, outline_width_pt=None)
    elif mode == "gradient":
        grad = fill.get("gradient") or {}
        colors = grad.get("colors", ["#FFFFFF","#F5F7FB"])
        c1 = colors[0]; c2 = colors[1] if len(colors) > 1 else colors[0]
        angle = float(grad.get("angle", 0))
        add_gradient_background(slide, color_start=c1, color_end=c2, angle_deg=angle, method="ooxml")
    elif mode == "image":
        img = fill.get("image") or {}
        src = img.get("src")
        if src:
            add_image_filled_shape(slide, image_path=src, left=0.0, top=0.0, width=CANVAS_W, height=CANVAS_H,
                                   shape_type=MSO_SHAPE.RECTANGLE, mode="stretch")

def _bg_pattern(slide, layer: Dict[str, Any], policy: Policy, content_blocks: List[BlockNode]):
    from pptx_tools.add_free_shape import add_solid_shape, add_line
    pat = layer.get("pattern") or {}
    primitive = pat.get("primitive","dot")
    tile = pat.get("tile") or {"w":0.5,"h":0.5}
    repeat = pat.get("repeat") or {}
    cols = int(repeat.get("cols", int(CANVAS_W / tile["w"])))
    rows = int(repeat.get("rows", int(CANVAS_H / tile["h"])))
    origin = pat.get("origin") or {"x":0.0,"y":0.0}
    gap = pat.get("gap") or {"x":0.0,"y":0.0}
    jitter = pat.get("jitter") or {"x":0.0,"y":0.0}
    st = pat.get("style") or {}
    size = float(st.get("size", 0.06))
    weight_pt = _pt_from_in(float(st.get("weight", 0.02)))
    color = st.get("color", "#AAB7C7")
    opacity = _cap_opacity(float(st.get("opacity", 0.15)), policy)
    dash = _dash_map(st.get("dash", "solid"))
    shape = st.get("shape", "circle")

    y = origin.get("y",0.0)
    for r in range(rows):
        x = origin.get("x",0.0)
        for c in range(cols):
            if primitive in ("dot","diamond","triangle"):
                bx = _round2(x + 0.5*tile["w"] - 0.5*size + (jitter.get("x",0.0) if (r+c)%2==0 else -jitter.get("x",0.0)))
                by = _round2(y + 0.5*tile["h"] - 0.5*size + (jitter.get("y",0.0) if (r+c)%2==0 else -jitter.get("y",0.0)))
                b = BBox(bx, by, _round2(size), _round2(size))
                b2 = _apply_safe_insets(b, policy.safe_insets)
                if _should_skip(b2, content_blocks, policy): 
                    x += tile["w"] + gap.get("x",0.0); continue
                add_solid_shape(slide, left=b2.x, top=b2.y, width=b2.w, height=b2.h,
                                shape_type=_shape_type("circle" if shape=="circle" else ("rounded" if shape=="rounded" else "rect")),
                                fill_color=color, fill_alpha=opacity, outline_color=None, outline_width_pt=None)
            elif primitive in ("stripe","line","chevron"):
                x1 = _round2(x); y1 = _round2(y + 0.5*tile["h"])
                x2 = _round2(x + tile["w"]); y2 = y1
                bb = _apply_safe_insets(BBox(min(x1,x2), min(y1,y2), abs(x2-x1), abs(y2-y1)), policy.safe_insets)
                if _should_skip(bb, content_blocks, policy):
                    x += tile["w"] + gap.get("x",0.0); continue
                add_line(slide, x1, y1, x2, y2, color=color, width_pt=weight_pt, dash_style=dash)
            else:
                bx = _round2(x + 0.5*tile["w"] - 0.5*size)
                by = _round2(y + 0.5*tile["h"] - 0.5*size)
                b = BBox(bx, by, _round2(size), _round2(size))
                b2 = _apply_safe_insets(b, policy.safe_insets)
                if _should_skip(b2, content_blocks, policy): 
                    x += tile["w"] + gap.get("x",0.0); continue
                add_solid_shape(slide, left=b2.x, top=b2.y, width=b2.w, height=b2.h,
                                shape_type=MSO_SHAPE.OVAL, fill_color=color, fill_alpha=opacity,
                                outline_color=None, outline_width_pt=None)
            x += tile["w"] + gap.get("x",0.0)
        y += tile["h"] + gap.get("y",0.0)

def _bg_grid(slide, layer: Dict[str, Any], policy: Policy, content_blocks: List[BlockNode]):
    from pptx_tools.add_free_shape import add_line
    grid = layer.get("grid") or {}
    cols = int(grid.get("cols", 12))
    rows = int(grid.get("rows", 6))
    stroke = grid.get("stroke") or {}
    color = stroke.get("color","#DCE4EE")
    weight_pt = _pt_from_in(float(stroke.get("weight",0.02)))
    dash = _dash_map(stroke.get("dash","dot"))
    _ = _cap_opacity(float(stroke.get("opacity",0.2)), policy)  # ppt_tools doesn't support alpha directly
    inner = grid.get("inner_margin") or {"x":0.0,"y":0.0}
    iw = CANVAS_W - 2*inner.get("x",0.0)
    ih = CANVAS_H - 2*inner.get("y",0.0)
    cw = iw / max(1, cols)
    rh = ih / max(1, rows)
    x0 = inner.get("x",0.0)
    y0 = inner.get("y",0.0)

    for c in range(cols+1):
        x = _round2(x0 + c*cw)
        bb = _apply_safe_insets(BBox(x, y0, 0.0, ih), policy.safe_insets)
        if _should_skip(bb, content_blocks, policy): continue
        add_line(slide, x, y0, x, y0+ih, color=color, width_pt=weight_pt, dash_style=dash)
    for r in range(rows+1):
        y = _round2(y0 + r*rh)
        bb = _apply_safe_insets(BBox(x0, y, iw, 0.0), policy.safe_insets)
        if _should_skip(bb, content_blocks, policy): continue
        add_line(slide, x0, y, x0+iw, y, color=color, width_pt=weight_pt, dash_style=dash)

def _bg_ornaments(slide, layer: Dict[str, Any], policy: Policy, content_blocks: List[BlockNode]):
    from pptx_tools.add_free_shape import add_solid_shape, add_gradient_shape, add_image_filled_shape
    orns = layer.get("ornaments") or []
    for item in orns:
        shape_name = item.get("shape","rounded")
        bbox = _apply_safe_insets(_bbox(item.get("bbox_in")), policy.safe_insets)
        if _should_skip(bbox, content_blocks, policy): continue
        fill = item.get("fill") or {}
        mode = fill.get("mode","solid")
        op = _cap_opacity(float(item.get("opacity",0.1)), policy)
        if mode == "solid":
            add_solid_shape(slide, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                            shape_type=_shape_type(shape_name), fill_color=fill.get("color","#000000"),
                            fill_alpha=op, outline_color=None, outline_width_pt=None)
        elif mode == "gradient":
            grad = fill.get("gradient") or {}
            colors = grad.get("colors", ["#FFFFFF","#EEEEEE"])
            c1 = colors[0]; c2 = colors[1] if len(colors)>1 else colors[0]
            angle = float(grad.get("angle",0))
            add_gradient_shape(slide, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                               shape_type=_shape_type(shape_name),
                               color_start=c1, color_end=c2, angle_deg=angle, method="ooxml")
        elif mode == "image":
            img = fill.get("image") or {}
            src = img.get("src")
            if src:
                add_image_filled_shape(slide, image_path=src, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                                       shape_type=_shape_type(shape_name), mode=img.get("mode","stretch"))

def _bg_shapes(slide, layer: Dict[str, Any], policy: Policy, content_blocks: List[BlockNode]):
    from pptx_tools.add_free_shape import add_solid_shape, add_gradient_shape, add_image_filled_shape, add_line
    elems = layer.get("elements") or []
    for e in elems:
        kind = e.get("kind","shape")
        if kind == "shape":
            shape_name = e.get("shape","rect")
            bbox = _bbox(e.get("bbox_in"))
            bbox = _apply_safe_insets(BBox(max(0.0, bbox.x), max(0.0, bbox.y),
                                min(bbox.w, CANVAS_W - max(0.0, bbox.x)),
                                min(bbox.h, CANVAS_H - max(0.0, bbox.y))), policy.safe_insets)
            if bbox.w == 0 or bbox.h == 0: continue
            if _should_skip(bbox, content_blocks, policy): continue
            fill = e.get("fill") or {}
            mode = fill.get("mode","solid")
            if mode == "solid":
                add_solid_shape(slide, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                                shape_type=_shape_type(shape_name), fill_color=fill.get("color","#FFFFFF"),
                                fill_alpha=_cap_opacity(float(e.get("opacity",1.0)), policy),
                                outline_color=None, outline_width_pt=None)
            elif mode == "gradient":
                grad = fill.get("gradient") or {}
                colors = grad.get("colors", ["#FFFFFF","#EEEEEE"])
                c1, c2 = colors[0], (colors[1] if len(colors)>1 else colors[0])
                angle = float(grad.get("angle",0))
                add_gradient_shape(slide, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                                   shape_type=_shape_type(shape_name),
                                   color_start=c1, color_end=c2, angle_deg=angle, method="ooxml")
            elif mode == "image":
                img = fill.get("image") or {}
                src = img.get("src")
                if src:
                    add_image_filled_shape(slide, image_path=src, left=bbox.x, top=bbox.y, width=bbox.w, height=bbox.h,
                                           shape_type=_shape_type(shape_name), mode=img.get("mode","stretch"))
        elif kind == "line":
            ln = e.get("line") or {}
            x1, y1 = float(ln.get("x1",0.0)), float(ln.get("y1",0.0))
            x2, y2 = float(ln.get("x2",0.0)), float(ln.get("y2",0.0))
            bb = _apply_safe_insets(BBox(min(x1,x2), min(y1,y2), abs(x2-x1), abs(y2-y1)), policy.safe_insets)
            if _should_skip(bb, content_blocks, policy): continue
            st = e.get("stroke") or {}
            add_line(slide, x1, y1, x2, y2,
                     color=st.get("color","#E6EAF0"),
                     width_pt=_pt_from_in(float(st.get("weight",0.02))),
                     dash_style=_dash_map(st.get("dash","solid")))

# ---- Blocks & content ----
def _render_block_bg(slide, b: BlockNode):
    from pptx_tools.add_free_shape import add_solid_shape, add_gradient_shape, add_image_filled_shape
    bg = b.background
    if bg.type == "solid" and bg.color:
        add_solid_shape(slide, left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                        shape_type=_shape_type(bg.shape), fill_color=bg.color, fill_alpha=bg.opacity,
                        outline_color=(b.border.color if b.border.on else None),
                        outline_width_pt=(_pt_from_in(b.border.weight) if b.border.on else None))
    elif bg.type == "gradient":
        colors = (bg.details.get("colors") or ["#FFFFFF","#EEEEEE"])
        c1, c2 = colors[0], (colors[1] if len(colors)>1 else colors[0])
        ang = float(bg.details.get("angle",0))
        add_gradient_shape(slide, left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                           shape_type=_shape_type(bg.shape),
                           color_start=c1, color_end=c2, angle_deg=ang, method="ooxml")
    elif bg.type == "shape":
        if "image" in bg.details:
            add_image_filled_shape(slide, image_path=bg.details["image"],
                                   left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                                   shape_type=_shape_type(bg.shape), mode=bg.details.get("mode","stretch"))
        else:
            add_solid_shape(slide, left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                            shape_type=_shape_type(bg.shape),
                            fill_color=(bg.color or "#FFFFFF"), fill_alpha=bg.opacity,
                            outline_color=(b.border.color if b.border.on else None),
                            outline_width_pt=(_pt_from_in(b.border.weight) if b.border.on else None))

def _render_separators(slide, seps: List[Separator]):
    from pptx_tools.add_free_shape import add_line
    for s in seps:
        x1, y1, x2, y2 = s.bbox.as_line()
        add_line(slide, x1, y1, x2, y2,
                 color=s.style.color,
                 width_pt=_pt_from_in(s.style.weight),
                 dash_style=_dash_map(s.style.dash))

def _content_lookup(content: Dict[str, Any], path: List[str]) -> Optional[Dict[str, Any]]:
    cur = content
    for k in path:
        if isinstance(cur, dict) and k in cur:
            cur = cur[k]
        else:
            return None
    return cur if isinstance(cur, dict) else None

def _render_block_content(slide, b: BlockNode, content: Dict[str, Any], id_path: List[str]):
    from pptx_tools.add_free_shape import add_textbox, add_image_filled_shape
    myc = _content_lookup(content, id_path+[b.id]) or content.get(b.id)
    if isinstance(myc, dict):
        if "text" in myc:
            st = myc.get("style", {})
            padx = round(0.10 * b.bbox.w, 2)
            pady = round(0.10 * b.bbox.h, 2)
            lx = round(b.bbox.x + padx, 2)
            ty = round(b.bbox.y + pady, 2)
            ww = round(max(0.01, b.bbox.w - 2*padx), 2)
            hh = round(max(0.01, b.bbox.h - 2*pady), 2)
            add_textbox(
                slide=slide,
                text=myc.get("text",""),
                # left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                left=lx, top=ty, width=ww, height=hh,
                font_size=int(st.get("size",18)),
                font_name=st.get("font","Calibri"),
                font_color=st.get("color","#000000"),
                has_base_box=bool(st.get("border", False)),
                border_color=(b.border.color if b.border.on else "#000000"),
                border_width=int(_pt_from_in(b.border.weight) if b.border.on else 0),
                corner_radius=float(st.get("radius", 0.0)),
                has_border=bool(st.get("border", False)),
                fill_color=st.get("fill_color", None),
                inner_padding=float(st.get("padding", 0.10))
            )
        if "image" in myc:
            img = myc.get("image")
            if img and img != "none":
                st = myc.get("style", {})
                mode = "stretch"  # cover/contain can be approximated inside your helper if available
                add_image_filled_shape(slide=slide, image_path=img,
                                       left=b.bbox.x, top=b.bbox.y, width=b.bbox.w, height=b.bbox.h,
                                       shape_type=MSO_SHAPE.RECTANGLE, mode=mode, crop=None,
                                       tile_scale=1.0, remove_line=True)
    for c in b.children:
        _render_block_content(slide, c, content, id_path+[b.id])

# ---- Saving helper ----
def _save_via_slide_package(slide, path: str) -> bool:
    """Try python-pptx internals: slide.part.package.save(path) or package.presentation.save(path)."""
    try:
        part = getattr(slide, "part", None)
        if part is None: return False
        package = getattr(part, "package", None)
        if package is not None and hasattr(package, "save"):
            package.save(path); return True
        pres = getattr(package, "presentation", None) if package is not None else None
        if pres is not None and hasattr(pres, "save"):
            pres.save(path); return True
    except Exception:
        return False
    return False

# ---- Public API ----
def design_to_slide(slide, design_str: str, path: str | None = None):
    parsed = parse_design_string(design_str)
    machine = parsed.get("background_machine") or {}
    policy = _policy_from_machine(machine)

    # Layout
    nodes = _parse_block_tree(parsed["layout"])
    _validate(nodes)
    flat = sorted(_collect(nodes), key=lambda n: n.z)

    # Avoidance set
    content_blocks = flat

    # Background layers
    for layer in sorted((machine.get("layers") or []), key=lambda l: l.get("z",0)):
        t = layer.get("type")
        if t == "fill": _bg_fill(slide, layer)
        elif t == "pattern": _bg_pattern(slide, layer, policy, content_blocks)
        elif t == "grid": _bg_grid(slide, layer, policy, content_blocks)
        elif t == "ornaments": _bg_ornaments(slide, layer, policy, content_blocks)
        elif t == "shapes": _bg_shapes(slide, layer, policy, content_blocks)

    # Separators
    seps: List[Separator] = []
    for s in (parsed["layout"].get("separators") or []):
        bb = _bbox(s.get("bbox_in"))
        style = Border(**(s.get("style") or {})) if s.get("style") else Border(on=True, weight=0.02, color=s.get("color","#E6EAF0"), dash="solid")
        seps.append(Separator(type=s.get("type","line"), bbox=bb, color=s.get("color","#E6EAF0"), style=style))
    _render_separators(slide, seps)

    # Block backgrounds
    for b in flat:
        _render_block_bg(slide, b)

    # Content
    for b in flat:
        _render_block_content(slide, b, parsed["content"], id_path=[])

    # Optional save
    if path:
        ok = _save_via_slide_package(slide, path)
        if not ok:
            raise RuntimeError("Unable to save via slide context. "
                               "Call Presentation.save(path) yourself or ensure slide.part.package.save is available.")

    return True



# ---- High-level convenience API ----
def render_design_to_ppt(design_str: str, path: str | None = None):
    """
    Create a new PowerPoint Presentation, render the given v3 design onto a blank slide,
    optionally save to `path`, and return the Presentation object.

    Returns:
        prs (Presentation): the in-memory Presentation object.
    Raises:
        ValueError/RuntimeError on invalid spec or save failure.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Reuse the slide-level renderer (no saving at slide level)
    design_to_slide(slide, design_str, path=None)
    if path:
        prs.save(path)
    return prs
