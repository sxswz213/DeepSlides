from typing import Literal
import json

from langchain.chat_models import init_chat_model
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.runnables import RunnableConfig
# from goto import goto, label
from langgraph.constants import Send
from langgraph.graph import START, END, StateGraph
from langgraph.types import interrupt, Command

import subprocess
import os
import datetime
import asyncio
import io
import httpx
from urllib.parse import urlparse
import re
import shutil
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx_tools.design2ppt import *
from pathlib import Path

import base64
from typing import Optional, List, Dict, Any, TypedDict


from open_deep_research.state import (
    ReportStateInput,
    ReportStateOutput,
    Sections,
    ReportState,
    SectionState,
    SectionOutputState,
    Queries,
    Feedback,
    PPTSection,
    PPTOutline,
    PPTSections,
    PPTSectionState,
    PPTSlideState,
    PPTSlide,
    PPTSlideOutputState,
    PPTSectionOutputState,
    SearchQuery
)

from open_deep_research.prompts import (
    report_planner_query_writer_instructions,
    report_planner_instructions,
    query_writer_instructions, 
    section_writer_instructions,
    final_section_writer_instructions,
    section_grader_instructions,
    section_writer_inputs,
    query_writer4PPT_instructions,
    ppt_tools_prompt,
    design_formatting_prompt,
    evaluation_design,
    evaluation_aesthetics,
    evaluation_complete,
    style_plan_prompt,
    eval_cover,
    color_examples_prompt,
    code_prefix
)

from open_deep_research.configuration import Configuration
from open_deep_research.utils import (
    format_sections, 
    get_config_value, 
    get_search_params, 
    select_and_execute_search,
    set_openai_api_base,
    generate_image_caption,
    generate_image_caption_v2,
    generate_image_caption_v3
)
MODE="openai" # ["openai","azure"]
COMP_MODE="llm" # ["llm","tools"]
# Ablation study switches — set via env vars: ABLATE_DESIGN=1 / ABLATE_SCORING=1
ABLATE_DESIGN  = os.getenv("ABLATE_DESIGN",  "0") == "1"
ABLATE_SCORING = os.getenv("ABLATE_SCORING", "0") == "1"
# Save root directory — set via env var: SAVES_ROOT=saves_ablation_design
SAVES_ROOT = os.getenv("SAVES_ROOT", "saves_sonnet")
## Nodes -- 

async def process_image_input(state: ReportState, config: RunnableConfig):
    """处理图像输入并生成描述文本。
    
    这个节点：
    1. 检查状态中是否包含图像路径
    2. 如果有图像，使用Vision API生成详细描述
    3. 将描述作为额外上下文添加到后续步骤
    4. 如果未提供topic，则使用图像描述作为topic
    
    Args:
        state: 当前图状态，包含可选的图像路径
        config: 配置参数
        
    Returns:
        包含图像描述和更新topic的状态
    """
    # 检查是否提供了图像路径
    image_path = state.get("image_path")
    
    # 如果没有提供图像路径，直接返回状态
    if not image_path:
        # 确保topic存在，即使为空字符串
        if "topic" not in state:
            return {"topic": ""}
        return {}
    
    # 确保设置了正确的API基础URL
    set_openai_api_base()
    
    try:
        if not state.get("topic"):
            topic = "The user did not provide a topic."
        else:
            topic = state["topic"]
        # 生成图像描述
        image_result = await generate_image_caption_v3(image_path, topic)
        image_result = json.loads(image_result)
        # print(image_result)
        caption, user_intent, topic = image_result["caption"], image_result["user_intent"], image_result["topic"]

        return {"caption": caption, "user_intent": user_intent, "topic": topic}

    except Exception as e:
        print(f"Error processing image: {str(e)}")
        # 返回错误信息作为caption，并确保topic存在
        if "topic" not in state:
            return {"image_caption": f"Unable to process image:{str(e)}", "topic": ""}
        return {"image_caption": f"Unable to process image:{str(e)}"}

async def generate_report_plan(state: ReportState, config: RunnableConfig):
    """Generate the initial report plan with sections, including image caption and user intent.

    This node:
    1. Gets configuration for the report structure and search parameters
    2. Generates search queries to gather context for planning
    3. Performs web searches using those queries
    4. Uses an LLM to generate a structured plan with sections
    5. Includes image caption and user intent in planning

    Args:
        state: Current graph state containing the report topic, image caption, and user intent
        config: Configuration for models, search APIs, etc.

    Returns:
        Dict containing the generated sections
    """

    # 获取topic、图片caption和用户意图
    topic = state["topic"]
    caption = state.get("image_caption", "")
    user_intent = state.get("user_intent", "")
    feedback = state.get("feedback_on_report_plan", None)

    # Get configuration
    configurable = Configuration.from_runnable_config(config)
    report_structure = configurable.report_structure
    number_of_queries = configurable.number_of_queries
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}
    params_to_pass = get_search_params(search_api, search_api_config)

    if isinstance(report_structure, dict):
        report_structure = str(report_structure)

    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    structured_llm = writer_model.with_structured_output(Queries)

    # Format system instructions including caption and user intent
    system_instructions_query = report_planner_query_writer_instructions.format(
        topic=topic,
        caption=caption,
        user_intent=user_intent,
        report_organization=report_structure,
        number_of_queries=number_of_queries
    )

    results = await structured_llm.ainvoke([
        SystemMessage(content=system_instructions_query),
        HumanMessage(content="Generate search queries that help plan each section of the report.")
    ])

    query_list = [query.search_query for query in results.queries][:number_of_queries]
    print(query_list)
    source_str = await select_and_execute_search(search_api, query_list, params_to_pass)

    image_path = state.get("image_path")
    if image_path:
        try:
            # 单独调用图片搜索API
            image_search_result = await select_and_execute_search("image_search", [image_path], params_to_pass)
            # 将图片搜索结果与之前的搜索结果合并
            source_str += f"\n\nImage search results:\n{image_search_result}"
        except Exception as e:
            print(f"Error calling the image-search API:{str(e)}")

    # Include caption and user intent in report planner instructions
    system_instructions_sections = report_planner_instructions.format(
        topic=topic,
        caption=caption,
        user_intent=user_intent,
        report_organization=report_structure,
        context=source_str,
        feedback=feedback
    )

    planner_provider = get_config_value(configurable.planner_provider)
    planner_model = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})

    planner_message = """Generate the sections of the report. Your reply must contain a 'sections' field that lists the sections.  
Each section must include: name, description, plan, research, and content fields."""

    set_openai_api_base()

    if planner_model == "claude-3-7-sonnet-latest":
        planner_llm = init_chat_model(
            model=planner_model,
            model_provider=planner_provider,
            max_tokens=20_000,
            thinking={"type": "enabled", "budget_tokens": 16_000}
        )
    else:
        # With other models, thinking tokens are not specifically allocated
        if MODE == "openai":
            planner_llm = init_chat_model(model=planner_model, 
                                      model_provider=planner_provider,
                                      model_kwargs={})
        else:
            planner_llm = AzureChatOpenAI(
            model=configurable.writer_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],  # Azure's API base
            deployment_name=planner_model_kwargs["azure_deployment"],  # Azure's deployment name
            openai_api_version=planner_model_kwargs["openai_api_version"],  # Azure's API version
            temperature=0,
            max_tokens=4096
        )

    # Generate the report sections
    structured_llm = planner_llm.with_structured_output(Sections)
    report_sections = await structured_llm.ainvoke([SystemMessage(content=system_instructions_sections),
                                             HumanMessage(content=planner_message)])

    sections = report_sections.sections

    return {"sections": sections}

def human_feedback(state: ReportState, config: RunnableConfig) -> Command[Literal["generate_report_plan","build_section_with_web_research"]]:
    """Get human feedback on the report plan and route to next steps.
    
    This node:
    1. Formats the current report plan for human review
    2. Gets feedback via an interrupt
    3. Routes to either:
       - Section writing if plan is approved
       - Plan regeneration if feedback is provided
    
    Args:
        state: Current graph state with sections to review
        config: Configuration for the workflow
        
    Returns:
        Command to either regenerate plan or start section writing
    """

    # Get sections
    topic = state["topic"]
    sections = state['sections']
    sections_str = "\n\n".join(
        f"Section: {section.name}\n"
        f"Description: {section.description}\n"
        f"Research needed: {'Yes' if section.research else 'No'}\n"
        for section in sections
    )

    # # Get feedback on the report plan from interrupt
    # interrupt_message = f"""Please provide feedback on the following report plan.\n\n{sections_str}\n
    # Does this report plan meet your needs?\nPass in 'true' to approve the report plan.\nAlternatively, provide feedback to regenerate the report plan:"""
    
    # feedback_dict = interrupt(interrupt_message)
    # if type(feedback_dict) == bool:
    #     # If the feedback is a boolean, treat it as approval
    #     feedback = feedback_dict
    # else:
    #     feedback = list(feedback_dict.values())[0]
    # # print(feedback)
    # # If the user approves the report plan, kick off section writing
    # if isinstance(feedback, bool) and feedback is True:
    #     # Treat this as approve and kick off section writing
    #     return Command(goto=[
    #         Send("build_section_with_web_research", {"topic": topic, "section": s, "search_iterations": 0}) 
    #         for s in sections 
    #         if s.research
    #     ])
    
    # # If the user provides feedback, regenerate the report plan 
    # elif isinstance(feedback, str):
    #     # Treat this as feedback
    #     return Command(goto="generate_report_plan", 
    #                    update={"feedback_on_report_plan": feedback})
    # else:
    #     raise TypeError(f"Interrupt value of type {type(feedback)} is not supported.")
    # tooooooodoooooo 此处为忽略人类反馈直接进行下一步
    return Command(
        goto=[
            Send(
                "build_section_with_web_research",
                {"topic": topic, "section": section, "search_iterations": 0},
            )
            for section in sections
            if section.research
        ]
    )
    
async def generate_queries(state: SectionState, config: RunnableConfig):
    """Generate search queries for researching a specific section.
    
    This node uses an LLM to generate targeted search queries based on the 
    section topic and description.
    
    Args:
        state: Current state containing section details
        config: Configuration including number of queries to generate
        
    Returns:
        Dict containing the generated search queries
    """

    # Get state 
    topic = state["topic"]
    section = state["section"]

    # Get configuration
    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate queries 
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    structured_llm = writer_model.with_structured_output(Queries)

    # Format system instructions
    system_instructions = query_writer_instructions.format(topic=topic, 
                                                           section_topic=section.description, 
                                                           number_of_queries=number_of_queries)

    # Generate queries
      
    queries = await structured_llm.ainvoke([SystemMessage(content=system_instructions),
                                     HumanMessage(content="Generate search queries for the given topic.")])
    print(queries)
    return {"search_queries": queries.queries}

async def search_web(state: SectionState, config: RunnableConfig):
    """Execute web searches for the section queries.
    
    This node:
    1. Takes the generated queries
    2. Executes searches using configured search API
    3. Formats results into usable context
    
    Args:
        state: Current state with search queries
        config: Search API configuration
        
    Returns:
        Dict with search results and updated iteration count
    """

    # Get state
    search_queries = state["search_queries"]
    # Get configuration

    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}  # Get the config dict, default to empty
    params_to_pass = get_search_params(search_api, search_api_config)  # Filter parameters

    # Web search
    query_list = [query.search_query for query in search_queries][:number_of_queries]

    # Search the web with parameters
    source_str = await select_and_execute_search(search_api, query_list, params_to_pass)

    return {"source_str": source_str, "search_iterations": state["search_iterations"] + 1}

async def write_section(state: SectionState, config: RunnableConfig) -> Command[Literal[END, "search_web"]]:
    """Write a section of the report and evaluate if more research is needed.
    
    This node:
    1. Writes section content using search results
    2. Evaluates the quality of the section
    3. Either:
       - Completes the section if quality passes
       - Triggers more research if quality fails
    
    Args:
        state: Current state with search results and section info
        config: Configuration for writing and evaluation
        
    Returns:
        Command to either complete section or do more research
    """

    # Get state 
    topic = state["topic"]
    section = state["section"]
    source_str = state["source_str"]
    
    # 提取图像信息
    images_data = []
    # TODO: make it configurable
    max_images = 6  # 最大图像数量限制

    if "--- IMAGES ---" in source_str:
        try:
            # 提取图像部分
            images_section = source_str.split("--- IMAGES ---")[1].split("-" * 80)[0]
            image_blocks = images_section.strip().split("IMAGE ")[1:]  # 跳过第一个空元素
            
            for i, block in enumerate(image_blocks):
                lines = block.strip().split("\n")
                image_url = ""
                image_description = ""
                
                for line in lines:
                    if line.startswith("URL:"):
                        image_url = line.replace("URL:", "").strip()
                    elif line.startswith("DESCRIPTION:"):
                        image_description = line.replace("DESCRIPTION:", "").strip()
                
                if image_url:  # 只添加有URL的图像
                    images_data.append({
                        "index": i,
                        "url": image_url,
                        "description": image_description
                    })
        except Exception as e:
            print(f"Error extracting image information: {str(e)}")

    # 达到最大图像数量后停止处理
    image_num_available = len(images_data)
    if image_num_available >= max_images:
        images_data = images_data[:max_images]

    # 将图像数据格式化为JSON字符串
    images_json = json.dumps(images_data, ensure_ascii=False, indent=2) if images_data else "[]"
    
    if images_data:
        # print(f"Extracted {image_num_available} images (maximum limit: {max_images})")
        pass

    # Get configuration
    configurable = Configuration.from_runnable_config(config)

    # Format system instructions
    section_writer_inputs_formatted = section_writer_inputs.format(
        topic=topic, 
        section_name=section.name, 
        section_topic=section.description, 
        context=source_str, 
        section_content=section.content,
        images_data=images_json
    )

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate section  
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )

    # TODO: Native image input
    section_content = await writer_model.ainvoke([SystemMessage(content=section_writer_instructions),
                                           HumanMessage(content=section_writer_inputs_formatted)])
    
    # 处理返回的内容，提取图像选择信息
    content = section_content.content

    # 检查是否包含图像选择信息
    if "```image_selection" in content and images_data:
        try:
            # 提取图像选择JSON
            image_selection_text = content.split("```image_selection")[1].split("```")[0].strip()
            image_selection = json.loads(image_selection_text)
            
            selected_index = image_selection.get("selected_image_index", -1)
            if 0 <= selected_index < len(images_data):
                selected_image = images_data[selected_index]
                selected_image["caption"] = image_selection.get("caption", "")
                
                # 处理内容分割，保留标记前后的内容
                before_image_selection = content.split("```image_selection")[0].strip()
                after_image_selection = ""
                if "```" in content.split("```image_selection")[1]:
                    after_image_selection = content.split("```image_selection")[1].split("```", 1)[1].strip()
                
                # 重新组合内容，将图像选择部分替换为图像和标题标记
                image_content = f"\n\n![{selected_image['caption']}]({selected_image['url']})\n*{selected_image['caption']}*\n"
                if after_image_selection:
                    content = f"{before_image_selection}\n\n{image_content}\n\n{after_image_selection}"
                else:
                    content = f"{before_image_selection}\n\n{image_content}"
            else:
                # selected_index不合法，移除图像选择部分
                before_image_selection = content.split("```image_selection")[0].strip()
                after_image_selection = ""
                if "```" in content.split("```image_selection")[1]:
                    after_image_selection = content.split("```image_selection")[1].split("```", 1)[1].strip()
                
                # 重新组合内容，移除图像选择部分但保留其前后内容
                if after_image_selection:
                    content = f"{before_image_selection}\n\n{after_image_selection}"
                else:
                    content = before_image_selection
        except Exception as e:
            print(f"Error processing image selection information: {str(e)}")
    
    # Write content to the section object  
    section.content = content
    section.source_str = source_str  # Store the source string in the section

    # Grade prompt 
    section_grader_message = (
        "Grade the report and consider follow-up questions for any missing information. "
        "If the grade is 'pass', all subsequent queries should return an empty string. "
        "If the grade is 'fail', provide specific search queries to collect the missing information."
    )

    
    section_grader_instructions_formatted = section_grader_instructions.format(topic=topic, 
                                                                               section_topic=section.description,
                                                                               section=section.content, 
                                                                               number_of_follow_up_queries=configurable.number_of_queries)

    # Use planner model for reflection
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})

    # Set OpenAI API Base URL
    set_openai_api_base()

    if planner_model == "claude-3-7-sonnet-latest":
        # Allocate a thinking budget for claude-3-7-sonnet-latest as the planner model
        reflection_model = init_chat_model(model=planner_model, 
                                           model_provider=planner_provider, 
                                           max_tokens=20_000, 
                                           thinking={"type": "enabled", "budget_tokens": 16_000}).with_structured_output(Feedback)
    else:
        if MODE == "openai":
            reflection_model = init_chat_model(model=planner_model, 
                                            model_provider=planner_provider, model_kwargs={}).with_structured_output(Feedback)
        else:
            reflection_model = AzureChatOpenAI(
            model=configurable.writer_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],  # Azure's API base
            deployment_name=planner_model_kwargs["azure_deployment"],  # Azure's deployment name
            openai_api_version=planner_model_kwargs["openai_api_version"],  # Azure's API version
            # temperature=0,
            max_tokens=2048
        ).with_structured_output(Feedback)
    # Generate feedback
    feedback = await reflection_model.ainvoke([SystemMessage(content=section_grader_instructions_formatted),
                                        HumanMessage(content=section_grader_message)])

    # If the section is passing or the max search depth is reached, publish the section to completed sections 
    if feedback.grade == "pass" or state["search_iterations"] >= configurable.max_search_depth:
        # Publish the section to completed sections 
        return  Command(
        update={"completed_sections": [section]},
        goto=END
    )

    # Update the existing section with new content and update search queries
    else:
        return  Command(
        update={"search_queries": feedback.follow_up_queries, "section": section},
        goto="search_web"
        )
    
async def write_final_sections(state: SectionState, config: RunnableConfig):
    """Write sections that don't require research using completed sections as context.
    
    This node handles sections like conclusions or summaries that build on
    the researched sections rather than requiring direct research.
    
    Args:
        state: Current state with completed sections as context
        config: Configuration for the writing model
        
    Returns:
        Dict containing the newly written section
    """

    # Get configuration
    configurable = Configuration.from_runnable_config(config)

    # Get state 
    topic = state["topic"]
    section = state["section"]
    completed_report_sections = state["report_sections_from_research"]
    
    # Format system instructions
    system_instructions = final_section_writer_instructions.format(topic=topic, section_name=section.name, section_topic=section.description, context=completed_report_sections)

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate section  
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    
    section_content = await writer_model.ainvoke([SystemMessage(content=system_instructions),
                                           HumanMessage(content="Generate the report sections based on the provided information.")])
    
    # Write content to section 
    section.content = section_content.content

    # Write the updated section to completed sections
    return {"completed_sections": [section]}

def gather_completed_sections(state: ReportState):
    """Format completed sections as context for writing final sections.
    
    This node takes all completed research sections and formats them into
    a single context string for writing summary sections.
    
    Args:
        state: Current state with completed sections
        
    Returns:
        Dict with formatted sections as context
    """

    # List of completed sections
    completed_sections = state["completed_sections"]

    # Format completed section to str to use as context for final sections
    completed_report_sections = format_sections(completed_sections)

    return {"report_sections_from_research": completed_report_sections}

async def compile_final_report(state: ReportState):
    """Compile all sections into the final report.
    
    This node:
    1. Gets all completed sections
    2. Orders them according to original plan
    3. Combines them into the final report
    
    Args:
        state: Current state with all completed sections
        
    Returns:
        Dict containing the complete report
    """

    # Get sections
    topic = state["topic"]
    sections = state["sections"]
    completed_sections = {s.name: s.content for s in state["completed_sections"]}

    # Update sections with completed content while maintaining original order
    for section in sections:
        section.content = completed_sections[section.name]

    # Compile final report
    all_sections = "\n\n".join([s.content for s in sections])

    save_dir = os.path.join(".", SAVES_ROOT, topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    outline_path = os.path.join(save_dir, "final_report.md")


    try:
        await asyncio.to_thread(
            lambda p=outline_path, d=all_sections: open(p, "w", encoding="utf-8").write(
                json.dumps(d, ensure_ascii=False, indent=2)
            )
        )
    except Exception as exc:
        outline_path = ""

    return {"final_report": all_sections}

def initiate_final_section_writing(state: ReportState):
    """Create parallel tasks for writing non-research sections.
    
    This edge function identifies sections that don't need research and
    creates parallel writing tasks for each one.
    
    Args:
        state: Current state with all sections and research context
        
    Returns:
        List of Send commands for parallel section writing
    """

    # Kick off section writing in parallel via Send() API for any sections that do not require research
    return [
        Send("write_final_sections", {"topic": state["topic"], "section": s, "report_sections_from_research": state["report_sections_from_research"]}) 
        for s in state["sections"] 
        if not s.research
    ]

async def generate_ppt_outline(state: ReportState, config: RunnableConfig):
    """
    根据演讲时长确定推荐的PPT页数，根据风格和故事线重新生成章节划分，最后基于此生成PPT大纲。

    Args:
        state: 当前状态，包含 final_report 等信息
        config: 配置参数

    Returns:
        包含PPT大纲的状态字典
    """
    configurable = Configuration.from_runnable_config(config)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.planner_model,
        azure_endpoint=planner_model_kwargs["openai_api_base"],
        deployment_name=planner_model_kwargs["azure_deployment"],
        openai_api_version=planner_model_kwargs["openai_api_version"],
        temperature=0,
        max_tokens=4096
    )

    topic = state.get("topic", "none")
    presentation_minutes = state.get("presentation_minutes", "10")
    style = state.get("style", "none")

    if style == "none":
        prefix = (
            "Reference presentation styles: professional business, modern tech, minimalist, "
            "creative lively, academically rigorous, storytelling narrative, magazine visual, "
            "illustration cartoon, retro nostalgic, data visualization."
        )
    else:
        prefix = f"User's expected presentation style/storyline: {style}"

    storyline_prompt = f"""
    You are an experienced presentation expert tasked with creating a presentation PPT. Now you need to determine the presentation's storyline.
    Presentation topic: {topic}
    Presentation duration: {presentation_minutes} minutes
    {prefix}

    Reference storyline templates:
    - Problem-Solution: Clearly identify a core problem and provide clear, specific solutions.
    - Situation-Conflict-Resolution-Outcome: First set up a scene, describe the challenge, offer the solution, and finally present positive results.
    - SCQA (Situation-Complication-Question-Answer): Provide background information, introduce the complication, state the key question clearly, and give the answer.
    - Timeline (Past-Present-Future): Present past events, current status, and future goals in chronological order.
    - Contrast (Current vs. Future): Clearly contrast existing problems with the ideal future state, highlighting how to achieve the transformation.
    - Pyramid: Start with the conclusion and unfold arguments layer by layer from top to bottom, reinforcing the core idea with rigorous, clear logic.
    - Research Report: Introduce the topic within its field background, systematically outline the current status, methods, and challenges, and finally present future research trends and directions.

    Return **only JSON format**:
    ```json
    {{
        "storyline": "Storyline type",
    }}
    ```
    """

    storyline_response = await writer_model.ainvoke([
        SystemMessage(content=storyline_prompt),
        HumanMessage(content="Please recommend a suitable storyline.")
    ])
    response_content = json.loads(storyline_response.content.split("```json")[-1].split('```')[0])
    # style = response_content["style"]
    storyline = response_content["storyline"]
    # main_color = response_content["main_color"]
    # accent_color = response_content["accent_color"]
    # background_tone = response_content["background_tone"]
    # heading_font_color = response_content["heading_font_color"]
    # body_font_color = response_content["body_font_color"]




    ppt_length_prompt = f"""
    You are an experienced presentation expert.

    Topic: {topic}
    Presentation duration: {presentation_minutes} minutes

    Please suggest an appropriate number of PPT slides for this duration
    (each slide should have a moderate amount of content, not overcrowded;
    one slide generally corresponds to about 1-2 minutes of presentation time).

    **OUTPUT ONLY JSON format**: 
    ```json
    {{\"recommended_slides\": 10}}
    ```
    """


    ppt_length_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_length_prompt),
        HumanMessage(content="Please provide the recommended number of PPT slides.")] 
    )
    print(f"PPT length response: {ppt_length_response.content}")
    recommended_slides = json.loads(ppt_length_response.content.split("```json")[-1].split('```')[0])["recommended_slides"]
    # print(f"Recommended number of slides: {recommended_slides}")

    ppt_section_distribution_prompt = f"""
    Presentation topic: {topic}
    Style: {style}
    Storyline: {storyline}
    Recommended total slides: {recommended_slides}

    Attention: Do NOT create a "Q&A" or "Closing" section.

    Based on the above information, re-plan the PPT section structure and allocate the number of slides for each section. Return in JSON format, for example:
    {{
        "section_distribution": {{
            "Introduction": 2,
            "Methodology": 3,
            "Results": 3,
            "Conclusion": 2
        }}
    }}
    """


    ppt_distribution_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_section_distribution_prompt),
        HumanMessage(content="Please plan the section structure and allocate the number of pages.")
    ])

    section_distribution = json.loads(ppt_distribution_response.content)["section_distribution"]
    # print(f"Section distribution: {section_distribution}")
    ppt_outline_prompt = f"""
    You excel at designing slide outlines for presentations.
    Presentation topic: {topic}
    Storyline: {storyline}

    Reference material for the presentation:
    {state["final_report"]}

    The slide allocation for each section is as follows:
    {json.dumps(section_distribution, ensure_ascii=False, indent=2)}

    Please generate a PPT outline that adheres to the above slide allocation. Each slide should include:
    - A title
    - key points 
    - Layout (single block / top-bottom / left-right / n horizontal blocks / card grid / ...); encourage diverse and innovative layouts. 

    Important: **Do NOT make every slide contain exactly 3 or 4 key points.**  
    Ensure variety by creating some slides with **5** key points, and some with **6**.

    JSON format:
    Return JSON only, e.g.:
    ```json
    {{
        "ppt_sections": [
            {{
            "name": "Introduction",
            "allocated_slides": 2,
            "slides": [
                {{"title":"Sample 6-point slide","points":["A","B","C","D","E","F"],"layout":"..."}},
                {{"title":"Sample 4-point slide","points":["A","B","C","D"],"layout":"..."}},
                {{"title":"Sample 5-point slide","points":["A","B","C","D","E"],"layout":"..."}},
                {{"title":"Sample 3-point slide","points":["A","B","C"],"layout":"..."}}
            ]
            }}
        ]
    }}
    ```
    """


    ppt_outline_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_outline_prompt),
        HumanMessage(content="Generate a PPT outline.")
    ])
    # print(f"PPT outline response: {ppt_outline_response.content}")
    ppt_sections_data = json.loads(ppt_outline_response.content.split("```json")[-1].split('```')[0])["ppt_sections"]
    # print(f"PPT outline sections: {ppt_sections_data}")
    for section in ppt_sections_data:
        for slide in section.get("slides", []):
            slide.setdefault("codes", [])
            slide.setdefault("detail", "")
            slide.setdefault("enriched_points", "")
            slide.setdefault("path", "")
            

    ppt_sections = [PPTSection(**section) for section in ppt_sections_data]
    ppt_outline = PPTOutline(ppt_sections=PPTSections(sections=ppt_sections))
    save_dir = os.path.join(".", SAVES_ROOT, "outlines", topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    outline_path = os.path.join(save_dir, "ppt_outline.json")

    outline_payload = {
        "recommended_slides": recommended_slides,
        "section_distribution": section_distribution,
        "ppt_outline": [section.model_dump() for section in ppt_sections],
    }

    try:
        await asyncio.to_thread(
            lambda p=outline_path, d=outline_payload: open(p, "w", encoding="utf-8").write(
                json.dumps(d, ensure_ascii=False, indent=2)
            )
        )
    except Exception as exc:
        print(f"[WARN] Error saving PPT report: {exc}")
        outline_path = ""
    return {
            "recommended_ppt_slides": recommended_slides,
            "section_distribution": section_distribution,
            "ppt_outline": ppt_outline,
            "ppt_sections": ppt_sections,
            "storyline": storyline,
        }
    # return Command(
    #     update={
    #         "recommended_ppt_slides": recommended_slides,
    #         "section_distribution": section_distribution,
    #         "ppt_outline": ppt_outline,
    #         "ppt_sections": ppt_sections,
    #         "storyline": storyline,
    #         "style": style,
    #         "main_color": main_color,  
    #         "accent_color": accent_color
    #     },
    #     goto=[
    #         Send("generate_ppt_sections", {"topic": topic, "ppt_section": ppt_section, "style": style, "main_color": main_color, "accent_color": accent_color})
    #         for ppt_section in ppt_sections
    #     ]
    # )

async def generate_ppt_styles(state: ReportState, config: RunnableConfig):
    """
    根据PPT大纲生成整体风格建议，包括颜色搭配、字体选择、设计元素等。

    Args:
        state: 当前状态，包含 ppt_outline 等信息
        config: 配置参数

    Returns:
        包含风格建议的状态字典
    """
    configurable = Configuration.from_runnable_config(config)
    topic = state.get("topic", "none")
    style = state.get("style", "none")
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_model_name = get_config_value(configurable.coder_model)
    coder_provider = get_config_value(configurable.coder_provider)
    coder_base_url = get_config_value(configurable.coder_base_url)

    designer_model_kwargs = get_config_value(configurable.designer_model_kwargs or {})
    designer_model_name = get_config_value(configurable.designer_model)
    designer_provider = get_config_value(configurable.designer_provider)
    designer_base_url = get_config_value(configurable.designer_base_url)
    openai_key = os.getenv("OPENAI_API_KEY")
    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
        designer_model = init_chat_model(model=designer_model_name, model_provider=designer_provider, model_kwargs={})
    else:
        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        # temperature=0.7,
        # max_tokens=4096
        max_completion_tokens=40960,
        reasoning_effort = coder_model_kwargs["reasoning_effort"]
    )

    style_prompt= f"""
    You are an experienced presentation expert. You need to recommend a suitable style and color scheme for a presentation PPT.
    Presentation topic: {topic}
    User's expected presentation style: {style}
    Please select the appropriate style and color scheme for the topic based on the provided information. And it is also important to ensure the harmony between the main color and the secondary color.
{color_examples_prompt}
Return **only JSON format**:
    {{
        "style": "Style",
        "main_color": "Recommended primary color",
        "accent_color": "Recommended accent color",
        "background_tone": "light/dark dominated + background color description",
        "heading_font_color": "recommended heading font color",
        "body_font_color": "recommended body font color"
        "font_name": "recommended font name"
    }}
"""
    try:
        style_response = await designer_model.ainvoke([
            SystemMessage(content=style_prompt),
            HumanMessage(content="Please recommend a suitable style and color scheme.")
        ])
    except Exception as e:
        # Diagnostic logging to help identify misconfiguration (model id / base_url)
        print(f"[WARN] designer_model.ainvoke failed: {e}")
        try:
            print(f"[DEBUG] designer_model_name={designer_model_name}, designer_provider={designer_provider}, designer_base_url={designer_base_url}")
        except Exception:
            pass

        # Try a lightweight fallback using the OpenAI client wrapper if available
        try:
            from langchain_openai import ChatOpenAI
            # Try to call OpenAI official API as a fallback (no custom base_url)
            try:
                openai_fallback = ChatOpenAI(
                    model=designer_model_name,
                    openai_api_key=os.getenv("OPENAI_API_KEY"),
                    max_tokens=40960,
                )
                print("[INFO] Attempting fallback with ChatOpenAI (official API)")
                style_response = await openai_fallback.ainvoke([
                    SystemMessage(content=style_prompt),
                    HumanMessage(content="Please recommend a suitable style and color scheme.")
                ])
            except Exception as e2:
                print(f"[ERROR] ChatOpenAI fallback failed: {e2}")
                raise
        except Exception as e3:
            # If fallback import or call fails, re-raise original exception for visibility
            print(f"[ERROR] No fallback available or fallback failed: {e3}")
            raise
    style_content = json.loads(style_response.content.split("```json")[-1].split('```')[0])
    style = style_content["style"]
    main_color = style_content["main_color"]
    accent_color = style_content["accent_color"]
    background_tone = style_content["background_tone"]
    heading_font_color = style_content["heading_font_color"]
    body_font_color = style_content["body_font_color"]
    font_name = style_content["font_name"]

    return {            
            "style": style,
            "main_color": main_color,  
            "accent_color": accent_color,
            "background_tone": background_tone,
            "heading_font_color": heading_font_color,
            "body_font_color": body_font_color,
            "font_name": font_name
            }


async def manage_ppt_templates(state: ReportState, config: RunnableConfig)-> Command[Literal["generate_ppt_sections"]]:
    template = state.get("template", "")
    topic = state.get("topic", "none")
    ppt_sections = state.get("ppt_sections", [])
    style = state.get("style", "none")
    main_color = state.get("main_color", "blue")
    accent_color = state.get("accent_color", "orange")
    background_tone = state.get("background_tone", "light")
    heading_font_color = state.get("heading_font_color", "black")
    body_font_color = state.get("body_font_color", "darkgray")
    font_name = state.get("font_name", "Arial")
    style_summary = state.get("style_summary", "")

    # template是一个PPT模板文件路径，若为空，则先生成模板



    return Command(

        goto=[
            Send("generate_ppt_sections", {"topic": topic, "ppt_section": ppt_section, "style": style, "main_color": main_color, "accent_color": accent_color
                                           , "background_tone": background_tone, "heading_font_color": heading_font_color, "body_font_color": body_font_color, "style_summary": style_summary, "font_name": font_name})
            for ppt_section in ppt_sections
        ]
    )


async def save_image_from_url(image_url, image_name, topic, ppt_section_name, slide_index):
    """
    异步下载图片并保存到本地，自动根据URL后缀决定图片类型。

    Args:
        image_url: 图片URL
        image_name: 保存的图片名称
        topic: 主题目录名称
        ppt_section_name: PPT章节名称
        slide_index: 幻灯片索引

    Returns:
        str: 图片保存路径或空字符串（失败时）
    """
    try:
        # 异步网络请求下载图片
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(image_url)
            response.raise_for_status()

            # 提取图片后缀
            parsed_url = urlparse(image_url)
            ext = os.path.splitext(parsed_url.path)[-1].lower()

            # 默认后缀为 .jpg，如果没有或后缀不标准时使用默认
            if ext not in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']:
                ext = '.jpg'

            # 异步创建目录（避免阻塞）
            save_dir = os.path.join(
                ".", SAVES_ROOT, topic, "images", ppt_section_name, f"slide_{slide_index+1}"
            )
            await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)

            # 构建完整图片路径
            image_path = os.path.join(save_dir, f"{image_name}{ext}")

            # 异步保存文件（避免阻塞）
            await asyncio.to_thread(
                lambda: open(image_path, 'wb').write(response.content)
            )

            # print(f"Image successfully saved to: {image_path}")
            return image_path

    except httpx.HTTPError as e:
        # print(f"HTTP request error: {str(e)}")
        return ""
    except Exception as e:
        print(f"Other error: {str(e)}")
        return ""

async def truncate_by_characters(text, max_chars=250000):
    if len(text) > max_chars:
        return text[:max_chars]
    return text

async def enrich_slide_content(state: PPTSlideState, config: RunnableConfig):
    """
    第一部分：根据幻灯片的标题和要点生成详细的内容描述，并生成幻灯片布局描述。

    Args:
        state: 包含当前PPT幻灯片信息的状态。
        config: 配置参数。

    Returns:
        Tuple: 包含生成的详细内容和幻灯片布局描述。
    """
    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries_for_ppt
    writer_model_kwargs = configurable.writer_model_kwargs or {}
    
    # 获取当前幻灯片信息
    topic = state["topic"]
    ppt_section = state["ppt_section"]
    slide_index = state["slide_index"]
    style = state.get("style")
    main_color = state.get("main_color")
    accent_color = state.get("accent_color")
    background_tone = state.get("background_tone")
    heading_font_color = state.get("heading_font_color")
    body_font_color = state.get("body_font_color")
    font_name = state.get("font_name")
    style_summary = state.get("style_summary", "")
    design_suggestions = state.get("design_suggestions", "")
    aestheitcs_suggestions = state.get("aestheitcs_suggestions", "")
    slide = ppt_section.slides[slide_index]
    slide_title = slide.title
    slide_points = slide.points
    slide_layout = slide.layout

    # 设置OpenAI API的基础URL
    set_openai_api_base()

    # 使用Azure模型生成查询语句
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    if MODE == "openai":
        writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs={})
    else:
        writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure的API基础URL
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure的部署名称
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure的API版本
        temperature=0,
        max_tokens=4096,
        # streaming=False
    )
    structured_llm = writer_model.with_structured_output(Queries)

    # 格式化系统指令
    system_instructions = query_writer4PPT_instructions.format(
        topic=topic,
        section_topic=ppt_section.name,
        slide_title=slide_title,
        slide_points=", ".join(slide_points),
        number_of_queries=number_of_queries
     )#+ """
# You must respond with a single valid JSON object that matches the following schema:

# {
#   "queries": [
#     {
#       // one SearchQuery object
#       // fields must strictly follow the SearchQuery schema defined in the system
#       // do not add extra fields
#     },
#     ...
#   ]
# }

# Constraints:
# - "queries" MUST be a JSON array.
# - The length of "queries" MUST equal `number_of_queries` provided above.
# - Do not add any commentary, explanation, or text outside the JSON object.
# """

    # 生成查询语句
    query_response = await structured_llm.ainvoke([
        SystemMessage(content=system_instructions),
        HumanMessage(content="Generate relevant search queries based on the information above.")
    ])
    # print(f"Generated search queries response: {query_response}")
    # data = json.loads(query_response.content)
    # queries_obj = Queries.parse_obj(data)

    # 从返回的查询结果中提取查询语句
    search_queries = [query.search_query for query in query_response.queries][:number_of_queries]
    # search_queries: list[SearchQuery] = queries_obj.queries

    # 执行Web搜索
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}
    params_to_pass = get_search_params(search_api, search_api_config)

    source_str = await select_and_execute_search(search_api, search_queries, params_to_pass)

    # 提取图像信息
    images_data = []
    # TODO: make it configurable
    max_images = 6  # 最大图像数量限制

    if "--- IMAGES ---" in source_str:
        try:
            # 提取图像部分
            images_section = source_str.split("--- IMAGES ---")[1].split("-" * 80)[0]
            image_blocks = images_section.strip().split("IMAGE ")[1:]  # 跳过第一个空元素
            
            for i, block in enumerate(image_blocks):
                lines = block.strip().split("\n")
                image_url = ""
                image_description = ""
                
                for line in lines:
                    if line.startswith("URL:"):
                        image_url = line.replace("URL:", "").strip()
                    elif line.startswith("DESCRIPTION:"):
                        image_description = line.replace("DESCRIPTION:", "").strip()
                
                if image_url:  # 只添加有URL的图像
                    # 保存图像并获取保存路径
                    image_path = await save_image_from_url(image_url, f"image_{i+1}", topic, ppt_section.name, slide_index)
                    images_data.append({
                        "index": i,
                        "url": image_url,
                        "description": image_description,
                        "local_path": image_path  # 添加本地图片路径
                    })
        except Exception as e:
            print(f"Error extracting image information: {str(e)}")

    # 达到最大图像数量后停止处理
    image_num_available = len(images_data)
    if image_num_available >= max_images:
        images_data = images_data[:max_images]

    # 将图像数据格式化为JSON字符串
    images_json = json.dumps(images_data, ensure_ascii=False, indent=2) if images_data else "[]"
    
    if images_data:
        # print(f"Extracted {image_num_available} images (maximum limit: {max_images})")
        pass

    images_data = images_data[:max_images] if len(images_data) >= max_images else images_data
    embedded_images = []
    for img in images_data:
        local_path = img.get("local_path", "")
        if local_path and os.path.exists(local_path):
            try:
                img_bytes = await asyncio.to_thread(
                    lambda path=local_path: open(path, "rb").read()
                )

                img_base64 = base64.b64encode(img_bytes).decode()
                ext = os.path.splitext(local_path)[1].lower().replace('.', '')
                mime_type = f'image/{ext if ext != "jpg" else "jpeg"}'

                embedded_images.append({
                    "index": img["index"],
                    "path": local_path,
                    "description": img["description"],
                    "height": img.get("height", ""),
                    "width": img.get("width", ""),
                    # "base64": f"data:{mime_type};base64,{img_base64}"
                })
            except Exception as e:
                print(f"Error loading image: {str(e)}")
    
    images_json_embedded = json.dumps(embedded_images, ensure_ascii=False, indent=2) if embedded_images else "[]"
    
    if embedded_images:
        # print(f"Successfully loaded {len(embedded_images)} images in Base64 format.")
        pass
    
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_model_name = get_config_value(configurable.coder_model)
    coder_provider = get_config_value(configurable.coder_provider)
    coder_base_url = get_config_value(configurable.coder_base_url)
    designer_model_kwargs = get_config_value(configurable.designer_model_kwargs or {})
    designer_model_name = get_config_value(configurable.designer_model)
    designer_provider = get_config_value(configurable.designer_provider)
    designer_base_url = get_config_value(configurable.designer_base_url)
    openai_key = os.getenv("OPENAI_API_KEY")
    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
        designer_model = init_chat_model(model=designer_model_name, model_provider=designer_provider, model_kwargs={})
    else:
        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        # temperature=0.7,
        # max_tokens=4096
        max_completion_tokens=40960,
        reasoning_effort = coder_model_kwargs["reasoning_effort"]
    )
    # 扩展幻灯片内容
    content_enrichment_prompt = f"""
    Based on the following points, expand the slide content. Provide a detailed description for each point and ensure it is linked to the search results.

    Slide topic: {topic}
    Slide section: {ppt_section.name}
    Slide title: {slide_title}
    Points: {', '.join(slide_points)}

    Search results:
    {source_str}
    Please expand and describe each point in detail, suitable for a presentation. Keep the language concise yet informative—ideally, each expanded point should not exceed **10 words**. The shorter, the better. Return the result in JSON format only; do not output any other text:

    {{
        "enriched_points": [
            {{"point_title": "Point Title 1", "expanded_content": "Expanded content 1"}},
            {{"point_title": "Point Title 2", "expanded_content": "Expanded content 2"}},
            {{"point_title": "Point Title 3", "expanded_content": "Expanded content 3"}},
            ...
        ]
    }}
    """


    enrichment_response = await writer_model.ainvoke([
        SystemMessage(content=content_enrichment_prompt),
        HumanMessage(content="Please expand the slide content. Only return Json, no additional text.")
    ])
    print("JSOOOOOOOOOOON")
    # enriched_points = json.loads(enrichment_response.content.replace("```json", "").replace("```", "").strip())["enriched_points"]
    enriched_points = enrichment_response.content

    save_dir = os.path.join(".", SAVES_ROOT, "outlines", topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    safe_section_name = ppt_section.name.replace(" ", "_")
    file_path = os.path.join(save_dir, f"{safe_section_name}_slide{slide_index+1}.json")

    data_to_save = {
        "query_results": source_str,
        "enriched_points": enriched_points,
    }

    # try:
    #     await asyncio.to_thread(
    #         lambda p=file_path, d=data_to_save: open(p, "w", encoding="utf-8").write(
    #             json.dumps(d, ensure_ascii=False, indent=2)
    #         )
    #     )
    #     # print(f"[INFO] Saved query & enriched points to: {file_path}")
    # except Exception as exc:
    #     print(f"[WARN] Error saving output file: {exc}")
    #     file_path = ""

    # 用于收集数据的返回值，不生成完整的幻灯片，只生成内容 TOOOOOOOOCHANGE
    # generated_slide = PPTSlide(
    #         title=slide_title,
    #         points=slide_points,
    #         codes=['no code'],
    #         enriched_points=json.dumps(enriched_points, ensure_ascii=False, indent=2),
    #         detail="no detail",
    #     )
    # return Command(
    #         update={"completed_slides": [generated_slide]},
    #         goto=END
    # )

    # 第二阶段：使用coder_model生成幻灯片布局描述
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    # writer_provider = get_config_value(configurable.writer_provider)
    # writer_model_name = get_config_value(configurable.writer_model)
    # writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    # coder_provider = get_config_value(configurable.coder_provider)
    # coder_model_name = get_config_value(configurable.coder_model)
    # coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    # if MODE == "openai":
    #     coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
    # else:
    #     coder_model = AzureChatOpenAI(
    #     model=configurable.coder_model,
    #     azure_endpoint=coder_model_kwargs["openai_api_base"],
    #     deployment_name=coder_model_kwargs["azure_deployment"],
    #     openai_api_version=coder_model_kwargs["openai_api_version"],
    #     # temperature=0.7,
    #     # max_tokens=4096
    #     max_completion_tokens=40960,
    #     reasoning_effort = coder_model_kwargs["reasoning_effort"]
    # )

    # OpenAI官方API兜底
    try:
        from langchain_openai import ChatOpenAI
        openai_model = ChatOpenAI(
            model=configurable.coder_model,
            api_key=os.getenv("OPENAI_API_KEY2"),
            bsae_url = os.getenv("OPENAI_API_BASE2"),
            max_tokens=40960,
        )
    except Exception:
        openai_model = None

    detail_prompt = f"""
    You are a seasoned slide designer responsible for designing slide layouts.
    Generate a slide layout description in JSON format based on the following details:

    Title: {slide_title}
    Detailed points: {enriched_points}
    Slide style: {style}
    Slide layout: {slide_layout}
    Primary color: {main_color}
    Accent color: {accent_color}
    Background tone: {background_tone}
    Heading font color: {heading_font_color}
    Body font color: {body_font_color}
    Font name: {font_name}
    Style summary: {style_summary}

    Here are some images you may optionally use in the PPT:
    <Image list>
    {images_json_embedded}
    </Image list>

    {design_formatting_prompt}

    <Suggestions for improving design>
    {design_suggestions}
    {aestheitcs_suggestions}

    """


    # [Ablation 1] ABLATE_DESIGN: skip designer model, use empty slide_detail
    if not ABLATE_DESIGN:
        try:
            detail_response = await asyncio.wait_for(
                designer_model.ainvoke([
                SystemMessage(content=detail_prompt),
                HumanMessage(content="Please output the JSON for the slide layout. Return JSON only, no additional text.")
                ]),
                timeout=300
            )
        except Exception as e:
            print(f"[WARN] Azure coder_model failed: {e}, fallback to OpenAI API.")
            if openai_model:
                try:
                    detail_response = await asyncio.wait_for(
                        openai_model.ainvoke([
                        SystemMessage(content=detail_prompt),
                        HumanMessage(content="Please output the JSON for the slide layout. Return JSON only, no additional text.")
                        ]),
                        timeout=300
                    )
                except Exception as e2:
                    print(f"[ERROR] Both Azure and OpenAI failed: {e2}")
                    slide_detail = ""
            else:
                slide_detail = ""
        else:
            slide_detail = detail_response.content
    else:
        print("[Ablation] ABLATE_DESIGN=True: skipping design step, slide_detail set to empty.")
        slide_detail = ""

    # [Ablation 2] ABLATE_SCORING / ABLATE_DESIGN: skip design quality scoring
    design_score = 0.0
    if not ABLATE_DESIGN and not ABLATE_SCORING:
        design_prompt = f"""
Requirements:
    Title: {slide_title}
    Detailed points: {enriched_points}
    Slide style: {style}
    Slide layout: {slide_layout}
    Primary color: {main_color}
    Accent color: {accent_color}
    Background tone: {background_tone}
    Heading font color: {heading_font_color}
    Body font color: {body_font_color}
    Font name: {font_name}
    Style summary: {style_summary}

Here are some images may optionally be used in the PPT:
    <Image list>
    {images_json_embedded}
    </Image list>

Agent's design:
{slide_detail}

{evaluation_design}
"""
        planner_model_kwargs = configurable.planner_model_kwargs or {}
        planner_provider = get_config_value(configurable.planner_provider)
        planner_model_name = get_config_value(configurable.planner_model)
        if MODE == "openai":
            planner_model = init_chat_model(model=planner_model_name, model_provider=planner_provider, model_kwargs={})
        else:
            planner_model = AzureChatOpenAI(
            model=configurable.planner_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],
            deployment_name=planner_model_kwargs["azure_deployment"],
            openai_api_version=planner_model_kwargs["openai_api_version"],
            max_tokens=4096
        )
        design_score_resp = await planner_model.ainvoke([
            SystemMessage(content=design_prompt),
            HumanMessage(content=f"Please evaluate its design quality.")
        ])

        design_score_value = design_score_resp.content.strip().split("```json")[-1].split("```")[0]
        try:
            design_score_json = json.loads(design_score_value)
            design_score = design_score_json.get("Total Score", 0)
            design_suggestions = design_score_json.get("Suggestions", "")
        except Exception as e:
            print(f"[WARN] Error parsing design score JSON: {e}")
    else:
        print(f"[Ablation] Skipping design scoring (ABLATE_DESIGN={ABLATE_DESIGN}, ABLATE_SCORING={ABLATE_SCORING}).")

    return {"enriched_points": enriched_points,
            "slide_detail": slide_detail,
            "design_score": design_score,
            "design_suggestions": design_suggestions,
            "image_data": images_json_embedded
            }


async def generate_slide_code_and_execute(state: PPTSlideState, config: RunnableConfig):
    enriched_points = state["enriched_points"]
    slide_detail = state["slide_detail"]
    images_json_embedded = state["image_data"]
    set_openai_api_base()
    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_base_url = get_config_value(configurable.coder_base_url)
    openai_key = os.getenv("OPENAI_API_KEY")
    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
    else:
        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960,
        reasoning_effort = coder_model_kwargs["reasoning_effort"]
    )

    # OpenAI官方API兜底
    try:
        from langchain_openai import ChatOpenAI
        openai_model = ChatOpenAI(
            model=configurable.coder_model,
            api_key=os.getenv("OPENAI_API_KEY2"),
            base_url = os.getenv("OPENAI_API_BASE2"),
            max_tokens=40960,
        )
    except Exception:
        openai_model = None

    topic = state["topic"]
    ppt_section = state["ppt_section"]
    slide_index = state["slide_index"]
    main_color = state.get("main_color")
    accent_color = state.get("accent_color")
    background_tone = state.get("background_tone")
    heading_font_color = state.get("heading_font_color")
    body_font_color = state.get("body_font_color")
    font_name = state.get("font_name")
    style_summary = state.get("style_summary", "")
    style = state.get("style")
    completeness_score = state.get("completeness_score", 0.0)
    slide_title = ppt_section.slides[slide_index].title
    slide_points = ppt_section.slides[slide_index].points

    # ✅ 使用 Path，避免 abspath/cwd；真正需要绝对路径时放到线程里 resolve
    save_dir = Path(SAVES_ROOT) / topic
    await asyncio.to_thread(save_dir.mkdir, parents=True, exist_ok=True)

    error_message = ""
    previous_code = ""
    python_code = None
    execution_successful = False

    async def run_script(script_path: Path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", str(script_path)],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    if COMP_MODE == "tools":
        path = f"{save_dir}/{ppt_section.name}_slide_{slide_index + 1}.pptx"
        try:
            # render_design_to_ppt(slide_detail, path=path)
            prs = await asyncio.to_thread(
                    render_design_to_ppt,
                    slide_detail,
                    path,    # 对应参数 path
                )
            assert await asyncio.to_thread(os.path.exists, path), "PPT file was not created."
            return {
                "codes": ["Rendered with tools"],
                "path": path,
                "title": slide_title,
                "points": slide_points,
                "completeness_score": 5.0,
            }
        except Exception as e:
            print(f"[ERROR] render_design_to_ppt failed: {e}")
            return {
                "codes": [],
                "path": "none",
                "title": slide_title,
                "points": slide_points,
                "completeness_score": 0.0,
                "completeness_suggestions": f"render_design_to_ppt failed: {e}"
            }


    for attempt in range(3):
        code_prompt = f"""
Generate Python code that creates slides using the python-pptx library based on the following detailed slide description:

Title: {slide_title}
Detailed bullet points: {enriched_points}
Slide description: {slide_detail}
Slide style: {style}
Primary color: {main_color} 
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}
Style summary: {style_summary}
{images_json_embedded}

{ppt_tools_prompt}

Code requirements:
1. Import the necessary libraries.
2. Create the slides and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
3. According to the detailed description, add the title, bullet points, and images at specified positions; set fonts and styles; explicitly set the size of each element to prevent overlap/occlusion; ensure text wraps automatically. The font size of the main text should be **at least 16**.
4. Only the provided image URLs can be used. Do not reserve any positions for any images that are not provided, and do not use text descriptions to fill the gaps. Or you can also manually create some flowcharts using various graphics, but don't just leave an empty space or just provide a textual description.
5. All the text should be placed on the top layer.
6. Save the file as: \"{save_dir}/{ppt_section.name}_slide_{slide_index + 1}.pptx\"

Previous code and errors (if any):
{previous_code}
{error_message}

Please provide complete, executable Python code based on this information. Note: output Python code only, do not output any other text.
Code will be save in utf-8 encoding.
        """

        try:
            code_response = await asyncio.wait_for(
                coder_model.ainvoke([
                SystemMessage(content=code_prompt),
                HumanMessage(content="Generate complete Python code.")
                ]),
                timeout=300
            )
            python_code = code_response.content.replace("```python", "").replace("```", "").strip()

        except Exception as e:
            print(f"[WARN] Azure coder_model code failed: {e}, fallback to OpenAI API.")
            if openai_model:
                try:
                    code_response = await asyncio.wait_for(
                        openai_model.ainvoke([
                        SystemMessage(content=code_prompt),
                        HumanMessage(content="Generate complete Python code.")
                        ]),
                    timeout=300
                    )
                    python_code = code_prefix + code_response.content.replace("```python", "").replace("```", "").strip()
                except Exception as e2:
                    print(f"[ERROR] Both Azure and OpenAI code failed: {e2}")
                    error_message = f"Both Azure and OpenAI code generation failed: {e2}"
                    continue
            else:
                error_message = f"Azure code generation failed: {e}, OpenAI API not available."
                continue

        # 保存生成的代码到文件
        script_path = save_dir / f"{ppt_section.name}_slide_{slide_index + 1}_attempt_{attempt}.py"
        try:
            await asyncio.to_thread(
                lambda: script_path.write_text(python_code, encoding="utf-8")
            )
            print(f"[INFO] Generated code saved to: {script_path}")
        except Exception as e:
            print(f"[ERROR] Failed to save code: {e}")
            error_message = f"Failed to save generated code: {e}"
            continue

        # 执行脚本
        print(f"[INFO] Executing script: {script_path}")
        returncode, stdout, stderr = await run_script(script_path)

        if returncode == 0:
            print(f"[INFO] Script executed successfully on attempt {attempt + 1}")
            execution_successful = True
            previous_code = python_code
            error_message = ""
            break
        else:
            error_message = f"Execution failed (attempt {attempt + 1}): {stderr or stdout}"
            print(f"[WARN] {error_message}")
            previous_code = python_code

    if not execution_successful:
        print(f"[ERROR] Failed to generate and execute code after 3 attempts")
        return {
            "codes": [python_code] if python_code else [],
            "path": "none",
            "title": slide_title,
            "points": slide_points,
            "completeness_score": 0.0,
        }

    # [Ablation 2] ABLATE_SCORING: skip completeness scoring
    completeness_score = 0.0
    completeness_suggestions = ""
    if not ABLATE_SCORING:
        planner_model_kwargs = configurable.planner_model_kwargs or {}
        planner_provider = get_config_value(configurable.planner_provider)
        planner_model_name = get_config_value(configurable.planner_model)
        if MODE == "openai":
            planner_model = init_chat_model(model=planner_model_name, model_provider=planner_provider, model_kwargs={})
        else:
            planner_model = AzureChatOpenAI(
            model=configurable.planner_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],
            deployment_name=planner_model_kwargs["azure_deployment"],
            openai_api_version=planner_model_kwargs["openai_api_version"],
            max_tokens=4096
        )
        completeness_prompt = f"""
# Design:
{slide_detail}

# Code:
{python_code}

# Task:
{evaluation_complete}"""
        completeness_response = await planner_model.ainvoke([
            SystemMessage(content=completeness_prompt),
            HumanMessage(content="Please evaluate the completeness of the generated slide.")
        ])
        completeness_value = completeness_response.content.strip().split("```json")[-1].split("```")[0]
        try:
            completeness_json = json.loads(completeness_value)
            completeness_score = completeness_json.get("Total Score", 0.0)
            completeness_suggestions = completeness_json.get("Suggestions", "")
        except Exception as e:
            print(f"[WARN] Error parsing completeness score JSON: {e}")
    else:
        print("[Ablation] ABLATE_SCORING=True: skipping completeness scoring.")
        completeness_score = 5.0

    pptx_path = save_dir / f"{ppt_section.name}_slide_{slide_index + 1}.pptx"
    # ✅ 绝对路径解析（会触发 cwd），放在线程里执行以避免阻塞
    pptx_abs = await asyncio.to_thread(lambda: str(pptx_path.resolve()))

    return {
        "codes": [python_code],
        "path": pptx_abs,
        "title": slide_title,
        "points": slide_points,
        "completeness_score": completeness_score,
        "completeness_suggestions": completeness_suggestions
    }


# def ppt_to_image(slide_ppt_path, image_path):
#     """
#     使用 unoconv 将PPT幻灯片导出为图片
#     """
#     print(f"Convert the slide {slide_ppt_path} to an image {image_path}")
    
#     # 使用 unoconv 将 PPT 转换为 PNG 格式
#     # try:
#         # 使用 unoconv 命令行工具来将 ppt 文件转换为图片
#     command = [
#             "C:\\Windows\\unoconv.bat",  # 调用 unoconv 命令
#             "-f", "png",  # 转换为 png 格式
#             "-o", image_path,  # 输出路径
#             slide_ppt_path  # 输入的 PPT 文件路径
#     ]
        
#         # 调用 unoconv 命令
#     subprocess.run(command, check=True)
#     print(f"Conversion successful: {slide_ppt_path} -> {image_path}")
    
#     # except subprocess.CalledProcessError as e:
#     #     print(f"转换失败：{str(e)}")

def ppt_to_image(slide_ppt_path, image_path, soffice_path=None, timeout=120):
    """
    在 macOS (conda 环境) 使用 LibreOffice(soffice) 将 PPT/PPTX 导出为 PNG。
    - 如果 PPT 仅 1 页 -> 生成 image_path 指定的 PNG。
    - 如果 PPT 多页 -> 第 1 页用 image_path，其余页在同目录生成 *_slide_02.png、*_slide_03.png...
    返回：保存的图片绝对路径列表（str）。
    """
    slide_ppt_path = Path(slide_ppt_path).expanduser().resolve()
    image_path = Path(image_path).expanduser()
    image_dir = image_path.parent
    image_stem = image_path.stem
    image_dir.mkdir(parents=True, exist_ok=True)

    # 优先级：函数参数 > 环境变量 > PATH 中的 soffice > 默认 App 路径
    candidates = []
    if soffice_path:
        candidates.append(Path(soffice_path))
    if os.environ.get("SOFFICE_PATH"):
        candidates.append(Path(os.environ["SOFFICE_PATH"]))
    which = shutil.which("soffice")
    if which:
        candidates.append(Path(which))
    candidates.append(Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"))

    soffice_bin = next((p for p in candidates if p and p.exists()), None)
    if soffice_bin is None:
        raise FileNotFoundError(
            "未找到 LibreOffice 的 'soffice' 可执行文件。请先安装 LibreOffice，"
            "或设置环境变量 SOFFICE_PATH 指向 soffice。"
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        cmd = [
            str(soffice_bin),
            "--headless", "--norestore", "--invisible",
            "--nodefault", "--nofirststartwizard", "--nolockcheck",
            "--convert-to", "png:impress_png_Export",
            "--outdir", str(tmpdir),
            str(slide_ppt_path),
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice 转换失败：{proc.stderr or proc.stdout}")

        generated = sorted(tmpdir.glob("*.png"))
        if not generated:
            raise FileNotFoundError("转换后未在临时目录发现 PNG 文件。")

        saved = []
        if len(generated) == 1:
            dst = image_path
            shutil.move(str(generated[0]), str(dst))
            saved.append(str(dst.resolve()))
        else:
            # 多页：第 1 页用给定文件名，其他页自动编号
            first_dst = image_path
            shutil.move(str(generated[0]), str(first_dst))
            saved.append(str(first_dst.resolve()))
            for i, src in enumerate(generated[1:], start=2):
                dst = image_dir / f"{image_stem}_slide_{i:02d}.png"
                shutil.move(str(src), str(dst))
                saved.append(str(dst.resolve()))

    return saved

async def ppt_slide_to_image_and_validate(state, config):
    """
    将生成的PPT幻灯片转换为图片，并使用大模型按新评分标准检查布局合理性。
    评分阈值：Total Score 严格 > 60 通过，否则 retry；最多重试 max_retry_count 次。
    """
    SCORE_THRESHOLD = 5  # 严格 > 60 才通过

    # ---- 新评分 Prompt ----
    REVIEW_PROMPT = f"""
{evaluation_aesthetics}
"""

    # ---- 配置 & 模型初始化 ----
    configurable = Configuration.from_runnable_config(config)
    planner_model_kwargs = configurable.planner_model_kwargs or {}
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})
    if MODE == "openai":
        planner_model = init_chat_model(model=planner_model, model_provider=planner_provider, model_kwargs={})
    else:

        planner_model = AzureChatOpenAI(
        model=configurable.planner_model,
        azure_endpoint=planner_model_kwargs["openai_api_base"],
        deployment_name=planner_model_kwargs["azure_deployment"],
        openai_api_version=planner_model_kwargs["openai_api_version"],
        temperature=0,
        max_tokens=2048
    )

    # ---- 读取状态 ----
    slide_ppt_path = state["path"]
    codes = state["codes"]
    title = state["title"]
    points = state["points"]
    enriched_points = state["enriched_points"]
    slide_detail = state["slide_detail"]
    layout = state.get("layout", "")
    max_retry_count = state.get("max_retry_count", 3)
    retry_count = state.get("retry_count", 0)
    path = state.get("path")
    design_score = state.get("design_score", 0.0)
    completeness_score = state.get("completeness_score", 0.0)

    print("Current Slide:", slide_ppt_path, " Current repetition count:", retry_count, "Max Retry:", max_retry_count)

    # ---- 达到最大重试次数 ----
    if retry_count >= max_retry_count:
        print("The maximum retry count has been reached. No further processing will be carried out.")
        generated_slide = PPTSlide(
            title=title,
            points=points,
            codes=codes,
            enriched_points=enriched_points,
            detail=slide_detail,
            layout=layout
        )
        return Command(update={"completed_slides": [generated_slide]}, goto=END)


    # ---- 无有效路径 ----
    if path == "none":
        return Command(update={"layout_valid": False, "retry_count": retry_count + 1}, goto="enrich_slide_content")

    # ---- 转换为图片 ----
    output_folder = os.path.dirname(slide_ppt_path)
    image_path = slide_ppt_path.replace(".pptx", ".png")
    print(f"Convert the slide {slide_ppt_path} to an image {image_path}")
    try:
        await asyncio.to_thread(ppt_to_image, slide_ppt_path, image_path)
    except Exception as e:
        print(f"Error converting PPT to image: {e}")
        return Command(update={"conversion_failed": True, "retry_count": retry_count + 1}, goto="enrich_slide_content")
    print(f"The slide transition to image was successful. The saved path is: {image_path}")

    # [Ablation 2] ABLATE_SCORING: skip image review and retry, accept first result
    if ABLATE_SCORING:
        print("[Ablation] ABLATE_SCORING=True: skipping aesthetics scoring, accepting slide as-is.")
        generated_slide = PPTSlide(
            title=title,
            points=points,
            codes=codes,
            enriched_points=enriched_points,
            detail=slide_detail,
            layout=layout
        )
        return Command(update={"completed_slides": [generated_slide]}, goto=END)

    # ---- 审查评分：>60 通过，否则 retry ----
    def _extract_json_dict(text: str):
        """从模型返回中提取JSON（兼容 ```json 代码块``` 或纯JSON）；失败时尽量解析 Total Score。"""
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
    response = await planner_model.ainvoke([
        SystemMessage(content=REVIEW_PROMPT),
        HumanMessage(content=[{
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
        }])
    ])
    print(f"[Review raw]: {response.content}")

    result_dict = _extract_json_dict(response.content or "")
    total_score_str = str(result_dict.get("Total Score", "-1"))
    try:
        total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
    except Exception:
        total_score = -1.0
    print(f"[Parsed total score]: {total_score}")
    aestheitcs_suggestions = ""
    if "Suggestions" in result_dict:
        aestheitcs_suggestions = result_dict['Suggestions']
        print(f"[Suggestions]: {result_dict['Suggestions']}")
    a=total_score
    total_score = design_score + (completeness_score-1)/4 * (total_score-3)
    print(f"[Final total score]: {total_score} (Design: {design_score} ; Completeness: {completeness_score} ; Aesthetics: {a})")
    if a>=3.8 and total_score > SCORE_THRESHOLD:
        generated_slide = PPTSlide(
            title=title,
            points=points,
            codes=codes,
            enriched_points=enriched_points,
            detail=slide_detail,
            layout=layout
        )
        return Command(update={"completed_slides": [generated_slide]}, goto=END)
    else:
        retry_count += 1
        if retry_count >= max_retry_count:
            generated_slide = PPTSlide(
                title=title,
                points=points,
                codes=codes,
                enriched_points=enriched_points,
                detail=slide_detail,
                layout=layout
            )
            return Command(update={"completed_slides": [generated_slide]}, goto=END)
        return Command(update={"layout_valid": False, "retry_count": retry_count, "aestheitcs_suggestions": aestheitcs_suggestions}, goto="enrich_slide_content")




async def generate_ppt_section_start(state: PPTSectionState):
    """
    为每个PPT章节开始生成幻灯片时，初始化一个空的幻灯片列表，并为每页PPT启动一个子图。

    Args:
        state: 包含章节信息、PPT章节信息的状态

    Returns:
        Dict: 包含空的 `generated_slides` 列表，并进入 `ppt_slide_subgraph`
    """
    # 从状态中获取章节信息和PPT章节信息
    topic = state["topic"]  # 主题
    # section = state["section"]  # 章节信息
    ppt_section = state["ppt_section"]  # PPT章节信息
    style = state.get("style")
    main_color = state.get("main_color")  # 主色
    accent_color = state.get("accent_color")  # 辅助
    background_tone = state.get("background_tone")
    heading_font_color = state.get("heading_font_color")
    body_font_color = state.get("body_font_color")
    font_name = state.get("font_name")
    style_summary = state.get("style_summary", "")

    # 初始化幻灯片列表，准备开始生成幻灯片
    generated_slides = []

    # 打印相关信息（用于调试）
    # print(f"开始生成PPT章节：{ppt_section.name}")

    # 获取该章节需要的页数
    num_slides = ppt_section.allocated_slides  # 获取分配的页数
    
    # 为每一页PPT启动一个子图
    return Command(
        update={"generated_slides": generated_slides},
        goto=[
            Send("generate_slide", {
                "topic": topic, 
                # "section": section, 
                "style": style,
                "ppt_section": ppt_section, 
                "slide_index": slide_index,  # 为每一页传递幻灯片的索引
                "main_color": main_color,
                "accent_color": accent_color,
                "background_tone": background_tone,
                "heading_font_color": heading_font_color,
                "body_font_color": body_font_color,
                "font_name": font_name,
                "style_summary": style_summary,
                "max_retry_count": 3,  # 设置最大重试次数
                "retry_count": 0  # 初始化重试计数
            })
            for slide_index in range(num_slides)  # 根据分配的页数启动相应数量的子图
        ]
    )


async def generate_ppt_section_end(state: PPTSectionState):
    ppt_section = state["ppt_section"]
    ppt_section.slides = state["completed_slides"]

    return Command(
        update={"completed_ppt_sections": [ppt_section]},
        goto=END
    )

async def generate_cover_slide(state, config):
    """生成封面幻灯片，包含布局评分检查，最多循环3次（Total Score > 60 通过）"""
    # ----------------- helper: extract JSON dict from LLM output -----------------
    def _extract_json_dict(text: str) -> Dict[str, Any]:
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    # ----------------- review prompt -----------------

    SCORE_THRESHOLD = 75  # strictly > 60 to pass

    # ----------------- state -----------------
    topic = state["topic"]
    style = state.get("style", "none")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")
    background_tone = state.get("background_tone", "Light")
    heading_font_color = state.get("heading_font_color", "#000000")
    body_font_color = state.get("body_font_color", "#000000")
    font_name = state.get("font_name", "Arial")

    save_dir = os.path.join(".", SAVES_ROOT, topic)
    cover_path = os.path.join(save_dir, "cover_slide.pptx")
    script_path = os.path.join(save_dir, "cover_slide.py")

    # ----------------- model init -----------------
    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_base_url = get_config_value(configurable.coder_base_url)
    designer_model_kwargs = configurable.designer_model_kwargs or {}
    designer_provider = get_config_value(configurable.designer_provider)
    designer_model_name = get_config_value(configurable.designer_model)
    designer_base_url = get_config_value(configurable.designer_base_url)
    planner_model_kwargs = configurable.planner_model_kwargs or {}
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model_name = get_config_value(configurable.planner_model)
    openai_key = os.getenv("OPENAI_API_KEY")
    set_openai_api_base()

    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider,  model_kwargs={})
        designer_model = init_chat_model(model=designer_model_name, model_provider=designer_provider, model_kwargs={})
        planner_model = init_chat_model(model=planner_model_name, model_provider=planner_provider, model_kwargs={})
    else:

        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960,
    )
    # OpenAI官方API兜底
    try:
        from langchain_openai import ChatOpenAI
        openai_model = ChatOpenAI(
            model=configurable.coder_model,
            api_key=os.getenv("OPENAI_API_KEY2"),
            base_url = os.getenv("OPENAI_API_BASE2"),
            max_tokens=40960,
        )
    except Exception:
        openai_model = None

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=600
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"Timeout exception: {str(e)}"
            except Exception as e:
                return -1, "", f"Exception: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    style_plan = f"""
You are a professional slide style designer. Your task is to create a comprehensive **style template** for an entire slide deck.

Below is the basic information about the slides:
Topic: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}

{style_plan_prompt}
"""
    # planner_model_kwargs = get_config_value(Configuration.from_runnable_config(config).planner_model_kwargs or {})
    # planner_model = AzureChatOpenAI(
    #     model=Configuration.from_runnable_config(config).planner_model,
    #     azure_endpoint=planner_model_kwargs["openai_api_base"],
    #     deployment_name=planner_model_kwargs["azure_deployment"],
    #     openai_api_version=planner_model_kwargs["openai_api_version"],
    #     max_tokens=4096
    # )
    style_response = await designer_model.ainvoke([
        SystemMessage(content=style_plan),
        HumanMessage(content="Please create a detailed slide style description.")
    ])

    style_summary = style_response.content.strip()
    suggestions = ""

    for attempt in range(3):
        # ----------------- ask model for layout description -----------------
        layout_prompt = f"""
Please design a slide cover layout with the title: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}
Style summary: {style_summary}

The slide cover should include the title, speaker name, and date as key information. 
You can design some background patterns and borders for the page to enhance its appearance, making it more **visually appealing** and sophisticated.
- Canvas: 13.33 × 7.5 inches (width × height)
- All coordinates and sizes are in **inches**, with 2 decimal places
- No elements may go out of bounds or overlap (except background textures/separators)
- In the Layout layer, every block uses **absolute positioning**: top-left `(x, y)`, size `(w, h)`;  
  `0 ≤ x ≤ 13.33`, `0 ≤ y ≤ 7.5`, `x + w ≤ 13.33`, `y + h ≤ 7.5`
- **If there is no images provided, do not reference images in the design.**
- Provide the layout description in JSON format.

{suggestions}
{f"Previous error message: {error_message}" if error_message else ""}
{f"Previously generated code: {previous_code}" if previous_code else ""}
        """
        try:
            layout_response = await asyncio.wait_for(
                designer_model.ainvoke([
                    {"role": "system", "content": layout_prompt},
                    {"role": "user", "content": "Generate a layout description"}
                ]),
                timeout=300
            )
            layout_description = layout_response.content.strip()
        except Exception as e:
            print(f"[WARN] Azure coder_model failed: {e}, fallback to OpenAI API.")
            if openai_model:
                try:
                    layout_response = await asyncio.wait_for(
                        openai_model.ainvoke([
                            {"role": "system", "content": layout_prompt},
                            {"role": "user", "content": "Generate a layout description"}
                        ]),
                        timeout=300
                    )
                    layout_description = layout_response.content.strip()
                except Exception as e2:
                    error_message = f"Both Azure and OpenAI layout failed: {e2}"
                    print(f"[ERROR] Both Azure and OpenAI layout failed: {e2}")
                    continue
            else:
                error_message = f"Azure layout failed: {e}, OpenAI API not available."
                continue
        
        
        if COMP_MODE == "tools":
            try:
                print("Cover comp start")
                # render_design_to_ppt(layout_description, path=cover_path)
                python_code = "None"
                prs = await asyncio.to_thread(
                    render_design_to_ppt,
                    layout_description,
                    cover_path,    # 对应参数 path
                )
                assert await asyncio.to_thread(os.path.exists, cover_path), "PPT file was not created."
                # goto .convert
                # return {
                #     "cover_slide_path": cover_path, "cover_layout_description": layout_description, "style_summary": style_summary
                # }
            except Exception as e:
                error_message = f"render_design_to_ppt failed: {e}"
                print(e)
                continue
        # ----------------- ask model to generate python-pptx code -----------------
        else:
            code_generation_prompt = f"""
Generate complete Python code using the python-pptx library to create a cover slide based on the following layout description:

Layout description: {layout_description}
{ppt_tools_prompt}
Code requirements:
1. Import the necessary libraries.
2. Create the slide and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
4. All the text should be placed on the top layer.
5. Save the file to: {cover_path}

{f"Previous error message: {error_message}" if error_message else ""}
{f"Previously generated code: {previous_code}" if previous_code else ""}

Please provide complete, executable Python code based on this information. Note: output Python code only, do not output any other text.
Code will be save in utf-8 encoding.
        """
            try:
                code_response = await asyncio.wait_for(
                    coder_model.ainvoke([
                        {"role": "system", "content": code_generation_prompt},
                        {"role": "user", "content": "Generate Python code"}
                    ]),
                    timeout=300
                )
                python_code = code_prefix + code_response.content.replace("```python", "").replace("```", "").strip()
            except Exception as e:
                print(f"[WARN] Azure coder_model code failed: {e}, fallback to OpenAI API.")
                if openai_model:
                    try:
                        code_response = await asyncio.wait_for(
                            openai_model.ainvoke([
                                {"role": "system", "content": code_generation_prompt},
                                {"role": "user", "content": "Generate Python code"}
                            ]),
                            timeout=300
                        )
                        python_code = code_response.content.replace("```python", "").replace("```", "").strip()
                    except Exception as e2:
                        error_message = f"Both Azure and OpenAI code failed: {e2}"
                        print(f"[ERROR] Both Azure and OpenAI code failed: {e2}")
                        continue
                else:
                    error_message = f"Azure code failed: {e}, OpenAI API not available."
                    continue

            await asyncio.to_thread(
                lambda: open(script_path, "w", encoding="utf-8").write(python_code)
            )

            # ----------------- run script -----------------
            returncode, stdout, stderr = await run_script(script_path)
            if returncode != 0:
                error_message = stderr
                previous_code = python_code
                print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
                continue
            
        
        # ----------------- convert to image -----------------
        image_path = cover_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, cover_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        # [Ablation 2] ABLATE_SCORING: skip review scoring, accept first successful result
        if ABLATE_SCORING:
            print("[Ablation] ABLATE_SCORING=True: skipping cover slide scoring.")
            return {"cover_slide_path": cover_path, "cover_layout_description": layout_description, "style_summary": style_summary}

        # ----------------- review with scoring -----------------
        REVIEW_PROMPT = f"""
Topic: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}
Style summary: {style_summary}

You are a Slide Review Expert. Please evaluate the cover slide design based on the following dimensions.
{eval_cover}
"""
        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        try:
            review = await asyncio.wait_for(
                planner_model.ainvoke([
                    SystemMessage(content=REVIEW_PROMPT),
                    HumanMessage(content=[{
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
                    }])
                ]),
                timeout=300
            )
        except Exception as e:
            print(f"[WARN] Azure coder_model review failed: {e}, fallback to OpenAI API.")
            previous_code = python_code
            if openai_model:
                try:
                    review = await asyncio.wait_for(
                        openai_model.ainvoke([
                            SystemMessage(content=REVIEW_PROMPT),
                            HumanMessage(content=[{
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
                            }])
                        ]),
                        timeout=300
                    )
                except Exception as e2:
                    error_message = f"Both Azure and OpenAI review failed: {e2}"
                    print(f"[ERROR] Both Azure and OpenAI review failed: {e2}")
                    continue
            else:
                error_message = f"Azure review failed: {e}, OpenAI API not available."
                continue

        print(f"[Cover Review raw] {review.content}")

        result = _extract_json_dict(review.content or "")
        total_score_str = str(result.get("Total Score", "-1"))
        try:
            total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
        except Exception:
            total_score = -1.0

        print(f"[Parsed score] {total_score} / 100")
        if "Suggestions" in result:
            suggestions = result['Suggestions']
            print(f"[Suggestions] {result['Suggestions']}")

        if total_score > SCORE_THRESHOLD:
            print(f"封面布局评分通过（{total_score}），路径：{cover_path}")
            return {"cover_slide_path": cover_path, "cover_layout_description": layout_description, "style_summary": style_summary}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"封面布局评分未通过，尝试重试：{attempt + 1}")

    # 最终返回（达到最大次数仍未通过）
    return {"cover_slide_path": cover_path, "cover_layout_description": layout_description, "style_summary": style_summary}



async def generate_section_cover_slides(state, config):
    """生成章节封面幻灯片，仅检查第一个章节的布局有效性，最多循环3次（Total Score > 60 通过）"""
    SCORE_THRESHOLD = 75  # strictly > 60 to pass

    def _extract_json_dict(text: str) -> Dict[str, Any]:
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}



    # ---- state ----
    topic = state["topic"]
    ppt_sections = state["ppt_sections"]
    style = state.get("style", "none")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")
    background_tone = state.get("background_tone", "Light")
    heading_font_color = state.get("heading_font_color", "#000000")
    body_font_color = state.get("body_font_color", "#000000")
    font_name = state.get("font_name", "Arial")
    style_summary = state.get("style_summary", "")
    suggesstions = ""

    save_dir = os.path.join(".", SAVES_ROOT, topic)
    # os.makedirs(save_dir, exist_ok=True)
    script_path = os.path.join(save_dir, "section_cover_slide.py")

    # ---- model init ----
    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_base_url = get_config_value(configurable.coder_base_url)
    designer_model_kwargs = configurable.designer_model_kwargs or {}
    designer_provider = get_config_value(configurable.designer_provider)
    designer_model_name = get_config_value(configurable.designer_model)
    designer_base_url = get_config_value(configurable.designer_base_url)
    planner_model_kwargs = configurable.planner_model_kwargs or {}
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model_name = get_config_value(configurable.planner_model)
    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
        designer_model = init_chat_model(model=designer_model_name, model_provider=designer_provider, model_kwargs={})
        planner_model = init_chat_model(model=planner_model_name, model_provider=planner_provider, model_kwargs={})
    else:

        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960,
        reasoning_effort = coder_model_kwargs["reasoning_effort"]
    )
    openai_model = None
    # OpenAI官方API兜底
    try:
        from langchain_openai import ChatOpenAI
        openai_model = ChatOpenAI(
            model=configurable.coder_model,
            api_key=os.getenv("OPENAI_API_KEY2"),
            base_url = os.getenv("OPENAI_API_BASE2"),
            max_tokens=40960,
        )
    except Exception:
        openai_model = None

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=120
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    chapters_list = [section.name for section in ppt_sections]

    for attempt in range(3):
        # ---- ask for a generic layout for all chapter covers ----
        layout_prompt = f"""
    Please design a universal layout for slide chapter covers, applicable to all chapters. You can design some background patterns and borders for the page to enhance its appearance, making it more **visually appealing** and sophisticated.
    Slide topic: {topic}
    Slide style: {style}
    Primary color: {main_color}
    Accent color: {accent_color}
    Background tone: {background_tone}
    Heading font color: {heading_font_color}
    Body font color: {body_font_color}
    Font name: {font_name}
    Style summary: {style_summary}
- Canvas: 13.33 × 7.5 inches (width × height)
- All coordinates and sizes are in **inches**, with 2 decimal places
- No elements may go out of bounds or overlap (except background textures/separators)
- In the Layout layer, every block uses **absolute positioning**: top-left `(x, y)`, size `(w, h)`;  
  `0 ≤ x ≤ 13.33`, `0 ≤ y ≤ 7.5`, `x + w ≤ 13.33`, `y + h ≤ 7.5`
- **If there is no images provided, do not reference images in the design.**

Chapter list: {chapters_list}
For one section cover page, there is one Chapter name.
To ensure the integrity of the overall structure, you can initially use "{chapters_list[0]}" as the placeholder for the chapter title content.


{suggesstions}

- Provide the layout description in JSON format.

    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}
        """
        try:
            layout_response = await designer_model.ainvoke([
                {"role": "system", "content": layout_prompt},
                {"role": "user", "content": "Generate a layout description"}
            ])
            layout_description = layout_response.content.strip()
        except Exception as e:
            print(f"[WARN] Azure coder_model layout failed: {e}, fallback to OpenAI API.")
            if openai_model:
                try:
                    layout_response = await openai_model.ainvoke([
                        {"role": "system", "content": layout_prompt},
                        {"role": "user", "content": "Generate a layout description"}
                    ])
                    layout_description = layout_response.content.strip()
                except Exception as e2:
                    error_message = f"Both Azure and OpenAI layout failed: {e2}"
                    print(f"[ERROR] Both Azure and OpenAI layout failed: {e2}")
                    continue
            else:
                error_message = f"Azure layout failed: {e}, OpenAI API not available."
                continue
        if COMP_MODE == "tools":
            try:
                # render_design_to_ppt(layout_description, path=f"{save_dir}/section_slide_1.pptx")
                prs = await asyncio.to_thread(
                    render_design_to_ppt,
                    layout_description,
                    f"{save_dir}/section_slide_1.pptx",    # 对应参数 path
                )
                assert await asyncio.to_thread(os.path.exists, f"{save_dir}/section_slide_1.pptx"), "PPT file was not created."
                for i in range(len(chapters_list)-1):
                    layout_description.replace(chapters_list[i],chapters_list[i+1])
                    # render_design_to_ppt(layout_description, path=f"{save_dir}/section_slide_{i+2}.pptx")
                    prs = await asyncio.to_thread(
                        render_design_to_ppt,
                        layout_description,
                        f"{save_dir}/section_slide_{i+2}.pptx",    # 对应参数 path
                    )
                # goto .convert_section
            except Exception as e:
                error_message = f"render_design_to_ppt failed: {e}"
                continue
        # ---- generate python script to create all chapter cover PPTX files ----
        else:
            code_generation_prompt = f"""
    Generate a Python script using the python-pptx library to create chapter cover slides based on the following layout description:

    Layout description: {layout_description}
    {ppt_tools_prompt}
    Script requirements:
    1. Import the necessary libraries.
    2. Create slides and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
    3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
    4. All the text should be placed on the top layer.
    5. Using the chapter list below, generate a separate PPTX file for each chapter. Save each file to {save_dir}/section_slide_{{chapter_index}}.pptx.

    Chapter list: {chapters_list}

    Please provide complete, executable Python code based on this information. Note: output Python code only; do not output any other text. Code will be save in utf-8 encoding.
        """
            try:
                code_response = await coder_model.ainvoke([
                    {"role": "system", "content": code_generation_prompt},
                    {"role": "user", "content": "Generate Python code"}
                ])
                python_code = code_response.content.replace("```python", "").replace("```", "").strip()
            except Exception as e:
                print(f"[WARN] Azure coder_model code failed: {e}, fallback to OpenAI API.")
                if openai_model:
                    try:
                        code_response = await openai_model.ainvoke([
                            {"role": "system", "content": code_generation_prompt},
                            {"role": "user", "content": "Generate Python code"}
                        ])
                        python_code = code_prefix + code_response.content.replace("```python", "").replace("```", "").strip()
                    except Exception as e2:
                        error_message = f"Both Azure and OpenAI code failed: {e2}"
                        print(f"[ERROR] Both Azure and OpenAI code failed: {e2}")
                        continue
                else:
                    error_message = f"Azure code failed: {e}, OpenAI API not available."
                    continue
            
            await asyncio.to_thread(lambda: open(script_path, "w", encoding="utf-8").write(python_code))

            # ---- run script to generate PPTs ----
            returncode, stdout, stderr = await run_script(script_path)
            if returncode != 0:
                error_message = stderr
                previous_code = python_code
                print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
                continue
            
        # label .convert_section
        # ---- validate ONLY the first chapter's PPT by scoring ----
        first_slide_path = os.path.join(save_dir, "section_slide_1.pptx")
        image_path = first_slide_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, first_slide_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        # [Ablation 2] ABLATE_SCORING: skip review scoring, accept first successful result
        if ABLATE_SCORING:
            print("[Ablation] ABLATE_SCORING=True: skipping section cover slide scoring.")
            return {"section_slides_path": save_dir, "layout_description": layout_description}
        REVIEW_PROMPT = f"""
Topic: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}
Style summary: {style_summary}
Chapter list: {chapters_list}

You are a Slide Review Expert. Please evaluate the section cover slide design for the first chapter based on the following dimensions.

{eval_cover}
"""
        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        try:
            review = await planner_model.ainvoke([
                SystemMessage(content=REVIEW_PROMPT),
                HumanMessage(content=[{
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
                }])
            ])
        except Exception as e:
            print(f"[WARN] Azure coder_model review failed: {e}, fallback to OpenAI API.")
            previous_code = python_code
            if openai_model:
                try:
                    review = await openai_model.ainvoke([
                        SystemMessage(content=REVIEW_PROMPT),
                        HumanMessage(content=[{
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
                        }])
                    ])
                except Exception as e2:
                    error_message = f"Both Azure and OpenAI review failed: {e2}"
                    print(f"[ERROR] Both Azure and OpenAI review failed: {e2}")
                    continue
            else:
                error_message = f"Azure review failed: {e}, OpenAI API not available."
                continue

        print(f"[Section Cover Review raw] {review.content}")

        result = _extract_json_dict(review.content or "")
        total_score_str = str(result.get("Total Score", "-1"))
        try:
            total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
        except Exception:
            total_score = -1.0

        print(f"[Parsed score] {total_score} / 100")
        if "Suggestions for Improvement" in result:
            suggesstions = result['Suggestions for Improvement']
            print(f"[Suggestions] {result['Suggestions for Improvement']}")

        if total_score > SCORE_THRESHOLD:
            print(f"第一个章节布局评分通过（{total_score}），目录：{save_dir}")
            return {"section_slides_path": save_dir, "layout_description": layout_description}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"第一个章节布局评分未通过，尝试重试：{attempt + 1}")

    # 最终返回（达到最大次数仍未通过）
    return {"section_slides_path": save_dir, "layout_description": layout_description}



async def generate_end_slide(state, config):
    """生成封底幻灯片，包含布局检查步骤，最多循环3次（评分>60通过，否则重试）"""
    SCORE_THRESHOLD = 75  # 严格“>60”才通过

    def _extract_json_dict(text: str) -> Dict[str, Any]:
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    topic = state["topic"]
    style = state.get("style", "专业商务")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")
    background_tone = state.get("background_tone", "Light")
    heading_font_color = state.get("heading_font_color", "#000000")
    body_font_color = state.get("body_font_color", "#000000")
    font_name = state.get("font_name", "Arial")
    style_summary = state.get("style_summary", "")
    suggestions = ""

    save_dir = os.path.join(".", SAVES_ROOT, topic)
    # os.makedirs(save_dir, exist_ok=True)
    end_path = os.path.join(save_dir, "end_slide.pptx")
    script_path = os.path.join(save_dir, "end_slide.py")

    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    coder_base_url = get_config_value(configurable.coder_base_url)
    designer_model_kwargs = configurable.designer_model_kwargs or {}
    designer_provider = get_config_value(configurable.designer_provider)
    designer_model_name = get_config_value(configurable.designer_model)
    designer_base_url = get_config_value(configurable.designer_base_url)
    planner_model_kwargs = configurable.planner_model_kwargs or {}
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model_name = get_config_value(configurable.planner_model)
    openai_key = os.getenv("OPENAI_API_KEY")
    if MODE == "openai":
        coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs={})
        designer_model = init_chat_model(model=designer_model_name, model_provider=designer_provider, model_kwargs={})
        planner_model = init_chat_model(model=planner_model_name, model_provider=planner_provider, model_kwargs={})
    else:

        coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960,
        reasoning_effort = coder_model_kwargs["reasoning_effort"]
    )

    # OpenAI官方API兜底
    try:
        from langchain_openai import ChatOpenAI
        openai_model = ChatOpenAI(
            model=configurable.coder_model,
            api_key=os.getenv("OPENAI_API_KEY2"),
            base_url = os.getenv("OPENAI_API_BASE2"),
            max_tokens=40960,
        )
    except Exception:
        openai_model = None

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    for attempt in range(3):
        layout_prompt = f"""
    Please design a layout for the slide's end (closing) page.
    Slide topic: {topic}
    Slide style: {style}
    Primary color: {main_color}
    Accent color: {accent_color}
    Background tone: {background_tone}
    Heading font color: {heading_font_color}
    Body font color: {body_font_color}
    Font name: {font_name}
    Style summary: {style_summary}

    {suggestions}
    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}
        """

        try:
            layout_response = await designer_model.ainvoke([
                {"role": "system", "content": layout_prompt},
                {"role": "user", "content": "Generate a layout description"}
            ])
            layout_description = layout_response.content.strip()
        except Exception as e:
            print(f"[WARN] Azure coder_model layout failed: {e}, fallback to OpenAI API.")
            if openai_model:
                try:
                    layout_response = await openai_model.ainvoke([
                        {"role": "system", "content": layout_prompt},
                        {"role": "user", "content": "Generate a layout description"}
                    ])
                    layout_description = layout_response.content.strip()
                except Exception as e2:
                    error_message = f"Both Azure and OpenAI layout failed: {e2}"
                    print(f"[ERROR] Both Azure and OpenAI layout failed: {e2}")
                    continue
            else:
                error_message = f"Azure layout failed: {e}, OpenAI API not available."
                continue
        
        if COMP_MODE == "tools":
            try:
                # render_design_to_ppt(layout_description, path=end_path)
                prs = await asyncio.to_thread(
                    render_design_to_ppt,
                    layout_description,
                    end_path,    # 对应参数 path
                )
                assert await asyncio.to_thread(os.path.exists, end_path), "PPT file was not created."
                # goto .convert_end
                # return {
                #     "cover_slide_path": cover_path, "cover_layout_description": layout_description, "style_summary": style_summary
                # }
            except Exception as e:
                error_message = f"render_design_to_ppt failed: {e}"
                continue
        else:
            code_generation_prompt = f"""
    Generate complete Python code using the python-pptx library to create the closing slide based on the following layout description:

    Layout description: {layout_description}
    {ppt_tools_prompt}
    Code requirements:
    1. Import the necessary libraries.
    2. Create the slide and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
    3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
    4. All the text should be placed on the top layer.
    5. Save the file to: {end_path}

    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}

    Please provide complete, executable Python code based on this information. Note: output Python code only, do not output any other text.Code will be save in utf-8 encoding.
        """

            try:
                code_response = await coder_model.ainvoke([
                    {"role": "system", "content": code_generation_prompt},
                    {"role": "user", "content": "Generate Python code"}
                ])
                python_code = code_response.content.replace("```python", "").replace("```", "").strip()
            except Exception as e:
                print(f"[WARN] Azure coder_model code failed: {e}, fallback to OpenAI API.")
                if openai_model:
                    try:
                        code_response = await openai_model.ainvoke([
                            {"role": "system", "content": code_generation_prompt},
                            {"role": "user", "content": "Generate Python code"}
                        ])
                        python_code = code_prefix + code_response.content.replace("```python", "").replace("```", "").strip()
                    except Exception as e2:
                        error_message = f"Both Azure and OpenAI code failed: {e2}"
                        print(f"[ERROR] Both Azure and OpenAI code failed: {e2}")
                        continue
                else:
                    error_message = f"Azure code failed: {e}, OpenAI API not available."
                    continue

            await asyncio.to_thread(
                lambda: open(script_path, "w", encoding="utf-8").write(python_code)
            )

            returncode, stdout, stderr = await run_script(script_path)
            if returncode != 0:
                error_message = stderr
                previous_code = python_code
                print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
                continue
        # label .convert_end
        image_path = end_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, end_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        # [Ablation 2] ABLATE_SCORING: skip review scoring, accept first successful result
        if ABLATE_SCORING:
            print("[Ablation] ABLATE_SCORING=True: skipping end slide scoring.")
            return {"end_slide_path": end_path, "end_layout_description": layout_description}

        # === 使用新的审查 Prompt 进行评分 ===
        REVIEW_PROMPT = f"""
Topic: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}
Background tone: {background_tone}
Heading font color: {heading_font_color}
Body font color: {body_font_color}
Font name: {font_name}
Style summary: {style_summary}

You are a Slide Review Expert. Please evaluate the end slide design based on the following dimensions.

{eval_cover}
"""
        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        review_response = await planner_model.ainvoke([
            SystemMessage(content=REVIEW_PROMPT),
            HumanMessage(content=[{
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
            }])
        ])

        print(f"[Review raw] {review_response.content}")

        # 解析评分JSON（严格>60通过）
        def _parse_total_score(payload: str) -> float:
            d = _extract_json_dict(payload)
            total_score_str = str(d.get("Total Score", "-1"))
            try:
                return float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
            except Exception:
                return -1.0

        total_score = _parse_total_score(review_response.content)
        print(f"[Parsed score] {total_score} / 100")
        try:
            if "Suggestions for Improvement" in review_response.content:
                suggestions = _extract_json_dict(review_response.content).get("Suggestions for Improvement", "")
                print(f"[Suggestions] {suggestions}")
        except Exception:
            pass

        if total_score > SCORE_THRESHOLD:
            print(f"布局评分通过（{total_score}），路径：{end_path}")
            return {"end_slide_path": end_path, "end_layout_description": layout_description}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"布局评分未通过，第 {attempt + 1} 次，准备重试。")

    return {"end_slide_path": end_path, "end_layout_description": layout_description}


import os
import asyncio

# 仅替换这个函数；其余代码保持不变
async def compile_ppt(state, config):
    """
    异步方式合并所有生成的单页 PPTX 为一个 PPTX（保留源格式与复杂元素）。
    依赖：pip install lxml
    """
    import re
    import posixpath as P
    from zipfile import ZipFile, ZIP_DEFLATED
    from lxml import etree as ET

    topic = state["topic"]
    ppt_sections = state["ppt_sections"]

    save_dir = os.path.join(".", SAVES_ROOT, topic)
    # os.makedirs(save_dir, exist_ok=True)
    final_ppt_path = os.path.join(save_dir, f"{topic}_final.pptx")

    # ----------- 构建待合并的单页文件列表（严格按顺序） -----------
    paths = ["cover_slide.pptx"]
    for idx, section in enumerate(ppt_sections):
        paths.append(f"section_slide_{idx+1}.pptx")
        paths.extend(
            [f"{section.name}_slide_{slide_index}.pptx"
             for slide_index, _ in enumerate(section.slides, start=1)]
        )
    paths.append("end_slide.pptx")

    input_files = [os.path.join(save_dir, p) for p in paths if os.path.exists(os.path.join(save_dir, p))]
    if not input_files:
        raise RuntimeError("未找到任何可合并的单页 PPTX。请确认生成流程是否已产出文件。")

    # ------------------ Open XML 深度无损合并实现 ------------------
    NS_CT = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
    NS_REL = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
    NS_P = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    FALLBACK_CT = {
        "ppt/slides": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
        "ppt/slideLayouts": "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
        "ppt/slideMasters": "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
        "ppt/theme": "application/vnd.openxmlformats-officedocument.theme+xml",
        "ppt/charts": "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
        "ppt/diagrams/data": "application/vnd.ms-office.drawingml.diagramData+xml",
        "ppt/diagrams/layout": "application/vnd.ms-office.drawingml.diagramLayout+xml",
        "ppt/diagrams/styles": "application/vnd.ms-office.drawingml.diagramStyle+xml",
        "ppt/diagrams/colors": "application/vnd.ms-office.drawingml.diagramColors+xml",
        "ppt/embeddings/ole": "application/vnd.openxmlformats-officedocument.oleObject",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tif": "image/tiff",
        ".tiff": "image/tiff",
        ".emf": "image/x-emf",
        ".wmf": "image/x-wmf",
        ".svg": "image/svg+xml",
        ".mp4": "video/mp4",
        ".mov": "video/quicktime",
        ".m4a": "audio/mp4",
        ".mp3": "audio/mpeg",
        ".wav": "audio/wav",
        ".bin": "application/vnd.openxmlformats-officedocument.oleObject",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xml": "application/xml",
    }

    def _read_zip_to_dict(path):
        with ZipFile(path, "r") as zf:
            return {name: zf.read(name) for name in zf.namelist()}

    def _write_dict_to_zip(files_dict, out_path):
        with ZipFile(out_path, "w", compression=ZIP_DEFLATED) as zf:
            for name, data in files_dict.items():
                zf.writestr(name, data)

    def _parse_xml(byte_content):
        return ET.fromstring(byte_content)

    def _xml_to_bytes(elem):
        return ET.tostring(elem, xml_declaration=True, encoding="UTF-8", standalone="yes")

    def _get_ct_tree(files_dict):
        return _parse_xml(files_dict["[Content_Types].xml"])

    def _ensure_default_ct(ct_tree, ext, ctype):
        exist = ct_tree.xpath(f'/ct:Types/ct:Default[@Extension="{ext.lstrip(".")}"]', namespaces=NS_CT)
        if not exist:
            node = ET.Element("{%s}Default" % NS_CT["ct"])
            node.set("Extension", ext.lstrip("."))
            node.set("ContentType", ctype)
            ct_tree.append(node)

    def _ensure_override_ct(ct_tree, part_name, ctype):
        part_name = part_name if part_name.startswith("/") else "/" + part_name
        exist = ct_tree.xpath(f'/ct:Types/ct:Override[@PartName="{part_name}"]', namespaces=NS_CT)
        if not exist:
            node = ET.Element("{%s}Override" % NS_CT["ct"])
            node.set("PartName", part_name)
            node.set("ContentType", ctype)
            ct_tree.append(node)

    def _src_content_type(src_ct_tree, part_name):
        part_name = part_name if part_name.startswith("/") else "/" + part_name
        ov = src_ct_tree.xpath(f'/ct:Types/ct:Override[@PartName="{part_name}"]', namespaces=NS_CT)
        if ov:
            return ov[0].get("ContentType")
        ext = P.splitext(part_name)[1].lower()
        if ext:
            dv = src_ct_tree.xpath(f'/ct:Types/ct:Default[@Extension="{ext.lstrip(".")}"]', namespaces=NS_CT)
            if dv:
                return dv[0].get("ContentType")
            if ext in FALLBACK_CT:
                return FALLBACK_CT[ext]
        return None

    def _next_index(existing_names, folder, prefix, suffix=".xml"):
        pat = re.compile(rf"^{re.escape(folder)}/{re.escape(prefix)}(\d+){re.escape(suffix)}$")
        m = [int(pat.match(n).group(1)) for n in existing_names if pat.match(n)]
        return (max(m) + 1) if m else 1

    def _rels_path_for(part_path):
        dirname = P.dirname(part_path)
        base = P.basename(part_path)
        return P.join(dirname, "_rels", base + ".rels")

    class PptxDeepMerger:
        def __init__(self, inputs, output):
            self.inputs = inputs
            self.output = output

            # 以第一个源为基底
            self.out_files = _read_zip_to_dict(self.inputs[0])
            self.out_ct = _get_ct_tree(self.out_files)
            self.out_pres = _parse_xml(self.out_files["ppt/presentation.xml"])
            self.out_pres_rels = _parse_xml(self.out_files["ppt/_rels/presentation.xml.rels"])

            names = set(self.out_files.keys())
            self.idx = {
                "slide": _next_index(names, "ppt/slides", "slide"),
                "layout": _next_index(names, "ppt/slideLayouts", "slideLayout"),
                "master": _next_index(names, "ppt/slideMasters", "slideMaster"),
                "theme": _next_index(names, "ppt/theme", "theme"),
                "chart": _next_index(names, "ppt/charts", "chart"),
                "diagram_data": _next_index(names, "ppt/diagrams", "data"),
                "diagram_layout": _next_index(names, "ppt/diagrams", "layout"),
                "diagram_styles": _next_index(names, "ppt/diagrams", "quickStyle"),
                "diagram_colors": _next_index(names, "ppt/diagrams", "colors"),
                "media": _next_index(names, "ppt/media", "image", suffix=""),
                "embed": _next_index(names, "ppt/embeddings", "oleObject", suffix=".bin"),
            }

            self.rid_base = self._max_rid(self.out_pres_rels) + 1
            self.sld_id_base = self._max_sldid(self.out_pres) + 1
            self.part_map = {}  # (src_zip_id, src_abs_part) -> tgt_abs_part

        @staticmethod
        def _max_rid(rels_tree):
            rid_nums = []
            for rel in rels_tree.xpath("/rel:Relationships/rel:Relationship", namespaces=NS_REL):
                rid = rel.get("Id", "")
                m = re.match(r"rId(\d+)$", rid)
                if m:
                    rid_nums.append(int(m.group(1)))
            return max(rid_nums) if rid_nums else 0

        @staticmethod
        def _max_sldid(pres_tree):
            ids = []
            for sldId in pres_tree.xpath("//p:sldId", namespaces=NS_P):
                try:
                    ids.append(int(sldId.get("id")))
                except Exception:
                    pass
            return max(ids) if ids else 255

        def _alloc_path(self, kind, src_path):
            if kind == "slide":
                p = f"ppt/slides/slide{self.idx['slide']}.xml"; self.idx["slide"] += 1; return p
            if kind == "layout":
                p = f"ppt/slideLayouts/slideLayout{self.idx['layout']}.xml"; self.idx["layout"] += 1; return p
            if kind == "master":
                p = f"ppt/slideMasters/slideMaster{self.idx['master']}.xml"; self.idx["master"] += 1; return p
            if kind == "theme":
                p = f"ppt/theme/theme{self.idx['theme']}.xml"; self.idx["theme"] += 1; return p
            if kind == "chart":
                p = f"ppt/charts/chart{self.idx['chart']}.xml"; self.idx["chart"] += 1; return p
            if kind == "diagram_data":
                p = f"ppt/diagrams/data{self.idx['diagram_data']}.xml"; self.idx["diagram_data"] += 1; return p
            if kind == "diagram_layout":
                p = f"ppt/diagrams/layout{self.idx['diagram_layout']}.xml"; self.idx["diagram_layout"] += 1; return p
            if kind == "diagram_styles":
                p = f"ppt/diagrams/quickStyle{self.idx['diagram_styles']}.xml"; self.idx["diagram_styles"] += 1; return p
            if kind == "diagram_colors":
                p = f"ppt/diagrams/colors{self.idx['diagram_colors']}.xml"; self.idx["diagram_colors"] += 1; return p
            if kind == "media":
                ext = P.splitext(src_path)[1].lower()
                p = f"ppt/media/image{self.idx['media']}{ext}"; self.idx["media"] += 1; return p
            if kind == "embed":
                p = f"ppt/embeddings/oleObject{self.idx['embed']}.bin"; self.idx["embed"] += 1; return p
            return src_path  # 兜底

        @staticmethod
        def _classify(path):
            if path.startswith("ppt/slides/"): return "slide"
            if path.startswith("ppt/slideLayouts/"): return "layout"
            if path.startswith("ppt/slideMasters/"): return "master"
            if path.startswith("ppt/theme/"): return "theme"
            if path.startswith("ppt/charts/"): return "chart"
            if path.startswith("ppt/diagrams/data"): return "diagram_data"
            if path.startswith("ppt/diagrams/layout"): return "diagram_layout"
            if path.startswith("ppt/diagrams/quickStyle"): return "diagram_styles"
            if path.startswith("ppt/diagrams/colors"): return "diagram_colors"
            if path.startswith("ppt/media/"): return "media"
            if path.startswith("ppt/embeddings/"): return "embed"
            return "other"

        def _add_override_from_src(self, tgt_path, src_ct_tree, src_path):
            ctype = _src_content_type(src_ct_tree, src_path)
            if ctype is None:
                kind = self._classify(src_path)
                if kind in ("slide", "layout", "master", "theme", "chart"):
                    base_dir = P.dirname(src_path)
                    if base_dir in FALLBACK_CT:
                        ctype = FALLBACK_CT[base_dir]
                if ctype is None:
                    ext = P.splitext(src_path)[1].lower()
                    if ext in FALLBACK_CT:
                        ctype = FALLBACK_CT[ext]
            if ctype:
                _ensure_override_ct(self.out_ct, tgt_path, ctype)
                ext = P.splitext(tgt_path)[1].lower()
                if ext and ext in FALLBACK_CT:
                    _ensure_default_ct(self.out_ct, ext, FALLBACK_CT[ext])

        def _copy_part_recursive(self, src_zip_dict, src_ct_tree, src_zip_id, src_abs_part):
            """
            递归复制指定部件（以及其 .rels 里指向的部件），返回目标路径。
            关键修复：在处理依赖之前，先把当前部件预登记到 part_map，打断 slideLayout<->slideMaster 等环。
           """
            key = (src_zip_id, src_abs_part)
            if key in self.part_map:
                return self.part_map[key]

            # 1) 先分配目标路径，并“预登记”到 part_map —— 这是打断递归环的关键
            kind = self._classify(src_abs_part)
            tgt_path = self._alloc_path(kind, src_abs_part)
            self.part_map[key] = tgt_path  # <-- pre-register to break cycles

            # 2) 写入主体内容 + 内容类型
            data = src_zip_dict[src_abs_part]
            self.out_files[tgt_path] = data
            self._add_override_from_src(tgt_path, src_ct_tree, src_abs_part)

            # 3) 处理 .rels 中的依赖（相对路径 → 绝对路径 → 递归复制 → 新相对路径）
            src_rels = self._rels_path_for(src_abs_part) if hasattr(self, "_rels_path_for") else None
            # 兼容：如果 _rels_path_for 是外部函数
            if src_rels is None:
                # 外层函数名为 _rels_path_for
                from posixpath import dirname, basename, join as pjoin
                def _rels_path_for(part_path):
                    return pjoin(dirname(part_path), "_rels", basename(part_path) + ".rels")
                src_rels = _rels_path_for(src_abs_part)

            if src_rels in src_zip_dict:
                rels_xml = ET.fromstring(src_zip_dict[src_rels])
                for rel in rels_xml.xpath("/rel:Relationships/rel:Relationship", namespaces=NS_REL):
                    target = rel.get("Target")
                    tmode = rel.get("TargetMode", "")
                    if tmode == "External":
                        # 外部链接不复制内容，只保留关系
                        continue

                    src_dir = P.dirname(src_abs_part)
                    abs_src_target = P.normpath(P.join(src_dir, target))
                    # 个别包会写成以 '/' 开头的绝对样式，zip 键没有前导斜杠，这里抹掉
                    if abs_src_target.startswith("/"):
                        abs_src_target = abs_src_target.lstrip("/")

                    # 某些空引用或 Office 注入的占位可能不存在；直接跳过
                    if abs_src_target not in src_zip_dict:
                        continue

                    # 递归复制依赖；由于已预登记，遇到环会直接返回
                    new_abs_target = self._copy_part_recursive(
                        src_zip_dict, src_ct_tree, src_zip_id, abs_src_target
                    )

                    # 回写新的相对路径
                    tgt_dir = P.dirname(tgt_path)
                    new_rel_target = P.relpath(new_abs_target, tgt_dir)
                    rel.set("Target", new_rel_target)

                # 写入更新后的 .rels
                new_rels_path = self._rels_path_for(tgt_path) if hasattr(self, "_rels_path_for") else _rels_path_for(tgt_path)
                self.out_files[new_rels_path] = ET.tostring(
                    rels_xml, xml_declaration=True, encoding="UTF-8", standalone="yes"
                )

            return tgt_path

        @staticmethod
        def _first_slide_path_in(src_zip_dict):
            pres = _parse_xml(src_zip_dict["ppt/presentation.xml"])
            pres_rels = _parse_xml(src_zip_dict["ppt/_rels/presentation.xml.rels"])
            sldId = pres.xpath("//p:sldId", namespaces=NS_P)
            if not sldId:
                raise RuntimeError("源文件未包含任何幻灯片")
            rid = sldId[0].get("{%s}id" % NS_P["r"])
            rel = pres_rels.xpath(f'/rel:Relationships/rel:Relationship[@Id="{rid}"]', namespaces=NS_REL)
            if not rel:
                raise RuntimeError("无法在源 rels 中解析到第一张幻灯片")
            target = rel[0].get("Target")  # 如 slides/slide1.xml
            return P.normpath(P.join("ppt", target))

        def _append_slide_to_presentation(self, new_slide_path):
            rid = f"rId{self.rid_base}"; self.rid_base += 1
            rel_node = ET.Element("{%s}Relationship" % NS_REL["rel"])
            rel_node.set("Id", rid)
            rel_node.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel_node.set("Target", P.relpath(new_slide_path, "ppt"))
            self.out_pres_rels.append(rel_node)

            sldIdLst = self.out_pres.xpath("/p:presentation/p:sldIdLst", namespaces=NS_P)
            if not sldIdLst:
                pres = self.out_pres.xpath("/p:presentation", namespaces=NS_P)[0]
                sldIdLst_elem = ET.Element("{%s}sldIdLst" % NS_P["p"])
                pres.append(sldIdLst_elem)
                sldIdLst = [sldIdLst_elem]
            sld_id_val = str(self.sld_id_base); self.sld_id_base += 1
            sldId_elem = ET.Element("{%s}sldId" % NS_P["p"])
            sldId_elem.set("id", sld_id_val)
            sldId_elem.set("{%s}id" % NS_P["r"], rid)
            sldIdLst[0].append(sldId_elem)

        def merge(self):
            # 其余单页逐一追加第一张幻灯片
            for idx, src in enumerate(self.inputs):
                if idx == 0:
                    continue
                src_zip = _read_zip_to_dict(src)
                src_ct = _get_ct_tree(src_zip)
                slide_abs = self._first_slide_path_in(src_zip)
                new_slide_path = self._copy_part_recursive(src_zip, src_ct, src_zip_id=idx, src_abs_part=slide_abs)
                self._append_slide_to_presentation(new_slide_path)

            # 写回三处关键XML
            self.out_files["[Content_Types].xml"] = _xml_to_bytes(self.out_ct)
            self.out_files["ppt/presentation.xml"] = _xml_to_bytes(self.out_pres)
            self.out_files["ppt/_rels/presentation.xml.rels"] = _xml_to_bytes(self.out_pres_rels)
            _write_dict_to_zip(self.out_files, self.output)

    def _merge_sync(out_path, inputs):
        merger = PptxDeepMerger(inputs, out_path)
        merger.merge()

    # 放到线程池，避免阻塞事件循环
    await asyncio.to_thread(_merge_sync, final_ppt_path, input_files)
    return {"final_ppt_path": final_ppt_path}






# PPT Slide sub-graph -- 
ppt_slide_graph = StateGraph(PPTSlideState, output=PPTSlideOutputState)
# ppt_slide_graph.add_node("generate_slide_content", generate_slide_content)
ppt_slide_graph.add_node("enrich_slide_content", enrich_slide_content)
ppt_slide_graph.add_node("generate_slide_code_and_execute", generate_slide_code_and_execute)
ppt_slide_graph.add_node("ppt_slide_to_image_and_validate", ppt_slide_to_image_and_validate)

ppt_slide_graph.add_edge(START, "enrich_slide_content")
ppt_slide_graph.add_edge("enrich_slide_content", "generate_slide_code_and_execute")
ppt_slide_graph.add_edge("generate_slide_code_and_execute", "ppt_slide_to_image_and_validate")
# ppt_slide_graph.add_edge("generate_slide_content", END)

ppt_slide_subgraph = ppt_slide_graph.compile()

# PPT Sction sub-graph --
ppt_section_graph = StateGraph(PPTSectionState,output=PPTSectionOutputState)
ppt_section_graph.add_node("generate_ppt_section_start", generate_ppt_section_start)
ppt_section_graph.add_node("generate_slide", ppt_slide_subgraph)
ppt_section_graph.add_node("generate_ppt_section_end", generate_ppt_section_end)

ppt_section_graph.add_edge(START, "generate_ppt_section_start")
# ppt_section_graph.add_edge("generate_ppt_section_start", "generate_slide")
ppt_section_graph.add_edge("generate_slide", "generate_ppt_section_end")
# ppt_section_graph.add_edge("generate_ppt_section_end", END)

ppt_section_subgraph = ppt_section_graph.compile()

# Report section sub-graph -- 

# Add nodes 
section_builder = StateGraph(SectionState, output=SectionOutputState)
section_builder.add_node("generate_queries", generate_queries)
section_builder.add_node("search_web", search_web)
section_builder.add_node("write_section", write_section)

# Add edges
section_builder.add_edge(START, "generate_queries")
section_builder.add_edge("generate_queries", "search_web")
section_builder.add_edge("search_web", "write_section")
# section_builder.add_edge("write_section", END)

# Outer graph for initial report plan compiling results from each section -- 

# Add nodes
builder = StateGraph(ReportState, input=ReportStateInput, output=ReportStateOutput, config_schema=Configuration)
builder.add_node("process_image_input", process_image_input)
builder.add_node("generate_report_plan", generate_report_plan)
builder.add_node("human_feedback", human_feedback)
builder.add_node("build_section_with_web_research", section_builder.compile())
builder.add_node("gather_completed_sections", gather_completed_sections)
builder.add_node("write_final_sections", write_final_sections)
builder.add_node("compile_final_report", compile_final_report)
builder.add_node("generate_ppt_outline", generate_ppt_outline)
builder.add_node("generate_ppt_styles", generate_ppt_styles)
builder.add_node("manage_ppt_templates", manage_ppt_templates)
builder.add_node("generate_ppt_sections", ppt_section_subgraph)
builder.add_node("generate_cover_slide", generate_cover_slide)
builder.add_node("generate_section_cover_slides", generate_section_cover_slides)
builder.add_node("generate_end_slide", generate_end_slide)
builder.add_node("compile_ppt", compile_ppt)

# Add edges
builder.add_edge(START, "process_image_input")
builder.add_edge("process_image_input", "generate_report_plan")
builder.add_edge("generate_report_plan", "human_feedback")
builder.add_edge("build_section_with_web_research", "gather_completed_sections")
builder.add_conditional_edges("gather_completed_sections", initiate_final_section_writing, ["write_final_sections"])
builder.add_edge("write_final_sections", "compile_final_report")
# builder.add_edge("compile_final_report", END)
builder.add_edge("compile_final_report", "generate_ppt_outline")
builder.add_edge("generate_ppt_outline", "generate_ppt_styles")
builder.add_edge("generate_ppt_styles", "generate_cover_slide")
builder.add_edge("generate_cover_slide", "generate_section_cover_slides")
builder.add_edge("generate_section_cover_slides", "generate_end_slide")
builder.add_edge("generate_end_slide", "manage_ppt_templates")
# builder.add_edge("generate_ppt_outline", "generate_ppt_sections")
# builder.add_edge("generate_ppt_sections", "compile_ppt")
builder.add_edge("generate_ppt_sections", "compile_ppt")
builder.add_edge("compile_ppt", END)

graph = builder.compile()
