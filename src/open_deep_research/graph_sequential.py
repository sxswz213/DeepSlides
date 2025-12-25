from typing import Literal
import json

from langchain.chat_models import init_chat_model
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.runnables import RunnableConfig

from langgraph.constants import Send
from langgraph.graph import START, END, StateGraph
from langgraph.types import interrupt, Command

import subprocess
import os
import datetime
import asyncio
import io
import httpx
from urllib.parse import urlparse
import re
import shutil
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pathlib import Path

import base64


from open_deep_research.state import (
    ReportStateInput,
    ReportStateOutput,
    Sections,
    ReportState,
    SectionState,
    SectionOutputState,
    Queries,
    Feedback,
    PPTSection,
    PPTOutline,
    PPTSections,
    PPTSectionState,
    PPTSlideState,
    PPTSlide,
    PPTSlideOutputState,
    PPTSectionOutputState
)

from open_deep_research.prompts import (
    report_planner_query_writer_instructions,
    report_planner_instructions,
    query_writer_instructions, 
    section_writer_instructions,
    final_section_writer_instructions,
    section_grader_instructions,
    section_writer_inputs,
    query_writer4PPT_instructions
)

from open_deep_research.configuration import Configuration
from open_deep_research.utils import (
    format_sections, 
    get_config_value, 
    get_search_params, 
    select_and_execute_search,
    set_openai_api_base,
    generate_image_caption,
    generate_image_caption_v2,
    generate_image_caption_v3
)

## Nodes -- 

async def process_image_input(state: ReportState, config: RunnableConfig):
    """处理图像输入并生成描述文本。
    
    这个节点：
    1. 检查状态中是否包含图像路径
    2. 如果有图像，使用Vision API生成详细描述
    3. 将描述作为额外上下文添加到后续步骤
    4. 如果未提供topic，则使用图像描述作为topic
    
    Args:
        state: 当前图状态，包含可选的图像路径
        config: 配置参数
        
    Returns:
        包含图像描述和更新topic的状态
    """
    # 检查是否提供了图像路径
    image_path = state.get("image_path")
    
    # 如果没有提供图像路径，直接返回状态
    if not image_path:
        # 确保topic存在，即使为空字符串
        if "topic" not in state:
            return {"topic": ""}
        return {}
    
    # 确保设置了正确的API基础URL
    set_openai_api_base()
    
    try:
        if not state.get("topic"):
            topic = "The user did not provide a topic."
        else:
            topic = state["topic"]
        # 生成图像描述
        image_result = await generate_image_caption_v3(image_path, topic)
        image_result = json.loads(image_result)
        # print(image_result)
        caption, user_intent, topic = image_result["caption"], image_result["user_intent"], image_result["topic"]

        return {"caption": caption, "user_intent": user_intent, "topic": topic}

    except Exception as e:
        print(f"Error processing image: {str(e)}")
        # 返回错误信息作为caption，并确保topic存在
        if "topic" not in state:
            return {"image_caption": f"Unable to process image:{str(e)}", "topic": ""}
        return {"image_caption": f"Unable to process image:{str(e)}"}

async def generate_report_plan(state: ReportState, config: RunnableConfig):
    """Generate the initial report plan with sections, including image caption and user intent.

    This node:
    1. Gets configuration for the report structure and search parameters
    2. Generates search queries to gather context for planning
    3. Performs web searches using those queries
    4. Uses an LLM to generate a structured plan with sections
    5. Includes image caption and user intent in planning

    Args:
        state: Current graph state containing the report topic, image caption, and user intent
        config: Configuration for models, search APIs, etc.

    Returns:
        Dict containing the generated sections
    """

    # 获取topic、图片caption和用户意图
    topic = state["topic"]
    caption = state.get("image_caption", "")
    user_intent = state.get("user_intent", "")
    feedback = state.get("feedback_on_report_plan", None)

    # Get configuration
    configurable = Configuration.from_runnable_config(config)
    report_structure = configurable.report_structure
    number_of_queries = configurable.number_of_queries
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}
    params_to_pass = get_search_params(search_api, search_api_config)

    if isinstance(report_structure, dict):
        report_structure = str(report_structure)

    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    structured_llm = writer_model.with_structured_output(Queries)

    # Format system instructions including caption and user intent
    system_instructions_query = report_planner_query_writer_instructions.format(
        topic=topic,
        caption=caption,
        user_intent=user_intent,
        report_organization=report_structure,
        number_of_queries=number_of_queries
    )

    results = await structured_llm.ainvoke([
        SystemMessage(content=system_instructions_query),
        HumanMessage(content="Generate search queries that help plan each section of the report.")
    ])

    query_list = [query.search_query for query in results.queries][:number_of_queries]

    source_str = await select_and_execute_search(search_api, query_list, params_to_pass)

    image_path = state.get("image_path")
    if image_path:
        try:
            # 单独调用图片搜索API
            image_search_result = await select_and_execute_search("image_search", [image_path], params_to_pass)
            # 将图片搜索结果与之前的搜索结果合并
            source_str += f"\n\nImage search results:\n{image_search_result}"
        except Exception as e:
            print(f"Error calling the image-search API:{str(e)}")

    # Include caption and user intent in report planner instructions
    system_instructions_sections = report_planner_instructions.format(
        topic=topic,
        caption=caption,
        user_intent=user_intent,
        report_organization=report_structure,
        context=source_str,
        feedback=feedback
    )

    planner_provider = get_config_value(configurable.planner_provider)
    planner_model = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})

    planner_message = """Generate the sections of the report. Your reply must contain a 'sections' field that lists the sections.  
Each section must include: name, description, plan, research, and content fields."""

    set_openai_api_base()

    if planner_model == "claude-3-7-sonnet-latest":
        planner_llm = init_chat_model(
            model=planner_model,
            model_provider=planner_provider,
            max_tokens=20_000,
            thinking={"type": "enabled", "budget_tokens": 16_000}
        )
    else:
        # With other models, thinking tokens are not specifically allocated
        # planner_llm = init_chat_model(model=planner_model, 
        #                               model_provider=planner_provider,
        #                               model_kwargs=planner_model_kwargs)
        planner_llm = AzureChatOpenAI(
            model=configurable.writer_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],  # Azure's API base
            deployment_name=planner_model_kwargs["azure_deployment"],  # Azure's deployment name
            openai_api_version=planner_model_kwargs["openai_api_version"],  # Azure's API version
            temperature=0,
            max_tokens=4096
        )

    # Generate the report sections
    structured_llm = planner_llm.with_structured_output(Sections)
    report_sections = await structured_llm.ainvoke([SystemMessage(content=system_instructions_sections),
                                             HumanMessage(content=planner_message)])

    sections = report_sections.sections

    return {"sections": sections}

def human_feedback(state: ReportState, config: RunnableConfig) -> Command[Literal["generate_report_plan","build_section_with_web_research"]]:
    """Get human feedback on the report plan and route to next steps.
    
    This node:
    1. Formats the current report plan for human review
    2. Gets feedback via an interrupt
    3. Routes to either:
       - Section writing if plan is approved
       - Plan regeneration if feedback is provided
    
    Args:
        state: Current graph state with sections to review
        config: Configuration for the workflow
        
    Returns:
        Command to either regenerate plan or start section writing
    """

    # Get sections
    topic = state["topic"]
    sections = state['sections']
    sections_str = "\n\n".join(
        f"Section: {section.name}\n"
        f"Description: {section.description}\n"
        f"Research needed: {'Yes' if section.research else 'No'}\n"
        for section in sections
    )

    # # Get feedback on the report plan from interrupt
    # interrupt_message = f"""Please provide feedback on the following report plan.\n\n{sections_str}\n
    # Does this report plan meet your needs?\nPass in 'true' to approve the report plan.\nAlternatively, provide feedback to regenerate the report plan:"""
    
    # feedback_dict = interrupt(interrupt_message)
    # if type(feedback_dict) == bool:
    #     # If the feedback is a boolean, treat it as approval
    #     feedback = feedback_dict
    # else:
    #     feedback = list(feedback_dict.values())[0]
    # # print(feedback)
    # # If the user approves the report plan, kick off section writing
    # if isinstance(feedback, bool) and feedback is True:
    #     # Treat this as approve and kick off section writing
    #     return Command(goto=[
    #         Send("build_section_with_web_research", {"topic": topic, "section": s, "search_iterations": 0}) 
    #         for s in sections 
    #         if s.research
    #     ])
    
    # # If the user provides feedback, regenerate the report plan 
    # elif isinstance(feedback, str):
    #     # Treat this as feedback
    #     return Command(goto="generate_report_plan", 
    #                    update={"feedback_on_report_plan": feedback})
    # else:
    #     raise TypeError(f"Interrupt value of type {type(feedback)} is not supported.")
    # tooooooodoooooo 此处为忽略人类反馈直接进行下一步
    return Command(
        goto=[
            Send(
                "build_section_with_web_research",
                {"topic": topic, "section": section, "search_iterations": 0},
            )
            for section in sections
            if section.research
        ]
    )
    
async def generate_queries(state: SectionState, config: RunnableConfig):
    """Generate search queries for researching a specific section.
    
    This node uses an LLM to generate targeted search queries based on the 
    section topic and description.
    
    Args:
        state: Current state containing section details
        config: Configuration including number of queries to generate
        
    Returns:
        Dict containing the generated search queries
    """

    # Get state 
    topic = state["topic"]
    section = state["section"]

    # Get configuration
    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate queries 
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    structured_llm = writer_model.with_structured_output(Queries)

    # Format system instructions
    system_instructions = query_writer_instructions.format(topic=topic, 
                                                           section_topic=section.description, 
                                                           number_of_queries=number_of_queries)

    # Generate queries  
    queries = await structured_llm.ainvoke([SystemMessage(content=system_instructions),
                                     HumanMessage(content="Generate search queries for the given topic.")])

    return {"search_queries": queries.queries}

async def search_web(state: SectionState, config: RunnableConfig):
    """Execute web searches for the section queries.
    
    This node:
    1. Takes the generated queries
    2. Executes searches using configured search API
    3. Formats results into usable context
    
    Args:
        state: Current state with search queries
        config: Search API configuration
        
    Returns:
        Dict with search results and updated iteration count
    """

    # Get state
    search_queries = state["search_queries"]
    # Get configuration

    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}  # Get the config dict, default to empty
    params_to_pass = get_search_params(search_api, search_api_config)  # Filter parameters

    # Web search
    query_list = [query.search_query for query in search_queries][:number_of_queries]

    # Search the web with parameters
    source_str = await select_and_execute_search(search_api, query_list, params_to_pass)

    return {"source_str": source_str, "search_iterations": state["search_iterations"] + 1}

async def write_section(state: SectionState, config: RunnableConfig) -> Command[Literal[END, "search_web"]]:
    """Write a section of the report and evaluate if more research is needed.
    
    This node:
    1. Writes section content using search results
    2. Evaluates the quality of the section
    3. Either:
       - Completes the section if quality passes
       - Triggers more research if quality fails
    
    Args:
        state: Current state with search results and section info
        config: Configuration for writing and evaluation
        
    Returns:
        Command to either complete section or do more research
    """

    # Get state 
    topic = state["topic"]
    section = state["section"]
    source_str = state["source_str"]
    
    # 提取图像信息
    images_data = []
    # dtyxs TODO: make it configurable
    max_images = 6  # 最大图像数量限制
    
    if "--- IMAGES ---" in source_str:
        try:
            # 提取图像部分
            images_section = source_str.split("--- IMAGES ---")[1].split("-" * 80)[0]
            image_blocks = images_section.strip().split("IMAGE ")[1:]  # 跳过第一个空元素
            
            for i, block in enumerate(image_blocks):
                lines = block.strip().split("\n")
                image_url = ""
                image_description = ""
                
                for line in lines:
                    if line.startswith("URL:"):
                        image_url = line.replace("URL:", "").strip()
                    elif line.startswith("DESCRIPTION:"):
                        image_description = line.replace("DESCRIPTION:", "").strip()
                
                if image_url:  # 只添加有URL的图像
                    images_data.append({
                        "index": i,
                        "url": image_url,
                        "description": image_description
                    })
        except Exception as e:
            print(f"Error extracting image information: {str(e)}")

    # 达到最大图像数量后停止处理
    image_num_available = len(images_data)
    if image_num_available >= max_images:
        images_data = images_data[:max_images]

    # 将图像数据格式化为JSON字符串
    images_json = json.dumps(images_data, ensure_ascii=False, indent=2) if images_data else "[]"
    
    if images_data:
        # print(f"Extracted {image_num_available} images (maximum limit: {max_images})")
        pass

    # Get configuration
    configurable = Configuration.from_runnable_config(config)

    # Format system instructions
    section_writer_inputs_formatted = section_writer_inputs.format(
        topic=topic, 
        section_name=section.name, 
        section_topic=section.description, 
        context=source_str, 
        section_content=section.content,
        images_data=images_json
    )

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate section  
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )

    # dtyxs TODO: Native image input
    section_content = await writer_model.ainvoke([SystemMessage(content=section_writer_instructions),
                                           HumanMessage(content=section_writer_inputs_formatted)])
    
    # 处理返回的内容，提取图像选择信息
    content = section_content.content

    # 检查是否包含图像选择信息
    if "```image_selection" in content and images_data:
        try:
            # 提取图像选择JSON
            image_selection_text = content.split("```image_selection")[1].split("```")[0].strip()
            image_selection = json.loads(image_selection_text)
            
            selected_index = image_selection.get("selected_image_index", -1)
            if 0 <= selected_index < len(images_data):
                selected_image = images_data[selected_index]
                selected_image["caption"] = image_selection.get("caption", "")
                
                # 处理内容分割，保留标记前后的内容
                before_image_selection = content.split("```image_selection")[0].strip()
                after_image_selection = ""
                if "```" in content.split("```image_selection")[1]:
                    after_image_selection = content.split("```image_selection")[1].split("```", 1)[1].strip()
                
                # 重新组合内容，将图像选择部分替换为图像和标题标记
                image_content = f"\n\n![{selected_image['caption']}]({selected_image['url']})\n*{selected_image['caption']}*\n"
                if after_image_selection:
                    content = f"{before_image_selection}\n\n{image_content}\n\n{after_image_selection}"
                else:
                    content = f"{before_image_selection}\n\n{image_content}"
            else:
                # selected_index不合法，移除图像选择部分
                before_image_selection = content.split("```image_selection")[0].strip()
                after_image_selection = ""
                if "```" in content.split("```image_selection")[1]:
                    after_image_selection = content.split("```image_selection")[1].split("```", 1)[1].strip()
                
                # 重新组合内容，移除图像选择部分但保留其前后内容
                if after_image_selection:
                    content = f"{before_image_selection}\n\n{after_image_selection}"
                else:
                    content = before_image_selection
        except Exception as e:
            print(f"Error processing image selection information: {str(e)}")
    
    # Write content to the section object  
    section.content = content
    section.source_str = source_str  # Store the source string in the section

    # Grade prompt 
    section_grader_message = (
        "Grade the report and consider follow-up questions for any missing information. "
        "If the grade is 'pass', all subsequent queries should return an empty string. "
        "If the grade is 'fail', provide specific search queries to collect the missing information."
    )

    
    section_grader_instructions_formatted = section_grader_instructions.format(topic=topic, 
                                                                               section_topic=section.description,
                                                                               section=section.content, 
                                                                               number_of_follow_up_queries=configurable.number_of_queries)

    # Use planner model for reflection
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})

    # Set OpenAI API Base URL
    set_openai_api_base()

    if planner_model == "claude-3-7-sonnet-latest":
        # Allocate a thinking budget for claude-3-7-sonnet-latest as the planner model
        reflection_model = init_chat_model(model=planner_model, 
                                           model_provider=planner_provider, 
                                           max_tokens=20_000, 
                                           thinking={"type": "enabled", "budget_tokens": 16_000}).with_structured_output(Feedback)
    else:
        # reflection_model = init_chat_model(model=planner_model, 
        #                                    model_provider=planner_provider, model_kwargs=planner_model_kwargs).with_structured_output(Feedback)
        reflection_model = AzureChatOpenAI(
            model=configurable.writer_model,
            azure_endpoint=planner_model_kwargs["openai_api_base"],  # Azure's API base
            deployment_name=planner_model_kwargs["azure_deployment"],  # Azure's deployment name
            openai_api_version=planner_model_kwargs["openai_api_version"],  # Azure's API version
            temperature=0,
            max_tokens=2048
        ).with_structured_output(Feedback)
    # Generate feedback
    feedback = await reflection_model.ainvoke([SystemMessage(content=section_grader_instructions_formatted),
                                        HumanMessage(content=section_grader_message)])

    # If the section is passing or the max search depth is reached, publish the section to completed sections 
    if feedback.grade == "pass" or state["search_iterations"] >= configurable.max_search_depth:
        # Publish the section to completed sections 
        return  Command(
        update={"completed_sections": [section]},
        goto=END
    )

    # Update the existing section with new content and update search queries
    else:
        return  Command(
        update={"search_queries": feedback.follow_up_queries, "section": section},
        goto="search_web"
        )
    
async def write_final_sections(state: SectionState, config: RunnableConfig):
    """Write sections that don't require research using completed sections as context.
    
    This node handles sections like conclusions or summaries that build on
    the researched sections rather than requiring direct research.
    
    Args:
        state: Current state with completed sections as context
        config: Configuration for the writing model
        
    Returns:
        Dict containing the newly written section
    """

    # Get configuration
    configurable = Configuration.from_runnable_config(config)

    # Get state 
    topic = state["topic"]
    section = state["section"]
    completed_report_sections = state["report_sections_from_research"]
    
    # Format system instructions
    system_instructions = final_section_writer_instructions.format(topic=topic, section_name=section.name, section_topic=section.description, context=completed_report_sections)

    # Set OpenAI API Base URL
    set_openai_api_base()

    # Generate section  
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure's API base
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure's deployment name
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure's API version
        temperature=0,
        max_tokens=2048
    )
    
    section_content = await writer_model.ainvoke([SystemMessage(content=system_instructions),
                                           HumanMessage(content="Generate the report sections based on the provided information.")])
    
    # Write content to section 
    section.content = section_content.content

    # Write the updated section to completed sections
    return {"completed_sections": [section]}

def gather_completed_sections(state: ReportState):
    """Format completed sections as context for writing final sections.
    
    This node takes all completed research sections and formats them into
    a single context string for writing summary sections.
    
    Args:
        state: Current state with completed sections
        
    Returns:
        Dict with formatted sections as context
    """

    # List of completed sections
    completed_sections = state["completed_sections"]

    # Format completed section to str to use as context for final sections
    completed_report_sections = format_sections(completed_sections)

    return {"report_sections_from_research": completed_report_sections}

async def compile_final_report(state: ReportState):
    """Compile all sections into the final report.
    
    This node:
    1. Gets all completed sections
    2. Orders them according to original plan
    3. Combines them into the final report
    
    Args:
        state: Current state with all completed sections
        
    Returns:
        Dict containing the complete report
    """

    # Get sections
    topic = state["topic"]
    sections = state["sections"]
    completed_sections = {s.name: s.content for s in state["completed_sections"]}

    # Update sections with completed content while maintaining original order
    for section in sections:
        section.content = completed_sections[section.name]

    # Compile final report
    all_sections = "\n\n".join([s.content for s in sections])

    save_dir = os.path.join(".", "saves", topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    outline_path = os.path.join(save_dir, "final_report.md")


    try:
        await asyncio.to_thread(
            lambda p=outline_path, d=all_sections: open(p, "w", encoding="utf-8").write(
                json.dumps(d, ensure_ascii=False, indent=2)
            )
        )
    except Exception as exc:
        outline_path = ""

    return {"final_report": all_sections}

def initiate_final_section_writing(state: ReportState):
    """Create parallel tasks for writing non-research sections.
    
    This edge function identifies sections that don't need research and
    creates parallel writing tasks for each one.
    
    Args:
        state: Current state with all sections and research context
        
    Returns:
        List of Send commands for parallel section writing
    """

    # Kick off section writing in parallel via Send() API for any sections that do not require research
    return [
        Send("write_final_sections", {"topic": state["topic"], "section": s, "report_sections_from_research": state["report_sections_from_research"]}) 
        for s in state["sections"] 
        if not s.research
    ]

async def generate_ppt_outline(state: ReportState, config: RunnableConfig)-> Command[Literal["generate_ppt_sections"]]:
    """
    根据演讲时长确定推荐的PPT页数，根据风格和故事线重新生成章节划分，最后基于此生成PPT大纲。

    Args:
        state: 当前状态，包含 final_report 等信息
        config: 配置参数

    Returns:
        包含PPT大纲的状态字典
    """
    configurable = Configuration.from_runnable_config(config)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.planner_model,
        azure_endpoint=planner_model_kwargs["openai_api_base"],
        deployment_name=planner_model_kwargs["azure_deployment"],
        openai_api_version=planner_model_kwargs["openai_api_version"],
        temperature=0,
        max_tokens=4096
    )

    topic = state.get("topic", "none")
    presentation_minutes = state.get("presentation_minutes", "45")
    style = state.get("style", "none")

    if style == "none":
        prefix = (
            "Reference presentation styles: professional business, modern tech, minimalist, "
            "creative lively, academically rigorous, storytelling narrative, magazine visual, "
            "illustration cartoon, retro nostalgic, data visualization."
        )
    else:
        prefix = f"User's expected presentation style/storyline: {style}"

    storyline_prompt = f"""
    You are an experienced presentation expert tasked with creating a presentation PPT. Now you need to determine the presentation's style, storyline, and theme colors (primary color + accent color).
    Presentation topic: {topic}
    Presentation duration: {presentation_minutes} minutes
    {prefix}

    Reference storyline templates:
    - Problem-Solution: Clearly identify a core problem and provide clear, specific solutions.
    - Situation-Conflict-Resolution-Outcome: First set up a scene, describe the challenge, offer the solution, and finally present positive results.
    - SCQA (Situation-Complication-Question-Answer): Provide background information, introduce the complication, state the key question clearly, and give the answer.
    - Timeline (Past-Present-Future): Present past events, current status, and future goals in chronological order.
    - Contrast (Current vs. Future): Clearly contrast existing problems with the ideal future state, highlighting how to achieve the transformation.
    - Pyramid: Start with the conclusion and unfold arguments layer by layer from top to bottom, reinforcing the core idea with rigorous, clear logic.
    - Research Report: Introduce the topic within its field background, systematically outline the current status, methods, and challenges, and finally present future research trends and directions.

    Return JSON format:
    {{
        "style": "Style",
        "storyline": "Storyline type",
        "main_color": "Recommended primary color",
        "accent_color": "Recommended accent color"
    }}
    """

    storyline_response = await writer_model.ainvoke([
        SystemMessage(content=storyline_prompt),
        HumanMessage(content="Please recommend a suitable storyline.")
    ])
    response_content = json.loads(storyline_response.content)
    style = response_content["style"]
    storyline = response_content["storyline"]
    main_color = response_content["main_color"]
    accent_color = response_content["accent_color"]

    ppt_length_prompt = f"""
    You are an experienced presentation expert.

    Topic: {topic}
    Presentation duration: {presentation_minutes} minutes

    Please suggest an appropriate number of PPT slides for this duration
    (each slide should have a moderate amount of content, not overcrowded;
    one slide generally corresponds to about 1-2 minutes of presentation time).
    JSON format: {{\"recommended_slides\": 10}}
    """


    ppt_length_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_length_prompt),
        HumanMessage(content="Please provide the recommended number of PPT slides.")] 
    )

    recommended_slides = json.loads(ppt_length_response.content)["recommended_slides"]
    # print(f"Recommended number of slides: {recommended_slides}")

    ppt_section_distribution_prompt = f"""
    Presentation topic: {topic}
    Style: {style}
    Storyline: {storyline}
    Recommended total slides: {recommended_slides}

    Attention: Do NOT create a "Q&A" or "Closing" section.

    Based on the above information, re-plan the PPT section structure and allocate the number of slides for each section. Return in JSON format, for example:
    {{
        "section_distribution": {{
            "Introduction": 2,
            "Methodology": 3,
            "Results": 3,
            "Conclusion": 2
        }}
    }}
    """


    ppt_distribution_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_section_distribution_prompt),
        HumanMessage(content="Please plan the section structure and allocate the number of pages.")
    ])

    section_distribution = json.loads(ppt_distribution_response.content)["section_distribution"]
    # print(f"Section distribution: {section_distribution}")
    ppt_outline_prompt = f"""
    You excel at designing slide outlines for presentations.
    Presentation topic: {topic}
    Style: {style}
    Storyline: {storyline}

    Reference material for the presentation:
    {state["final_report"]}

    The slide allocation for each section is as follows:
    {json.dumps(section_distribution, ensure_ascii=False, indent=2)}

    Please generate a PPT outline that adheres to the above slide allocation. Each slide should include:
    - A title
    - key points 

    Important: **Do NOT make every slide contain exactly 3 or 4 key points.**  
    Ensure variety by creating some slides with **5** key points, and some with **6**.

    JSON format:
    Return JSON only, e.g.:
    {{
        "ppt_sections": [
            {{
            "name": "Introduction",
            "allocated_slides": 2,
            "slides": [
                {{"title":"Sample 6-point slide","points":["A","B","C","D","E","F"]}},
                {{"title":"Sample 4-point slide","points":["A","B","C","D"]}},
                {{"title":"Sample 5-point slide","points":["A","B","C","D","E"]}},
                {{"title":"Sample 3-point slide","points":["A","B","C"]}}
            ]
            }}
        ]
    }}
    """


    ppt_outline_response = await writer_model.ainvoke([
        SystemMessage(content=ppt_outline_prompt),
        HumanMessage(content="Generate a PPT outline.")
    ])
    # print(f"PPT outline response: {ppt_outline_response.content}")
    ppt_sections_data = json.loads(ppt_outline_response.content.split("```json")[-1].split('```')[0])["ppt_sections"]
    # print(f"PPT outline sections: {ppt_sections_data}")
    for section in ppt_sections_data:
        for slide in section.get("slides", []):
            slide.setdefault("codes", [])
            slide.setdefault("detail", "")
            slide.setdefault("enriched_points", "")
            slide.setdefault("path", "")

    ppt_sections = [PPTSection(**section) for section in ppt_sections_data]
    ppt_outline = PPTOutline(ppt_sections=PPTSections(sections=ppt_sections))
    save_dir = os.path.join(".", "saves", "outlines", topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    outline_path = os.path.join(save_dir, "ppt_outline.json")

    outline_payload = {
        "recommended_slides": recommended_slides,
        "section_distribution": section_distribution,
        "ppt_outline": [section.model_dump() for section in ppt_sections],
    }

    try:
        await asyncio.to_thread(
            lambda p=outline_path, d=outline_payload: open(p, "w", encoding="utf-8").write(
                json.dumps(d, ensure_ascii=False, indent=2)
            )
        )
    except Exception as exc:
        print(f"[WARN] Error saving PPT report: {exc}")
        outline_path = ""
    return Command(
        update={
            "recommended_ppt_slides": recommended_slides,
            "section_distribution": section_distribution,
            "ppt_outline": ppt_outline,
            "ppt_sections": ppt_sections,
            "storyline": storyline,
            "style": style,
            "main_color": main_color,  
            "accent_color": accent_color
        },
        goto=[
            Send("generate_ppt_sections", {"topic": topic, "ppt_section": ppt_section, "style": style, "main_color": main_color, "accent_color": accent_color})
            for ppt_section in ppt_sections
        ]
    )

async def save_image_from_url(image_url, image_name, topic, ppt_section_name, slide_index):
    """
    异步下载图片并保存到本地，自动根据URL后缀决定图片类型。

    Args:
        image_url: 图片URL
        image_name: 保存的图片名称
        topic: 主题目录名称
        ppt_section_name: PPT章节名称
        slide_index: 幻灯片索引

    Returns:
        str: 图片保存路径或空字符串（失败时）
    """
    try:
        # 异步网络请求下载图片
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(image_url)
            response.raise_for_status()

            # 提取图片后缀
            parsed_url = urlparse(image_url)
            ext = os.path.splitext(parsed_url.path)[-1].lower()

            # 默认后缀为 .jpg，如果没有或后缀不标准时使用默认
            if ext not in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']:
                ext = '.jpg'

            # 异步创建目录（避免阻塞）
            save_dir = os.path.join(
                ".", "saves", topic, "images", ppt_section_name, f"slide_{slide_index+1}"
            )
            await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)

            # 构建完整图片路径
            image_path = os.path.join(save_dir, f"{image_name}{ext}")

            # 异步保存文件（避免阻塞）
            await asyncio.to_thread(
                lambda: open(image_path, 'wb').write(response.content)
            )

            # print(f"Image successfully saved to: {image_path}")
            return image_path

    except httpx.HTTPError as e:
        # print(f"HTTP request error: {str(e)}")
        return ""
    except Exception as e:
        print(f"Other error: {str(e)}")
        return ""

async def truncate_by_characters(text, max_chars=250000):
    if len(text) > max_chars:
        return text[:max_chars]
    return text

async def enrich_slide_content(state: PPTSlideState, config: RunnableConfig):
    """
    第一部分：根据幻灯片的标题和要点生成详细的内容描述，并生成幻灯片布局描述。

    Args:
        state: 包含当前PPT幻灯片信息的状态。
        config: 配置参数。

    Returns:
        Tuple: 包含生成的详细内容和幻灯片布局描述。
    """
    configurable = Configuration.from_runnable_config(config)
    number_of_queries = configurable.number_of_queries_for_ppt
    writer_model_kwargs = configurable.writer_model_kwargs or {}
    
    # 获取当前幻灯片信息
    topic = state["topic"]
    ppt_section = state["ppt_section"]
    slide_index = state["slide_index"]
    style = state.get("style")
    main_color = state.get("main_color")
    accent_color = state.get("accent_color")
    slide = ppt_section.slides[slide_index]
    slide_title = slide.title
    slide_points = slide.points

    # 设置OpenAI API的基础URL
    set_openai_api_base()

    # 使用Azure模型生成查询语句
    writer_provider = get_config_value(configurable.writer_provider)
    writer_model_name = get_config_value(configurable.writer_model)
    writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    writer_model = AzureChatOpenAI(
        model=configurable.writer_model,
        azure_endpoint=writer_model_kwargs["openai_api_base"],  # Azure的API基础URL
        deployment_name=writer_model_kwargs["azure_deployment"],  # Azure的部署名称
        openai_api_version=writer_model_kwargs["openai_api_version"],  # Azure的API版本
        temperature=0,
        max_tokens=4096
    )

    # 使用结构化输出生成查询语句
    structured_llm = writer_model.with_structured_output(Queries)

    # 格式化系统指令
    system_instructions = query_writer4PPT_instructions.format(
        topic=topic,
        section_topic=ppt_section.name,
        slide_title=slide_title,
        slide_points=", ".join(slide_points),
        number_of_queries=number_of_queries
    )

    # 生成查询语句
    query_response = await structured_llm.ainvoke([
        SystemMessage(content=system_instructions),
        HumanMessage(content="Generate relevant search queries based on the information above.")
    ])

    # 从返回的查询结果中提取查询语句
    search_queries = [query.search_query for query in query_response.queries][:number_of_queries]

    # 执行Web搜索
    search_api = get_config_value(configurable.search_api)
    search_api_config = configurable.search_api_config or {}
    params_to_pass = get_search_params(search_api, search_api_config)

    source_str = await select_and_execute_search(search_api, search_queries, params_to_pass)

    # 提取图像信息
    images_data = []
    # dtyxs TODO: make it configurable
    max_images = 6  # 最大图像数量限制
    
    if "--- IMAGES ---" in source_str:
        try:
            # 提取图像部分
            images_section = source_str.split("--- IMAGES ---")[1].split("-" * 80)[0]
            image_blocks = images_section.strip().split("IMAGE ")[1:]  # 跳过第一个空元素
            
            for i, block in enumerate(image_blocks):
                lines = block.strip().split("\n")
                image_url = ""
                image_description = ""
                
                for line in lines:
                    if line.startswith("URL:"):
                        image_url = line.replace("URL:", "").strip()
                    elif line.startswith("DESCRIPTION:"):
                        image_description = line.replace("DESCRIPTION:", "").strip()
                
                if image_url:  # 只添加有URL的图像
                    # 保存图像并获取保存路径
                    image_path = await save_image_from_url(image_url, f"image_{i+1}", topic, ppt_section.name, slide_index)
                    images_data.append({
                        "index": i,
                        "url": image_url,
                        "description": image_description,
                        "local_path": image_path  # 添加本地图片路径
                    })
        except Exception as e:
            print(f"Error extracting image information: {str(e)}")

    # 达到最大图像数量后停止处理
    image_num_available = len(images_data)
    if image_num_available >= max_images:
        images_data = images_data[:max_images]

    # 将图像数据格式化为JSON字符串
    images_json = json.dumps(images_data, ensure_ascii=False, indent=2) if images_data else "[]"
    
    if images_data:
        # print(f"Extracted {image_num_available} images (maximum limit: {max_images})")
        pass

    images_data = images_data[:max_images] if len(images_data) >= max_images else images_data
    embedded_images = []
    for img in images_data:
        local_path = img.get("local_path", "")
        if local_path and os.path.exists(local_path):
            try:
                img_bytes = await asyncio.to_thread(
                    lambda path=local_path: open(path, "rb").read()
                )

                img_base64 = base64.b64encode(img_bytes).decode()
                ext = os.path.splitext(local_path)[1].lower().replace('.', '')
                mime_type = f'image/{ext if ext != "jpg" else "jpeg"}'

                embedded_images.append({
                    "index": img["index"],
                    "path": local_path,
                    "description": img["description"],
                    # "base64": f"data:{mime_type};base64,{img_base64}"
                })
            except Exception as e:
                print(f"Error loading image: {str(e)}")
    
    images_json_embedded = json.dumps(embedded_images, ensure_ascii=False, indent=2) if embedded_images else "[]"
    
    if embedded_images:
        # print(f"Successfully loaded {len(embedded_images)} images in Base64 format.")
        pass
    
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    # coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs=coder_model_kwargs)
    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        # temperature=0.7,
        # max_tokens=4096
        max_completion_tokens=40960
    )
    # 扩展幻灯片内容
    content_enrichment_prompt = f"""
    Based on the following points, expand the slide content. Provide a detailed description for each point and ensure it is linked to the search results.

    Slide topic: {topic}
    Slide section: {ppt_section.name}
    Slide title: {slide_title}
    Points: {', '.join(slide_points)}

    Search results:
    {source_str}
    Please expand and describe each point in detail, suitable for a presentation. Keep the language concise yet informative—ideally, each expanded point should not exceed **30 words**. Return the result in JSON format only; do not output any other text:

    {{
        "enriched_points": [
            {{"point_title": "Point Title 1", "expanded_content": "Expanded content 1"}},
            {{"point_title": "Point Title 2", "expanded_content": "Expanded content 2"}},
            {{"point_title": "Point Title 3", "expanded_content": "Expanded content 3"}},
            ...
        ]
    }}
    """


    enrichment_response = await writer_model.ainvoke([
        SystemMessage(content=content_enrichment_prompt),
        HumanMessage(content="Please expand the slide content. Only return Json, no additional text.")
    ])
    # print("JSOOOOOOOOOOON")
    # enriched_points = json.loads(enrichment_response.content.replace("```json", "").replace("```", "").strip())["enriched_points"]
    enriched_points = enrichment_response.content

    save_dir = os.path.join(".", "saves", "outlines", topic)
    await asyncio.to_thread(os.makedirs, save_dir, exist_ok=True)
    safe_section_name = ppt_section.name.replace(" ", "_")
    file_path = os.path.join(save_dir, f"{safe_section_name}_slide{slide_index+1}.json")

    # data_to_save = {
    #     "query_results": source_str,
    #     "enriched_points": enriched_points,
    # }

    # try:
    #     await asyncio.to_thread(
    #         lambda p=file_path, d=data_to_save: open(p, "w", encoding="utf-8").write(
    #             json.dumps(d, ensure_ascii=False, indent=2)
    #         )
    #     )
    #     # print(f"[INFO] Saved query & enriched points to: {file_path}")
    # except Exception as exc:
    #     print(f"[WARN] Error saving output file: {exc}")
    #     file_path = ""

    # 用于收集数据的返回值，不生成完整的幻灯片，只生成内容 TOOOOOOOOCHANGE
    # generated_slide = PPTSlide(
    #         title=slide_title,
    #         points=slide_points,
    #         codes=['no code'],
    #         enriched_points=json.dumps(enriched_points, ensure_ascii=False, indent=2),
    #         detail="no detail",
    #     )
    # return Command(
    #         update={"completed_slides": [generated_slide]},
    #         goto=END
    # )

    # 第二阶段：使用coder_model生成幻灯片布局描述
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    # writer_provider = get_config_value(configurable.writer_provider)
    # writer_model_name = get_config_value(configurable.writer_model)
    # writer_model_kwargs = get_config_value(configurable.writer_model_kwargs or {})
    # writer_model = init_chat_model(model=writer_model_name, model_provider=writer_provider, model_kwargs=writer_model_kwargs)
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})
    # coder_model = init_chat_model(model=coder_model_name, model_provider=coder_provider, model_kwargs=coder_model_kwargs)
    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        # temperature=0.7,
        # max_tokens=4096
        max_completion_tokens=40960
    )

    detail_prompt = f"""
    You are a seasoned slide designer responsible for designing slide layouts.
    Generate a slide layout description in JSON format based on the following details:

    Title: {slide_title}
    Detailed points: {enriched_points}
    Slide style: {style}
    Primary color: {main_color}
    Accent color: {accent_color}

    Here are some images you may optionally use in the PPT:
    <Image list>
    {images_json_embedded}
    </Image list>

    Return JSON only, no additional text:
    {{
        "layout": "Layout type",
        "content_details": ["Layout for title, detailed content, images (positions, dimensions, paths, etc.)"],
        "design_style": "Design style"
    }}
    """


    detail_response = await coder_model.ainvoke([
        SystemMessage(content=detail_prompt),
        HumanMessage(content="Please output the JSON for the slide layout. Return JSON only, no additional text.")
    ])
    
    slide_detail = detail_response.content

    return {"enriched_points": enriched_points, 
            "slide_detail": slide_detail}


async def generate_slide_code_and_execute(state: PPTSlideState, config: RunnableConfig):
    enriched_points = state["enriched_points"]
    slide_detail = state["slide_detail"]

    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})

    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960,
    )

    topic = state["topic"]
    ppt_section = state["ppt_section"]
    slide_index = state["slide_index"]
    main_color = state.get("main_color")
    accent_color = state.get("accent_color")
    style = state.get("style")
    slide_title = ppt_section.slides[slide_index].title
    slide_points = ppt_section.slides[slide_index].points

    # ✅ 使用 Path，避免 abspath/cwd；真正需要绝对路径时放到线程里 resolve
    save_dir = Path("saves") / topic
    await asyncio.to_thread(save_dir.mkdir, parents=True, exist_ok=True)

    error_message = ""
    previous_code = ""
    python_code = None
    execution_successful = False

    async def run_script(script_path: Path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", str(script_path)],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    for attempt in range(3):
        code_prompt = f"""
Generate Python code that creates slides using the python-pptx library based on the following detailed slide description:

Title: {slide_title}
Detailed bullet points: {enriched_points}
Slide description: {slide_detail}
Slide style: {style}
Primary color: {main_color} 
Accent color: {accent_color}

Level-1 heading font size: 36
Other heading levels font size: 20
Body text font size: 20
Chinese font: 微软雅黑
English font: Arial

Code requirements:
1. Import the necessary libraries.
2. Create the slides and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
3. According to the detailed description, add the title, bullet points, and images at specified positions; set fonts and styles; explicitly set the size of each element to prevent overlap/occlusion; ensure text wraps automatically.
4. Use a Blank layout for all slides; place titles, content, and images as regular elements.
5. Save the file as: \"{save_dir}/{ppt_section.name}_slide_{slide_index + 1}.pptx\"

Please provide complete, executable Python code based on this information. Note: output Python code only, do not output any other text.
Code will be save in utf-8 encoding.
        """

        code_response = await coder_model.ainvoke([
            SystemMessage(content=code_prompt),
            HumanMessage(content="Generate complete Python code.")
        ])

        python_code = code_response.content.replace("```python", "").replace("```", "").strip()

        script_path = save_dir / f"{ppt_section.name}_slide_{slide_index + 1}.py"
        # ✅ 文件写入放到线程里
        await asyncio.to_thread(script_path.write_text, python_code, encoding="utf-8")

        returncode, stdout, stderr = await run_script(script_path)

        if returncode == 0:
            execution_successful = True
            break
        else:
            print(python_code)  # 仅用于调试
            error_message = stderr
            previous_code = python_code
            print(f"Code execution failed. Attempting again in {attempt + 1}/3 attempts. Error message: {stderr}")

    if not execution_successful:
        return {
            "codes": [python_code],
            "path": "none",
            "title": slide_title,
            "points": slide_points,
        }

    pptx_path = save_dir / f"{ppt_section.name}_slide_{slide_index + 1}.pptx"
    # ✅ 绝对路径解析（会触发 cwd），放在线程里执行以避免阻塞事件循环
    pptx_abs = await asyncio.to_thread(lambda: str(pptx_path.resolve()))

    return {
        "codes": [python_code],
        "path": pptx_abs,
        "title": slide_title,
        "points": slide_points,
    }


# def ppt_to_image(slide_ppt_path, image_path):
#     """
#     使用 unoconv 将PPT幻灯片导出为图片
#     """
#     print(f"Convert the slide {slide_ppt_path} to an image {image_path}")
    
#     # 使用 unoconv 将 PPT 转换为 PNG 格式
#     # try:
#         # 使用 unoconv 命令行工具来将 ppt 文件转换为图片
#     command = [
#             "C:\\Windows\\unoconv.bat",  # 调用 unoconv 命令
#             "-f", "png",  # 转换为 png 格式
#             "-o", image_path,  # 输出路径
#             slide_ppt_path  # 输入的 PPT 文件路径
#     ]
        
#         # 调用 unoconv 命令
#     subprocess.run(command, check=True)
#     print(f"Conversion successful: {slide_ppt_path} -> {image_path}")
    
#     # except subprocess.CalledProcessError as e:
#     #     print(f"转换失败：{str(e)}")

def ppt_to_image(slide_ppt_path, image_path, soffice_path=None, timeout=120):
    """
    在 macOS (conda 环境) 使用 LibreOffice(soffice) 将 PPT/PPTX 导出为 PNG。
    - 如果 PPT 仅 1 页 -> 生成 image_path 指定的 PNG。
    - 如果 PPT 多页 -> 第 1 页用 image_path，其余页在同目录生成 *_slide_02.png、*_slide_03.png...
    返回：保存的图片绝对路径列表（str）。
    """
    slide_ppt_path = Path(slide_ppt_path).expanduser().resolve()
    image_path = Path(image_path).expanduser()
    image_dir = image_path.parent
    image_stem = image_path.stem
    image_dir.mkdir(parents=True, exist_ok=True)

    # 优先级：函数参数 > 环境变量 > PATH 中的 soffice > 默认 App 路径
    candidates = []
    if soffice_path:
        candidates.append(Path(soffice_path))
    if os.environ.get("SOFFICE_PATH"):
        candidates.append(Path(os.environ["SOFFICE_PATH"]))
    which = shutil.which("soffice")
    if which:
        candidates.append(Path(which))
    candidates.append(Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"))

    soffice_bin = next((p for p in candidates if p and p.exists()), None)
    if soffice_bin is None:
        raise FileNotFoundError(
            "未找到 LibreOffice 的 'soffice' 可执行文件。请先安装 LibreOffice，"
            "或设置环境变量 SOFFICE_PATH 指向 soffice。"
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        cmd = [
            str(soffice_bin),
            "--headless", "--norestore", "--invisible",
            "--nodefault", "--nofirststartwizard", "--nolockcheck",
            "--convert-to", "png:impress_png_Export",
            "--outdir", str(tmpdir),
            str(slide_ppt_path),
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice 转换失败：{proc.stderr or proc.stdout}")

        generated = sorted(tmpdir.glob("*.png"))
        if not generated:
            raise FileNotFoundError("转换后未在临时目录发现 PNG 文件。")

        saved = []
        if len(generated) == 1:
            dst = image_path
            shutil.move(str(generated[0]), str(dst))
            saved.append(str(dst.resolve()))
        else:
            # 多页：第 1 页用给定文件名，其他页自动编号
            first_dst = image_path
            shutil.move(str(generated[0]), str(first_dst))
            saved.append(str(first_dst.resolve()))
            for i, src in enumerate(generated[1:], start=2):
                dst = image_dir / f"{image_stem}_slide_{i:02d}.png"
                shutil.move(str(src), str(dst))
                saved.append(str(dst.resolve()))

    return saved

async def ppt_slide_to_image_and_validate(state, config):
    """
    将生成的PPT幻灯片转换为图片，并使用大模型按新评分标准检查布局合理性。
    评分阈值：Total Score 严格 > 60 通过，否则 retry；最多重试 max_retry_count 次。
    """
    SCORE_THRESHOLD = 60  # 严格 > 60 才通过

    # ---- 新评分 Prompt ----
    REVIEW_PROMPT = """
You are a professional PowerPoint reviewer with strong logical expression and aesthetic judgment. Your task is to evaluate a given slide according to the following dimensions. Please score each dimension based on the allocated points and provide an overall evaluation. Be concise, objective, and follow the exact output format.

## 1. Content Logic & Expression (40 points)

* **Clear Structure**: Logical flow, well-organized, easy to follow.
* **Completeness & Focus**: Covers all key points, comprehensive yet without redundancy.
* **Concise Expression**: Theme is prominent, wording precise and succinct, highlights are clear.
* **Content Relevance**: Texts and visuals are appropriate and aligned with the slide’s theme.

---

## 2. Visual Design & Creativity (30 points)

* **Aesthetic Consistency**: Overall style is clean and coherent, layout comfortable, color scheme harmonious and aligned with brand/theme.
* **Balanced Layout**: Proper ratio of text to visuals, avoids clutter or emptiness.
* **Creative Appeal**: Innovative design and effective use of visual elements to engage the audience.

---

## 3. Technical Standards (30 points)

* **Text Standards**: Font sizes are appropriate, formatting consistent across levels, text fits within text boxes.
* **Image Standards**: Images are clear, with proper size and aspect ratio, no distortion.
* **Page Standards**: Correct slide dimensions and proportions, content fits within the page, no overlap between elements.

**Output Format:**
```json
{
  "Content Logic & Expression": "x/40",
  "Visual Design & Creativity": "x/30",
  "Technical Standards": "x/30",
  "Total Score": "x/100",
  "Overall Evaluation": "Excellent/Good/Average/Poor",
  "Suggestions for Improvement": "Your brief suggestions here"
}
```
"""

    # ---- 配置 & 模型初始化 ----
    configurable = Configuration.from_runnable_config(config)
    planner_model_kwargs = configurable.planner_model_kwargs or {}
    planner_provider = get_config_value(configurable.planner_provider)
    planner_model_name = get_config_value(configurable.planner_model)
    planner_model_kwargs = get_config_value(configurable.planner_model_kwargs or {})

    planner_model = AzureChatOpenAI(
        model=configurable.planner_model,
        azure_endpoint=planner_model_kwargs["openai_api_base"],
        deployment_name=planner_model_kwargs["azure_deployment"],
        openai_api_version=planner_model_kwargs["openai_api_version"],
        temperature=0,
        max_tokens=2048
    )

    # ---- 读取状态 ----
    slide_ppt_path = state["path"]
    codes = state["codes"]
    title = state["title"]
    points = state["points"]
    enriched_points = state["enriched_points"]
    slide_detail = state["slide_detail"]
    max_retry_count = state.get("max_retry_count", 3)
    retry_count = state.get("retry_count", 0)
    path = state.get("path")

    print("Current Slide:", slide_ppt_path, " Current repetition count:", retry_count, "Max Retry:", max_retry_count)

    # ---- 达到最大重试次数 ----
    if retry_count >= max_retry_count:
        print("The maximum retry count has been reached. No further processing will be carried out.")
        generated_slide = PPTSlide(
            title=title,
            points=points,
            codes=codes,
            enriched_points=enriched_points,
            detail=slide_detail
        )
        return Command(update={"completed_slides": [generated_slide]}, goto=END)

    # ---- 无有效路径 ----
    if path == "none":
        return Command(update={"layout_valid": False, "retry_count": retry_count + 1}, goto="enrich_slide_content")

    # ---- 转换为图片 ----
    output_folder = os.path.dirname(slide_ppt_path)
    image_path = slide_ppt_path.replace(".pptx", ".png")
    print(f"Convert the slide {slide_ppt_path} to an image {image_path}")
    try:
        await asyncio.to_thread(ppt_to_image, slide_ppt_path, image_path)
    except Exception as e:
        print(f"Error converting PPT to image: {e}")
        return Command(update={"conversion_failed": True, "retry_count": retry_count + 1}, goto="enrich_slide_content")
    print(f"The slide transition to image was successful. The saved path is: {image_path}")

    # ---- 审查评分：>60 通过，否则 retry ----
    def _extract_json_dict(text: str):
        """从模型返回中提取JSON（兼容 ```json 代码块``` 或纯JSON）；失败时尽量解析 Total Score。"""
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
    response = await planner_model.ainvoke([
        SystemMessage(content=REVIEW_PROMPT),
        HumanMessage(content=[{
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
        }])
    ])
    print(f"[Review raw]: {response.content}")

    result_dict = _extract_json_dict(response.content or "")
    total_score_str = str(result_dict.get("Total Score", "-1"))
    try:
        total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
    except Exception:
        total_score = -1.0
    print(f"[Parsed total score]: {total_score}")
    if "Suggestions for Improvement" in result_dict:
        print(f"[Suggestions]: {result_dict['Suggestions for Improvement']}")

    if total_score > SCORE_THRESHOLD:
        generated_slide = PPTSlide(
            title=title,
            points=points,
            codes=codes,
            enriched_points=enriched_points,
            detail=slide_detail
        )
        return Command(update={"completed_slides": [generated_slide]}, goto=END)
    else:
        retry_count += 1
        if retry_count >= max_retry_count:
            generated_slide = PPTSlide(
                title=title,
                points=points,
                codes=codes,
                enriched_points=enriched_points,
                detail=slide_detail
            )
            return Command(update={"completed_slides": [generated_slide]}, goto=END)
        return Command(update={"layout_valid": False, "retry_count": retry_count}, goto="enrich_slide_content")




async def generate_ppt_section_start(state: PPTSectionState):
    """
    为每个PPT章节开始生成幻灯片时，初始化一个空的幻灯片列表，并为每页PPT启动一个子图。

    Args:
        state: 包含章节信息、PPT章节信息的状态

    Returns:
        Dict: 包含空的 `generated_slides` 列表，并进入 `ppt_slide_subgraph`
    """
    # 从状态中获取章节信息和PPT章节信息
    topic = state["topic"]  # 主题
    # section = state["section"]  # 章节信息
    ppt_section = state["ppt_section"]  # PPT章节信息
    style = state.get("style")
    main_color = state.get("main_color")  # 主色
    accent_color = state.get("accent_color")  # 辅助

    # 初始化幻灯片列表，准备开始生成幻灯片
    generated_slides = []

    # 打印相关信息（用于调试）
    # print(f"开始生成PPT章节：{ppt_section.name}")

    # 获取该章节需要的页数
    num_slides = ppt_section.allocated_slides  # 获取分配的页数
    
    # 为每一页PPT启动一个子图
    return Command(
        update={"generated_slides": generated_slides},
        goto=[
            Send("generate_slide", {
                "topic": topic, 
                # "section": section, 
                "style": style,
                "ppt_section": ppt_section, 
                "slide_index": slide_index,  # 为每一页传递幻灯片的索引
                "main_color": main_color,
                "accent_color": accent_color,
                "max_retry_count": 3,  # 设置最大重试次数
                "retry_count": 0  # 初始化重试计数
            })
            for slide_index in range(num_slides)  # 根据分配的页数启动相应数量的子图
        ]
    )


async def generate_ppt_section_end(state: PPTSectionState):
    ppt_section = state["ppt_section"]
    ppt_section.slides = state["completed_slides"]

    return Command(
        update={"completed_ppt_sections": [ppt_section]},
        goto=END
    )

async def generate_cover_slide(state, config):
    """生成封面幻灯片，包含布局评分检查，最多循环3次（Total Score > 60 通过）"""
    # ----------------- helper: extract JSON dict from LLM output -----------------
    def _extract_json_dict(text: str) -> Dict[str, Any]:
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    # ----------------- review prompt -----------------
    REVIEW_PROMPT = """
You are a professional PowerPoint reviewer with strong logical expression and aesthetic judgment. Your task is to evaluate a given slide according to the following dimensions. Please score each dimension based on the allocated points and provide an overall evaluation. Be concise, objective, and follow the exact output format.

## 1. Content Logic & Expression (40 points)

* **Clear Structure**: Logical flow, well-organized, easy to follow.
* **Completeness & Focus**: Covers all key points, comprehensive yet without redundancy.
* **Concise Expression**: Theme is prominent, wording precise and succinct, highlights are clear.
* **Content Relevance**: Texts and visuals are appropriate and aligned with the slide’s theme.

---

## 2. Visual Design & Creativity (30 points)

* **Aesthetic Consistency**: Overall style is clean and coherent, layout comfortable, color scheme harmonious and aligned with brand/theme.
* **Balanced Layout**: Proper ratio of text to visuals, avoids clutter or emptiness.
* **Creative Appeal**: Innovative design and effective use of visual elements to engage the audience.

---

## 3. Technical Standards (30 points)

* **Text Standards**: Font sizes are appropriate, formatting consistent across levels, text fits within text boxes.
* **Image Standards**: Images are clear, with proper size and aspect ratio, no distortion.
* **Page Standards**: Correct slide dimensions and proportions, content fits within the page, no overlap between elements.

**Output Format:**
```json
{
  "Content Logic & Expression": "x/40",
  "Visual Design & Creativity": "x/30",
  "Technical Standards": "x/30",
  "Total Score": "x/100",
  "Overall Evaluation": "Excellent/Good/Average/Poor",
  "Suggestions for Improvement": "Your brief suggestions here"
}
```
"""
    SCORE_THRESHOLD = 60  # strictly > 60 to pass

    # ----------------- state -----------------
    topic = state["topic"]
    style = state.get("style", "none")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")

    save_dir = os.path.join(".", "saves", topic)
    cover_path = os.path.join(save_dir, "cover_slide.pptx")
    script_path = os.path.join(save_dir, "cover_slide.py")

    # ----------------- model init -----------------
    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})

    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960
    )

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"Timeout exception: {str(e)}"
            except Exception as e:
                return -1, "", f"Exception: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    for attempt in range(3):
        # ----------------- ask model for layout description -----------------
        layout_prompt = f"""
Please design a slide cover layout with the title: {topic}
Slide style: {style}
Primary color: {main_color}
Accent color: {accent_color}

The slide cover should include the title, speaker name, and date as key information.

{f"Previous error message: {error_message}" if error_message else ""}
{f"Previously generated code: {previous_code}" if previous_code else ""}
        """
        layout_response = await coder_model.ainvoke([
            {"role": "system", "content": layout_prompt},
            {"role": "user", "content": "Generate a layout description"}
        ])
        layout_description = layout_response.content.strip()

        # ----------------- ask model to generate python-pptx code -----------------
        code_generation_prompt = f"""
Generate complete Python code using the python-pptx library to create a cover slide based on the following layout description:

Layout description: {layout_description}

Code requirements:
1. Import the necessary libraries.
2. Create the slide and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
4. Save the PPT file to: {cover_path}

{f"Previous error message: {error_message}" if error_message else ""}
{f"Previously generated code: {previous_code}" if previous_code else ""}

Please provide complete, executable Python code based on this information. Note: output Python code only; do not output any other text.
Code will be save in utf-8 encoding.
        """
        code_response = await coder_model.ainvoke([
            {"role": "system", "content": code_generation_prompt},
            {"role": "user", "content": "Generate Python code"}
        ])
        python_code = code_response.content.replace("```python", "").replace("```", "").strip()
        await asyncio.to_thread(lambda: open(script_path, "w", encoding="utf-8").write(python_code))

        # ----------------- run script -----------------
        returncode, stdout, stderr = await run_script(script_path)
        if returncode != 0:
            error_message = stderr
            previous_code = python_code
            print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
            continue

        # ----------------- convert to image -----------------
        image_path = cover_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, cover_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        # ----------------- review with scoring -----------------
        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        review = await coder_model.ainvoke([
            SystemMessage(content=REVIEW_PROMPT),
            HumanMessage(content=[{
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
            }])
        ])

        print(f"[Cover Review raw] {review.content}")

        result = _extract_json_dict(review.content or "")
        total_score_str = str(result.get("Total Score", "-1"))
        try:
            total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
        except Exception:
            total_score = -1.0

        print(f"[Parsed score] {total_score} / 100")
        if "Suggestions for Improvement" in result:
            print(f"[Suggestions] {result['Suggestions for Improvement']}")

        if total_score > SCORE_THRESHOLD:
            print(f"封面布局评分通过（{total_score}），路径：{cover_path}")
            return {"cover_slide_path": cover_path, "cover_layout_description": layout_description, "score": total_score}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"封面布局评分未通过，尝试重试：{attempt + 1}")

    # 最终返回（达到最大次数仍未通过）
    return {"cover_slide_path": cover_path, "cover_layout_description": layout_description, "score": total_score}


async def generate_section_cover_slides(state, config):
    """生成章节封面幻灯片，仅检查第一个章节的布局有效性，最多循环3次（Total Score > 60 通过）"""
    SCORE_THRESHOLD = 60  # strictly > 60 to pass

    def _extract_json_dict(text: str) -> Dict[str, Any]:
        """从模型返回中提取JSON（兼容 ```json 代码块``` 或纯JSON）；失败则尽量解析 Total Score。"""
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    REVIEW_PROMPT = """
You are a professional PowerPoint reviewer with strong logical expression and aesthetic judgment. Your task is to evaluate a given slide according to the following dimensions. Please score each dimension based on the allocated points and provide an overall evaluation. Be concise, objective, and follow the exact output format.

## 1. Content Logic & Expression (40 points)

* **Clear Structure**: Logical flow, well-organized, easy to follow.
* **Completeness & Focus**: Covers all key points, comprehensive yet without redundancy.
* **Concise Expression**: Theme is prominent, wording precise and succinct, highlights are clear.
* **Content Relevance**: Texts and visuals are appropriate and aligned with the slide’s theme.

---

## 2. Visual Design & Creativity (30 points)

* **Aesthetic Consistency**: Overall style is clean and coherent, layout comfortable, color scheme harmonious and aligned with brand/theme.
* **Balanced Layout**: Proper ratio of text to visuals, avoids clutter or emptiness.
* **Creative Appeal**: Innovative design and effective use of visual elements to engage the audience.

---

## 3. Technical Standards (30 points)

* **Text Standards**: Font sizes are appropriate, formatting consistent across levels, text fits within text boxes.
* **Image Standards**: Images are clear, with proper size and aspect ratio, no distortion.
* **Page Standards**: Correct slide dimensions and proportions, content fits within the page, no overlap between elements.

**Output Format:**
```json
{
  "Content Logic & Expression": "x/40",
  "Visual Design & Creativity": "x/30",
  "Technical Standards": "x/30",
  "Total Score": "x/100",
  "Overall Evaluation": "Excellent/Good/Average/Poor",
  "Suggestions for Improvement": "Your brief suggestions here"
}
```
"""

    # ---- state ----
    topic = state["topic"]
    ppt_sections = state["ppt_sections"]
    style = state.get("style", "none")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")

    save_dir = os.path.join(".", "saves", topic)
    os.makedirs(save_dir, exist_ok=True)
    script_path = os.path.join(save_dir, "section_cover_slide.py")

    # ---- model init ----
    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})

    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960
    )

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=120
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    chapters_list = [section.name for section in ppt_sections]

    for attempt in range(3):
        # ---- ask for a generic layout for all chapter covers ----
        layout_prompt = f"""
    Please design a universal layout for slide chapter covers, applicable to all chapters.
    Slide topic: {topic}
    Slide style: {style}
    Primary color: {main_color}
    Accent color: {accent_color}

    Please return a layout description including element positions, sizes, and font information.

    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}
        """
        layout_response = await coder_model.ainvoke([
            {"role": "system", "content": layout_prompt},
            {"role": "user", "content": "Generate a layout description"}
        ])
        layout_description = layout_response.content.strip()

        # ---- generate python script to create all chapter cover PPTX files ----
        code_generation_prompt = f"""
    Generate a Python script using the python-pptx library to create chapter cover slides based on the following layout description:

    Layout description: {layout_description}

    Script requirements:
    1. Import the necessary libraries.
    2. Create slides and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
    3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
    4. Using the chapter list below, generate a separate PPTX file for each chapter. Save each file to {save_dir}/section_slide_{{chapter_index}}.pptx.

    Chapter list: {chapters_list}

    Please provide complete, executable Python code based on this information. Note: output Python code only; do not output any other text. Code will be save in utf-8 encoding.
        """
        code_response = await coder_model.ainvoke([
            {"role": "system", "content": code_generation_prompt},
            {"role": "user", "content": "Generate Python code"}
        ])
        python_code = code_response.content.replace("```python", "").replace("```", "").strip()
        await asyncio.to_thread(lambda: open(script_path, "w", encoding="utf-8").write(python_code))

        # ---- run script to generate PPTs ----
        returncode, stdout, stderr = await run_script(script_path)
        if returncode != 0:
            error_message = stderr
            previous_code = python_code
            print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
            continue

        # ---- validate ONLY the first chapter's PPT by scoring ----
        first_slide_path = os.path.join(save_dir, "section_slide_1.pptx")
        image_path = first_slide_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, first_slide_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        review = await coder_model.ainvoke([
            SystemMessage(content=REVIEW_PROMPT),
            HumanMessage(content=[{
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
            }])
        ])

        print(f"[Section Cover Review raw] {review.content}")

        result = _extract_json_dict(review.content or "")
        total_score_str = str(result.get("Total Score", "-1"))
        try:
            total_score = float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
        except Exception:
            total_score = -1.0

        print(f"[Parsed score] {total_score} / 100")
        if "Suggestions for Improvement" in result:
            print(f"[Suggestions] {result['Suggestions for Improvement']}")

        if total_score > SCORE_THRESHOLD:
            print(f"第一个章节布局评分通过（{total_score}），目录：{save_dir}")
            return {"section_slides_path": save_dir, "layout_description": layout_description, "score": total_score}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"第一个章节布局评分未通过，尝试重试：{attempt + 1}")

    # 最终返回（达到最大次数仍未通过）
    return {"section_slides_path": save_dir, "layout_description": layout_description, "score": total_score}



async def generate_end_slide(state, config):
    """生成封底幻灯片，包含布局检查步骤，最多循环3次（评分>60通过，否则重试）"""
    SCORE_THRESHOLD = 60  # 严格“>60”才通过

    def _extract_json_dict(text: str) -> Dict[str, Any]:
        m = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.S)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
        try:
            return json.loads(text)
        except Exception:
            pass
        m = re.search(r'"Total Score"\s*:\s*"?(?P<num>\d+(\.\d+)?)', text)
        if m:
            return {"Total Score": m.group("num")}
        return {}

    REVIEW_PROMPT = """
You are a professional PowerPoint reviewer with strong logical expression and aesthetic judgment. Your task is to evaluate a given slide according to the following dimensions. Please score each dimension based on the allocated points and provide an overall evaluation. Be concise, objective, and follow the exact output format.

## 1. Content Logic & Expression (40 points)

* **Clear Structure**: Logical flow, well-organized, easy to follow.
* **Completeness & Focus**: Covers all key points, comprehensive yet without redundancy.
* **Concise Expression**: Theme is prominent, wording precise and succinct, highlights are clear.
* **Content Relevance**: Texts and visuals are appropriate and aligned with the slide’s theme.

---

## 2. Visual Design & Creativity (30 points)

* **Aesthetic Consistency**: Overall style is clean and coherent, layout comfortable, color scheme harmonious and aligned with brand/theme.
* **Balanced Layout**: Proper ratio of text to visuals, avoids clutter or emptiness.
* **Creative Appeal**: Innovative design and effective use of visual elements to engage the audience.

---

## 3. Technical Standards (30 points)

* **Text Standards**: Font sizes are appropriate, formatting consistent across levels, text fits within text boxes.
* **Image Standards**: Images are clear, with proper size and aspect ratio, no distortion.
* **Page Standards**: Correct slide dimensions and proportions, content fits within the page, no overlap between elements.

**Output Format:**
```json
{
  "Content Logic & Expression": "x/40",
  "Visual Design & Creativity": "x/30",
  "Technical Standards": "x/30",
  "Total Score": "x/100",
  "Overall Evaluation": "Excellent/Good/Average/Poor",
  "Suggestions for Improvement": "Your brief suggestions here"
}
```
"""

    topic = state["topic"]
    style = state.get("style", "专业商务")
    main_color = state.get("main_color", "#FFFFFF")
    accent_color = state.get("accent_color", "#000000")

    save_dir = os.path.join(".", "saves", topic)
    os.makedirs(save_dir, exist_ok=True)
    end_path = os.path.join(save_dir, "end_slide.pptx")
    script_path = os.path.join(save_dir, "end_slide.py")

    configurable = Configuration.from_runnable_config(config)
    coder_model_kwargs = configurable.coder_model_kwargs or {}
    coder_provider = get_config_value(configurable.coder_provider)
    coder_model_name = get_config_value(configurable.coder_model)
    coder_model_kwargs = get_config_value(configurable.coder_model_kwargs or {})

    coder_model = AzureChatOpenAI(
        model=configurable.coder_model,
        azure_endpoint=coder_model_kwargs["openai_api_base"],
        deployment_name=coder_model_kwargs["azure_deployment"],
        openai_api_version=coder_model_kwargs["openai_api_version"],
        max_completion_tokens=40960
    )

    async def run_script(script_path):
        def run():
            try:
                proc_result = subprocess.run(
                    ["python", script_path],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                return proc_result.returncode, proc_result.stdout, proc_result.stderr
            except subprocess.TimeoutExpired as e:
                return -1, "", f"超时异常: {str(e)}"
            except Exception as e:
                return -1, "", f"其他异常: {str(e)}"
        return await asyncio.to_thread(run)

    error_message = ""
    previous_code = ""

    for attempt in range(3):
        layout_prompt = f"""
    Please design a layout for the slide's end (closing) page.
    Slide topic: {topic}
    Slide style: {style}
    Primary color: {main_color}
    Accent color: {accent_color}

    Please return a layout description including element content, positions, sizes, and font information.

    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}
        """

        layout_response = await coder_model.ainvoke([
            {"role": "system", "content": layout_prompt},
            {"role": "user", "content": "Generate a layout description"}
        ])
        layout_description = layout_response.content.strip()

        code_generation_prompt = f"""
    Generate complete Python code using the python-pptx library to create the closing slide based on the following layout description:

    Layout description: {layout_description}

    Code requirements:
    1. Import the necessary libraries.
    2. Create the slide and ensure the widescreen standard aspect ratio: 16:9 (13.33 inches × 7.5 inches).
    3. Use a rectangle the same size as the page to set the background; do not set slide.background directly.
    4. Save the PPT file to: {end_path}

    {f"Previous error message: {error_message}" if error_message else ""}
    {f"Previously generated code: {previous_code}" if previous_code else ""}

    Please provide complete, executable Python code based on this information. Note: output Python code only; do not output any other text.Code will be save in utf-8 encoding.
        """

        code_response = await coder_model.ainvoke([
            {"role": "system", "content": code_generation_prompt},
            {"role": "user", "content": "Generate Python code"}
        ])
        python_code = code_response.content.replace("```python", "").replace("```", "").strip()

        await asyncio.to_thread(
            lambda: open(script_path, "w", encoding="utf-8").write(python_code)
        )

        returncode, stdout, stderr = await run_script(script_path)
        if returncode != 0:
            error_message = stderr
            previous_code = python_code
            print(f"尝试{attempt + 1}失败，错误信息：{stderr}")
            continue

        image_path = end_path.replace(".pptx", ".png")
        try:
            await asyncio.to_thread(ppt_to_image, end_path, image_path)
        except Exception as e:
            print(f"幻灯片转图片失败: {e}")
            error_message = f"幻灯片转图片失败: {e}"
            previous_code = python_code
            continue

        # === 使用新的审查 Prompt 进行评分 ===
        image_content = await asyncio.to_thread(lambda: open(image_path, "rb").read())
        review_response = await coder_model.ainvoke([
            SystemMessage(content=REVIEW_PROMPT),
            HumanMessage(content=[{
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{base64.b64encode(image_content).decode()}"}
            }])
        ])

        print(f"[Review raw] {review_response.content}")

        # 解析评分JSON（严格>60通过）
        def _parse_total_score(payload: str) -> float:
            d = _extract_json_dict(payload)
            total_score_str = str(d.get("Total Score", "-1"))
            try:
                return float(re.search(r"(\d+(\.\d+)?)", total_score_str).group(1))
            except Exception:
                return -1.0

        total_score = _parse_total_score(review_response.content)
        print(f"[Parsed score] {total_score} / 100")

        if total_score > SCORE_THRESHOLD:
            print(f"布局评分通过（{total_score}），路径：{end_path}")
            return {"end_slide_path": end_path, "end_layout_description": layout_description, "score": total_score}
        else:
            error_message = f"评分未达阈值（{total_score} ≤ 60）"
            previous_code = python_code
            print(f"布局评分未通过，第 {attempt + 1} 次，准备重试。")

    return {"end_slide_path": end_path, "end_layout_description": layout_description, "score": total_score}



async def compile_ppt(state: ReportState, config: RunnableConfig):
    """
    异步方式合并所有生成的pptx文件到一个pptx文件中，包含封面、章节封面和封底。

    Args:
        state: 当前状态，包含所有完成的幻灯片信息。

    Returns:
        Dict: 更新后的状态，包含最终合并的PPT路径。
    """
    topic = state["topic"]
    ppt_sections = state["ppt_sections"]

    save_dir = os.path.join(".", "saves", topic)
    final_ppt_path = os.path.join(save_dir, f"{topic}_final.pptx")

    def merge_ppts_sync():
        final_presentation = Presentation()
        final_presentation.slide_width = Inches(13.33)
        final_presentation.slide_height = Inches(7.5)

        paths = ["cover_slide.pptx"]
        for idx, section in enumerate(ppt_sections):
            # 加入章节封面
            paths.append(f"section_slide_{idx+1}.pptx")
            # 加入该章节所有幻灯片
            paths.extend([
                f"{section.name}_slide_{slide_index}.pptx"
                for slide_index, _ in enumerate(section.slides, start=1)
            ])
        paths.append("end_slide.pptx")

        for ppt_file in paths:
            ppt_path = os.path.join(save_dir, ppt_file)
            if not os.path.exists(ppt_path):
                continue

            presentation = Presentation(ppt_path)
            for slide in presentation.slides:
                new_slide = final_presentation.slides.add_slide(final_presentation.slide_layouts[6])
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # 图片
                        image_stream = io.BytesIO(shape.image.blob)
                        new_slide.shapes.add_picture(
                            image_stream, shape.left, shape.top, shape.width, shape.height
                        )
                    else:
                        new_slide.shapes._spTree.insert_element_before(shape.element, 'p:extLst')

        final_presentation.save(final_ppt_path)

    await asyncio.to_thread(merge_ppts_sync)

    return {"final_ppt_path": final_ppt_path}




# PPT Slide sub-graph -- 
ppt_slide_graph = StateGraph(PPTSlideState, output=PPTSlideOutputState)
# ppt_slide_graph.add_node("generate_slide_content", generate_slide_content)
ppt_slide_graph.add_node("enrich_slide_content", enrich_slide_content)
ppt_slide_graph.add_node("generate_slide_code_and_execute", generate_slide_code_and_execute)
ppt_slide_graph.add_node("ppt_slide_to_image_and_validate", ppt_slide_to_image_and_validate)

ppt_slide_graph.add_edge(START, "enrich_slide_content")
ppt_slide_graph.add_edge("enrich_slide_content", "generate_slide_code_and_execute")
ppt_slide_graph.add_edge("generate_slide_code_and_execute", "ppt_slide_to_image_and_validate")
# ppt_slide_graph.add_edge("generate_slide_content", END)

ppt_slide_subgraph = ppt_slide_graph.compile()

# PPT Sction sub-graph --
ppt_section_graph = StateGraph(PPTSectionState,output=PPTSectionOutputState)
ppt_section_graph.add_node("generate_ppt_section_start", generate_ppt_section_start)
ppt_section_graph.add_node("generate_slide", ppt_slide_subgraph)
ppt_section_graph.add_node("generate_ppt_section_end", generate_ppt_section_end)

ppt_section_graph.add_edge(START, "generate_ppt_section_start")
# ppt_section_graph.add_edge("generate_ppt_section_start", "generate_slide")
ppt_section_graph.add_edge("generate_slide", "generate_ppt_section_end")
# ppt_section_graph.add_edge("generate_ppt_section_end", END)

ppt_section_subgraph = ppt_section_graph.compile()

# Report section sub-graph -- 

# Add nodes 
section_builder = StateGraph(SectionState, output=SectionOutputState)
section_builder.add_node("generate_queries", generate_queries)
section_builder.add_node("search_web", search_web)
section_builder.add_node("write_section", write_section)

# Add edges
section_builder.add_edge(START, "generate_queries")
section_builder.add_edge("generate_queries", "search_web")
section_builder.add_edge("search_web", "write_section")
# section_builder.add_edge("write_section", END)

# Outer graph for initial report plan compiling results from each section -- 

# Add nodes
builder = StateGraph(ReportState, input=ReportStateInput, output=ReportStateOutput, config_schema=Configuration)
builder.add_node("process_image_input", process_image_input)
builder.add_node("generate_report_plan", generate_report_plan)
builder.add_node("human_feedback", human_feedback)
builder.add_node("build_section_with_web_research", section_builder.compile())
builder.add_node("gather_completed_sections", gather_completed_sections)
builder.add_node("write_final_sections", write_final_sections)
builder.add_node("compile_final_report", compile_final_report)
builder.add_node("generate_ppt_outline", generate_ppt_outline)
builder.add_node("generate_ppt_sections", ppt_section_subgraph)
builder.add_node("generate_cover_slide", generate_cover_slide)
builder.add_node("generate_section_cover_slides", generate_section_cover_slides)
builder.add_node("generate_end_slide", generate_end_slide)
builder.add_node("compile_ppt", compile_ppt)

# Add edges
builder.add_edge(START, "process_image_input")
builder.add_edge("process_image_input", "generate_report_plan")
builder.add_edge("generate_report_plan", "human_feedback")
builder.add_edge("build_section_with_web_research", "gather_completed_sections")
builder.add_conditional_edges("gather_completed_sections", initiate_final_section_writing, ["write_final_sections"])
builder.add_edge("write_final_sections", "compile_final_report")
# builder.add_edge("compile_final_report", END)
builder.add_edge("compile_final_report", "generate_ppt_outline")
# builder.add_edge("generate_ppt_outline", "generate_ppt_sections")
# builder.add_edge("generate_ppt_sections", "compile_ppt")
builder.add_edge("generate_ppt_sections", "generate_cover_slide")
builder.add_edge("generate_cover_slide", "generate_section_cover_slides")
builder.add_edge("generate_section_cover_slides", "generate_end_slide")
builder.add_edge("generate_end_slide", "compile_ppt")
builder.add_edge("compile_ppt", END)

graph = builder.compile()
