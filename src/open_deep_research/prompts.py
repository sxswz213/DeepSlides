report_planner_query_writer_instructions = """You are conducting research for a report.

<Report Topic>
{topic}
</Report Topic>

<Description of the user-provided image>
{caption}
</Description of the user-provided image>

<Possible User Intent>
{user_intent}
</Possible User Intent>

<Report Organization>
{report_organization}
</Report Organization>

<Task>
Your goal is to generate {number_of_queries} web search queries to help collect the information needed to plan each part of the report.

These queries should:

1. Be relevant to the report topic and the user's intent  
2. Be based on the image description, ensuring the queries are closely related to the image content  
3. Help meet the requirements specified in the report organization
4. **No more than 5 words** (each query)

Make sure the queries are specific enough to find high-quality, relevant resources while still covering the breadth required by the report structure.
</Task>

<Format>
Call the Queries tool
</Format>
"""


report_planner_instructions = """I need a concise, focused report plan.

<Report Topic>
The topic of the report is:
{topic}
</Report Topic>

<Description of the user-provided image>
{caption}
</Description of the user-provided image>

<Possible User Intent>
{user_intent}
</Possible User Intent>

<Report Organization>
The report should follow this structure:
{report_organization}
</Report Organization>

<Context>
Here is the contextual information for planning each part of the report:
{context}
</Context>

<Task>
Generate a list of sections for the report. Your plan should be compact and focused, avoiding overlapping sections or unnecessary filler.

For example, a good report structure might look like:
1/ Introduction
2/ Overview of Topic A
3/ Overview of Topic B
4/ Comparison of A and B
5/ Conclusion

Each section must include these fields:

* Name - the name of the report section.  
* Description - a brief overview of the main topics covered in that section.  
* Research - whether web research is required for this section. **Important:** Core sections (not Introduction/Conclusion) must have Research=True. A useful report must include at least 2-3 sections with Research=True.  
* Content - the section content, left blank for now.  
* Source_str - the search-query string related to this section, left blank for now.  

Guidelines for integration:

* Include examples and implementation details within topical sections, not as standalone sections.  
* Ensure each section has a clear purpose; avoid overlapping content.  
* Merge related concepts instead of treating them separately.  
* Critical: every section must directly relate to the topic and user intent.  
* Avoid off-topic or loosely related sections that are not directly connected to the core topic and intent.

Before submitting, review your structure to be sure there are no redundant sections and that it follows a logical flow.
</Task>

<Feedback>
Here is reviewer feedback on the report structure (if any):
{feedback}
</Feedback>

<Format>
Call the Sections tool
</Format>
"""



query_writer_instructions = """You are a professional technical writing expert creating targeted web search queries to gather comprehensive information for writing a technical report section.

<Report Topic>
{topic}
</Report Topic>

<Section Topic>
{section_topic}
</Section Topic>

<Task>
Your goal is to generate {number_of_queries} search queries to help collect thorough information on the section topic.

These queries should:

1. Be relevant to the topic  
2. Explore different facets of the topic
3. **No more than 5 words** (each query)

Ensure the queries are specific enough to find high-quality, relevant resources.
</Task>

<Format>
Call the Queries tool
</Format>
"""


section_writer_instructions = """Write one section of a research report.

<Task>
1. Carefully read the report topic, section name, and section theme.  
2. If existing section content is provided, review it.  
3. Then examine the provided sources.  
4. Decide which sources you will use to write the report section.  
5. If an image list is provided, evaluate each image's relevance to the section and select the single most suitable one.  
6. Write the report section and choose an illustration.  
7. List the sources for the section at the end.
</Task>

<Writing Guidelines>
- If the existing section content is empty, write it from scratch.  
- If existing content is present, integrate it with the sources.  
- Strictly limit the length to 150-200 words.  
- Use clear, simple language.  
- Use short paragraphs (no more than 2-3 sentences).  
- Use "##" for the section heading (Markdown format).
</Writing Guidelines>

<Image Handling Guidelines>
- If an image list is provided (up to 6 images), choose the one that best supports or illustrates the section content.  
- The return should include two parts: the normal section content and the image selection.  
- Image-selection format:

```image_selection
{
  "selected_image_index": index of the chosen image (starting at 0),
  "reason": "Why this image was chosen",
  "caption": "A brief caption for the image"
}
````

* If no suitable image exists or no images are provided, set "selected\_image\_index" to -1.
* The section content must stand alone; do not reference "as shown in the image" in the text.
* The image-selection block should come after the section content but before the source list.
  \</Image Handling Guidelines>

<Citation Rules>
- Assign a reference number to each unique URL.  
- End with ### Sources and list each source with its number.  
- Important: Number sources sequentially (1, 2, 3, 4 …) with no gaps, regardless of which are selected.  
- Example format:  
  [1] Source Title: URL  
  [2] Source Title: URL
</Citation Rules>

<Final Checks>
1. Verify each point is supported by the provided sources.  
2. Confirm each URL appears only once in the source list.  
3. Ensure sources are numbered in order (1, 2, 3 …) with no gaps.  
4. Make sure the image-selection block (if any) follows the required JSON format.
</Final Checks>
"""



section_writer_inputs = """
<Report Topic>
{topic}
</Report Topic>

<Section Name>
{section_name}
</Section Name>

<Section Theme>
{section_topic}
</Section Theme>

<Existing Section Content (if any)>
{section_content}
</Existing Section Content>

<Sources>
{context}
</Sources>

<Image List>
{images_data}
</Image List>
"""


section_grader_instructions = """Review a report section against the specified topic:

<Report Topic>
{topic}
</Report Topic>

<Section Theme>
{section_topic}
</Section Theme>

<Section Content>
{section}
</Section Content>

<Task>
Evaluate whether the section content adequately covers the section theme.

If the content does not sufficiently cover the theme, generate {number_of_follow_up_queries} follow-up search queries to gather the missing information.
</Task>

<Format>
Call the Feedback tool and output using the pattern below:

grade: Literal["pass", "fail"] = Field(
    description="The evaluation result indicating whether the response meets the requirements ('pass') or needs revision ('fail')."
)
follow_up_queries: List[SearchQuery] = Field(
    description="A list of follow-up search queries.",
)
</Format>
"""


final_section_writer_instructions = """You are a professional technical writing expert tasked with composing a chapter that synthesizes the remaining information for the report.

<Report Topic>
{topic}
</Report Topic>

<Section Name>
{section_name}
</Section Name>

<Section Theme>
{section_topic}
</Section Theme>

<Available Report Content>
{context}
</Available Report Content>

<Task>
1. Section-specific guidance:

For the **Introduction**:
- Use "#" for the report title (Markdown format).  
- Limit to **50-100 words**.  
- Use clear, straightforward language.  
- Focus on the core motivation of the report in 1-2 paragraphs.  
- Employ a clear narrative structure to introduce the report.  
- Do **not** include structural elements (no lists or tables).  
- No sources section required.

For the **Conclusion/Summary**:
- Use "##" for the section heading (Markdown format).  
- Limit to **100-150 words**.  
- For comparative reports:  
  * **Must** include a concise comparison table using Markdown table syntax.  
  * The table should distill the insights of the report.  
  * Keep table entries clear and succinct.  
- For non-comparative reports:  
  * Include **one** structural element **only if** it helps distill key points:  
    • Either a table (Markdown) focusing on items mentioned in the report, **or**  
    • A short list using proper Markdown list syntax:  
      - Use `*` or `-` for unordered lists.  
      - Use `1.` for ordered lists.  
      - Ensure correct indentation and spacing.  
- End with concrete next steps or implications.  
- No sources section required.

3. Writing approach:
- Use specific details rather than general statements.  
- Make every word count.  
- Concentrate on your most important points.
</Task>

<Quality Check>
- **Introduction**: 50-100 words, "#" as title, no structural elements, no sources section.  
- **Conclusion**: 100-150 words, "##" as heading, at most one structural element, no sources section.  
- Use Markdown formatting.  
- Do **not** include word counts or any preamble in your reply.
</Quality Check>
"""



## Supervisor
SUPERVISOR_INSTRUCTIONS = """
You are responsible for conducting investigative research for a report based on the user-provided topic.

### Your Responsibilities:

1. **Gather Background Information**  
   Use `enhanced_tavily_search` to collect relevant information on the user’s topic.  
   - You must perform **exactly one** search to gather comprehensive context.  
   - Craft highly targeted search queries to obtain the most valuable information.  
   - Take time to analyze and synthesize the search results before proceeding.  
   - Do not move on until you have a clear understanding of the topic.

2. **Clarify the Topic**  
   After preliminary research, engage with the user to clarify any open questions.  
   - Ask concrete follow-up questions based on what you learned from the search.  
   - Do not continue until you fully understand the topic, goals, constraints, and any preferences.  
   - Summarize what you have learned so far before asking questions.  
   - You must have at least one clarification exchange with the user before proceeding.

3. **Define the Report Structure**  
   Only after research and clarification are completed:  
   - Use the `Sections` tool to define a list of report sections.  
   - Each section should include: a section name and a research plan for that section.  
   - **Do not** include Introduction or Conclusion sections (we will add these later).  
   - Ensure each section’s scope is suitable for independent research.  
   - Base the sections on search results and user clarifications.  
   - Format the sections as a list of strings, each string describing the research scope of that section.

4. **Assemble the Final Report**  
   When all sections have been returned:  
   - **Important:** First check your previous messages to see what you have already completed.  
   - If you have not yet created an introduction, generate one with the `Introduction` tool:  
     - Set the content, beginning with a single‐level H1 title (`#`) for the report.  
     - Example: `# [Report Title]`  
   - After the introduction, summarize key insights with the `Conclusion` tool:  
     - Set the content, beginning with an H2 heading (`##`) titled "Conclusion".  
     - Example: `## Conclusion\n\n[Conclusion content…]`  
     - Include **exactly one** structural element *only if* it helps distill key points:  
       - A table (Markdown) focusing on items compared in the report, **or**  
       - A short list using proper Markdown list syntax:  
         - Use `*` or `-` for unordered lists.  
         - Use `1.` for ordered lists.  
         - Ensure correct indentation and spacing.  
   - Do **not** call the same tool twice—check your message history.

### Additional Notes:
- You are a reasoning model. Think step-by-step before acting.  
- **Important:** Do not rush to create the report structure. Thoroughly collect information first.  
- Use multiple searches to build a complete picture before drawing conclusions.  
- Maintain a clear, informative, and professional tone throughout.
"""


RESEARCH_INSTRUCTIONS = """
You are a researcher responsible for completing a specific section of a report.

### Your Goals:

1. **Understand the Section Scope**  
   First, review the section's scope of work. This defines your research focus. Treat it as your objective.

<Section Description>
{section_description}
</Section Description>

2. **Strategic Research Process**  
   Follow this precise research strategy:

   a) **First Query**: Start with a single, well-crafted search query using `enhanced_tavily_search` that directly targets the core of the section topic.  
      - Formulate a focused query that will yield the most valuable information.  
      - Avoid generating multiple similar queries (e.g., "benefits of X", "advantages of X", "why use X").  
      - Example: "model context protocol developer advantages and use cases" is better than separate queries for advantages and use cases.

   b) **Thoroughly Analyze Results**: After receiving search results:  
      - Read and analyze all provided content carefully.  
      - Identify aspects already well covered and those requiring more information.  
      - Assess how the current information addresses the section scope.

   c) **Follow-up Research**: If needed, conduct targeted follow-up searches:  
      - Create a follow-up query that targets specific missing information.  
      - Example: If general benefits are covered but technical details are missing, search "model context protocol technical implementation details".  
      - Avoid redundant queries that return similar information.

   d) **Complete the Research**: Continue this focused process until you have:  
      - Comprehensive information covering all aspects of the section scope.  
      - At least 3 high-quality sources offering different perspectives.  
      - Both breadth (all aspects covered) and depth (specific details).

3. **Use the Section Tool**  
   Only after thorough research, write a high-quality section using the Section tool:  
   - `name`: Section title  
   - `description`: The scope of research you completed (brief, 1-2 sentences).  
   - `content`: The full body of the section, which must:  
     - Begin with "## [Section Title]" (H2 level, use ##).  
     - Be formatted in Markdown style.  
     - Be no more than 200 words (strict limit).  
     - End with a "### Sources" subsection (H3 level, use ###) listing numbered URLs used.  
     - Use clear, concise language and bullet points where appropriate.  
     - Include relevant facts, statistics, or expert opinions.

Content format example:
```

## \[Section Title]

\[Body text, Markdown format, up to 200 words...]

### Sources

1. \[URL 1]
2. \[URL 2]
3. \[URL 3]

```

---

### Research Decision Framework

Before each search query or writing the section, think:

1. **What information do I already have?**  
   - Review all information collected so far.  
   - Identify key insights and facts already uncovered.

2. **What information is still missing?**  
   - Determine the specific knowledge gaps related to the section scope.  
   - Prioritize the most important missing information.

3. **What is the most effective next step?**  
   - Decide whether another search is needed (and the specific aspect to search).  
   - Or determine if sufficient information has been collected to write a comprehensive section.

---

### Notes:
- Emphasize the **quality** of searches over quantity.  
- Each search should have a clear, distinct purpose.  
- Do not write introductions or conclusions unless they are explicitly part of your section.  
- Maintain a professional, factual tone at all times.  
- Always follow Markdown formatting.  
- Keep the main content within the 200-word limit.
"""



query_writer4PPT_instructions = """
You are an assistant highly skilled at generating relevant search queries for presentation slides.

<Slide Topic>
{topic}
</Slide Topic>

<Current Section Topic>
{section_topic}
</Current Section Topic>

<Current Slide Topic>
Title: {slide_title}
</Current Slide Topic>

<Slide Bullet Points>
Slide points: {slide_points}
</Slide Bullet Points>

<Task>
Your goal is to generate {number_of_queries} web search queries to help us find additional material to expand the slide’s points. Each query should be no more than **10 words**.

These queries should:
1. Be tightly related to the slide’s topic.  
2. Focus on the key points and relevant details mentioned in the slide bullet points.  
3. Be specific and clear enough to retrieve high-quality, relevant resources.
4. **No more than 5 words** (each query)

<Format>
Call the Queries tool
</Format>
"""

ppt_tools_prompt = """
First, complete the background layer, then complete the layout layer, and finally complete the content layer. Use the add_textbox function to place all the text. Do not use slide.shapes.add_textbox.
These are utility functions you can choose to use to reduce repetitive work in your PowerPoint slides. You don't need to use all the functions, just select the appropriate ones to complete your task. Please strictly follow the designed layout when coding; do not change the design to accommodate the utility functions. 
**You must use the "add_textbox" function provided in the given tools bellow, DO NOT use "slide.shapes.add_textbox".**
Attention: You can still choose to write python code manually without using these functions, or write code manually to accomplish parts that these functions cannot handle.

Before using, please add the import statement:
from pptx_tools.add_free_shape import *

Here are the descriptions of each function. When using the function, please clearly specify each field, for example: Field_Name = Value.

add_gradient_shape: Add a shape with gradient fill to a PowerPoint slide.
    Required:
    - slide: Slide object to add the shape to
    - left/top/width/height: Position and size (inches, use only float values such as 3.0, do not use Inches(3.0) )
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - color_start/color_end: Gradient colors in "#RRGGBB" format (default: "#FFD54F"/"#9E9E9E")
    - angle_deg: Gradient angle, 0=left→right, 90=bottom→top (default: 0)
    - method: "ooxml" (native, recommended) or "picture" (compatible) (default: "ooxml")
    - alpha_start/alpha_end: Opacity 0~1 (ooxml mode only)
    Returns: The added shape object

add_gradient_background: Apply a gradient background to an entire slide.
    Required:
    - slide: Slide object to modify
    Optional:
    - color_start/color_end: Gradient colors in "#RRGGBB" format (default: "#FFD54F"/"#9E9E9E")
    - angle_deg: Gradient angle, 0=left→right, 90=bottom→top (default: 90)
    - method: "ooxml" (native) or "picture" (compatible) (default: "ooxml")
    - alpha_start/alpha_end: Opacity 0~1 (ooxml mode only)
    - stops: Custom gradient stops [(pos,rgb,alpha),...] (overrides colors if set)
    Effect: Modifies slide background directly, no return value

add_solid_shape: Add a shape with solid color fill to a slide.
    Required:
    - slide: Slide object to add the shape to
    - left/top/width/height: Position and size (inches, use only float values such as 3.0 )
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - fill_color: Fill color in "#RRGGBB" format (default: "#FFD54F")
    - fill_alpha: Opacity 0~1 (default: 1.0)
    - outline_color: Border color "#RRGGBB" or None for no border
    - outline_width_pt: Border width in points (None or 0 for no border)
    Returns: The added shape object

add_image_filled_shape: Add a shape filled with an image to a slide.
    Required:
    - slide: Slide object to add the shape to
    - image_path: Path to the image file
    - left/top/width/height: Position and size (inches, use only float values such as 3.0 )
    Optional:
    - shape_type: Shape type from MSO_SHAPE enum (default: RECTANGLE)
    - mode: "stretch" (default) or "tile"
    - crop: Cropping (left,top,right,bottom) in 0~1 range
    - tile_scale: Scale factor for tile mode (default: 1.0)
    - remove_line: Whether to remove shape outline (default: True)
    Returns: The added shape object

add_textbox: Add a text box with / without base box and auto resize font size to adapt the box size. (Do not resize font after using this function)
    Required:
    - slide: Slide object to add to
    - text: Text content
    - left/top/width/height: Position and size (inches, use only float values such as 3.0, do not use Inches(3.0) )
    - font_size: Target font size in points (will auto-shrink if needed)
    Optional:
    - font_name: Font name (default: "Calibri")
    - font_color: Text color (default: #000000)
    - has_base_box: Whether to draw a base box / border (default: False)
    - border_color: Base box border color (default: #000000)
    - border_width: Border thickness in points (default: 2)
    - corner_radius: Corner rounding 0~1 (default: 0.1)
    - has_border: Whether to show border (default: True)
    - fill_color: Base box background color or None for transparent
    - inner_padding: Text padding in inches (default: 0.10)
    Returns: Tuple of (border shape, text box shape)

add_line: Add a line to a slide.
    Required:
    - slide: Slide object to add to
    - x1/y1: Start point position in inches (inchs, use only float values)
    - x2/y2: End point position in inches (inchs, use only float values)
    Optional:
    - color: Line color (default: #000000)
    - width_pt: Line thickness in points (default: 2)
    - dash_style: Line dash pattern keyword ("solid", "dashed", "dash_dot", "dash_dot_dot", "long_dash", "long_dash_dot", "dot", "square_dot", "mixed")
    Returns: Connector shape object representing the line
"""

design_formatting_prompt = """
# Canvas & Units
- Canvas: 13.33 × 7.5 inches (width × height)
- All coordinates and sizes are in **inches**, with 2 decimal places
- No elements may go out of bounds or overlap (except background textures/separators)
- In the Layout layer, every block uses **absolute positioning**: top-left `(x, y)`, size `(w, h)`;  
  `0 ≤ x ≤ 13.33`, `0 ≤ y ≤ 7.5`, `x + w ≤ 13.33`, `y + h ≤ 7.5`
# General Requirements
- **Only the provided image URLs can be used. Do not reserve any positions for any images that are not provided, and do not use text descriptions to fill the gaps.**
- **If there is no images provided, do not reference images in the design and do not leave extra space.**
- **The use of pictures or the creation of flowcharts is encouraged.**

# Layered Structure (put only the appropriate info in each layer)
## Background Layer (weakly coupled with content/layout; ideally reusable)
- One **natural-language summary** (≤ 80 words) describing style and visual tone
- One **machine-readable JSON** (see “Background JSON template”)
- If Tone is light, use a light-colored background; If it is dark, use a dark background.
- Attention: The elements added to the background layer should not interfere with the readability of the content layer.
- The font size of the main text should be at least 16.

## Layout Layer (blocks and separators only; no concrete text/image content)
- Layout summary (single block / top-bottom / left-right / n horizontal blocks / 2×2 grid / card grid); encourage diverse and innovative layouts.
- One **natural-language description** (≤ 120 words) explaining layout logic, alignment/whitespace/emphasis
- A **structure JSON**: keys “Block1/Block2/…” with function, position, size; nested sub-blocks allowed (Block1.1)
- For each block, you can use either a solid color, gradient color, or only the border of a 【rectangle/rounded rectangle/other shapes】 as the background. You can also use a dividing line to separate each block.
- Blocks of the same level should not overlap with each other. The children blocks need to be inside the parent blocks. The divider line should not cover the blocks.
- If the internal blocks (such as text boxes) have borders, then the external blocks should not use the same borders. Instead, they can be changed to solid backgrounds or dividing lines.
- Generally speaking, the background of the content block should be lighter than the background color, or only the border should be colored while the middle part is filled with white. If the block is dark, the text inside should be light-colored.
- You can appropriately use emojis as icons to embellish the slides. Please simply provide the emojis instead of using text placeholders.
- If the background color of the block is dark, then the text should be light. On the contrary, the text is in a dark color.

## Content Layer (text and image bindings only, no positioning)
- For each block id from the Layout layer, provide **text/image** content.
- The content of the picture can only be based on the provided picture URL.
- The images should be scaled proportionally. If this is not possible, a deformation within 20% in the aspect ratio is acceptable.

# Output Format (strict; output nothing else)
[Background]
  Preset background: none / <filename or URL>
  Tone: light/dark/neutral
  Primary/Accent: {main_color} / {accent_color}
  Base texture: <color + style, or "none">
  Edge ornament: <color + style, or "none">
  Natural-language description: <≤80 words>

[Layout]
  Summary: single / top-bottom / left-right / three-horizontal / 2x2-grid / card-grid / ...
  Natural-language details: <≤120 words>
  Structure: {
    "blocks":{
      "Block1":{"function":"title/text/image...",
                "background":{"type":"solid/gradient/shape/none",
                              "color":"#RRGGBB",
                              "details":{...}},
                "bbox_in":{"x":...,"y":...,"w":...,"h":...},
                "children":{}},
      "Block2":{"function":"...",
                "bbox_in":{...},
                "background":{...},
                "description":"...",
                "children":{
                  "Block2.1":{"function":"text","bbox_in":{"x":...,"y":...,"w":...,"h":...},"background":{...},"font_color":"#RRGGBB", "description":"..."},
                  "Block2.2":{...}
                }},
      "Block3":{...}
    },
    "separators":[
      {"type":"line",
      "bbox_in":{...},
      "color":"..."}
    ]
  }

[Content]
  {
    "Block1":{"text":"..."},
    "Block2":{
      "Block2.1":{"text":"..."},
      "Block2.2":{"image":"<id/url from <Image list> or 'none'>","caption":"optional"}
    },
    "Block3":{...}
  }
"""

evaluation_design = """
You are a Slide Layout Review Expert. Please evaluate the slide design based on the following dimensions.

A. Completeness: Whether all requested design requirements are properly reflected, and whether text and visuals match the provided content.

A.1 Design Element Consistency
Assess whether the color, style, font, and other requirements are correctly applied throughout the design.

A.2 Content Fidelity
Check whether all referenced text and images appear in the provided source content without omission or fabrication.

B. Compliance: Whether the designed visual and text blocks follow structural rules such as non-overlap and proper spatial arrangement.

B.1 Overlap Ratio
Check whether any text or visual elements unintentionally overlap with each other, excluding decorative backgrounds or stylistic framing elements.

B.2 Page Occupancy Ratio
Evaluate whether the total occupied area of all elements is appropriate and whether the layout is visually balanced without leaning excessively toward one side.

B.3 Overflow Ratio
Identify portions of any element that extend beyond the slide boundary and quantify the exceeded area.

Provide a score from 1 to 5 for each dimension, where 1 = poor, 2 = fair, 3 = good, 4 = very good, 5 = excellent.
OUTPUT FORMAT:
{
   "A.1 Design Element Consistency": score,
   "A.2 Content Fidelity": score,
   "B.1 Overlap Ratio": score,
   "B.2 Page Occupancy Ratio": score,
   "B.3 Overflow Ratio": score,
   "Total Score": score (score = 0.7 * (A.1 + A.2) / 2 + 0.3 * (B.1 + B.2 + B.3) / 3, only return final score),
   "Suggesstions": improvement suggestions
}
"""

evaluation_complete="""
You are an expert evaluator for Slide Design-to-Code generation tasks. You must assess whether the produced code accurately matches the provided slide design according to the following criteria:
(1) Every designed element appears in the code and is correctly configured in terms of color, position, shape, size, and other relevant attributes;
(2) Every requirement specified in the design specification is faithfully implemented in the code.

OUTPUT FORMAT:
{
   "Element Match": discription of correctly implemented design elements,
   "Element Mismatch": discription of incorrectly implemented or missing design elements,
   "Element Match Score": score (1-5, 1=poor, 5=excellent),
   "Requirement Compliance": discription of met requirements,
   "Requirement Non-Compliance": discription of unmet requirements,
   "Requirement Compliance Score": score (1-5, 1=poor, 5=excellent),
   "Total Score": total_score (average of above scores),
   "Suggesstions": improvement suggestions
}
"""

evaluation_aesthetics="""
You are a Slide Aesthetics Expert. Please evaluate the slide purely from a visual and aesthetic perspective (ignore content accuracy) across the following dimensions:

1. Layout & Composition
Whether the spatial arrangement is balanced, alignment is consistent, and spacing between elements is appropriate.

2. Visual Hierarchy
Whether visual weight is properly distributed, key elements stand out, and the viewing flow feels natural.

3. Color & Contrast
Whether the color palette is harmonious, contrasts are sufficient, and overall color usage feels cohesive.

4. Typography
Whether font selection, sizes, spacing, and text layout are visually appealing and easy to read.

5. Whitespace & Balance
Whether negative space is appropriately used and the slide feels neither overcrowded nor empty.

6. Overall Aesthetic Consistency
Whether shapes, colors, fonts, and stylistic elements follow a coherent and unified aesthetic style.


Provide a score from 1 to 5 for each dimension, where 1 = poor, 2 = fair, 3 = good, 4 = very good, 5 = excellent.

OUTPUT FORMAT:
{
   "Layout & Composition": score,
   "Visual Hierarchy": score,
   "Color & Contrast": score,
   "Typography": score,
   "Whitespace & Balance": score,
   "Overall Aesthetic Consistency": score,
   "Total Score": total_score (average of above scores),
   "Suggesstions": improvement suggestions
}
"""

style_plan_prompt="""
### **Color Style Guide for Slide Design**

Effective slide design relies on a clear and disciplined color system. A well-constructed palette generally includes one dominant color, one to two accent colors, and a set of neutral tones for balance. The dominant color establishes the visual identity of the slide, while the accent colors highlight key information such as keywords, icons, or numerical results. Neutral grays should be used for backgrounds, text blocks, and low-priority elements to prevent visual clutter.

For gradients, choose subtle transitions within the same hue family rather than mixing unrelated colors; this preserves consistency and avoids visual noise. When applying multiple colors on a single slide, maintain a contrast ratio that ensures text readability, especially on large screens. Limit the use of saturated colors to essential parts only, and avoid using more than one highly saturated accent color at the same time.

To preserve stylistic coherence across slides, apply the same color hierarchy to titles, body text, shapes, and background regions. This ensures a stable rhythm and prevents the viewer’s attention from being scattered.


### Please provide a **detailed, cohesive description of the overall visual style** for the slide template, **no more than 200 words**.
Focus **only on stylistic and aesthetic elements** — **do NOT provide any guidance on layout or content placement.**

Your style description should cover the following aspects:

### **1. Overall Style Overview**

Describe the high-level visual tone, atmosphere, and aesthetic direction of the slide deck. 
Describe the background style of the slide, including whether it uses a gradient, the proportion of each color when no gradient is applied, and any patterns or textures incorporated.
The proportion of various colors on the function slides does not need to be consistent with that on the content slides.

### **2. Decorative Graphic Motifs**

Clearly specify the main decorative shapes used in this style - detail the colors and application methods of their color schemes, whether to use only the edges or the filled parts, whether to use solid color filling, gradient filling or with transparency, the relative size range, density and quantity, typical layout tendencies (such as edge or background decorations), and any dynamic or repetitive style patterns.
The functional slides and content slides are unified.

### 3. Style Guidelines for function Slides (Cover, Ending, and Section Break Slides) (**50 words**)

Consist with the overall style, describe specific stylistic elements for function slides.

### 4. Style Guidelines for Standard Content Slides (**50 words**)

Consist with the overall style, describe specific stylistic elements for standard content slides.

"""

eval_cover="""
A. Completeness: Whether all requested design requirements are properly reflected, and whether text and visuals match the provided content.

A.1 Design Element Consistency
Assess whether the color, style, font, and other requirements are correctly applied throughout the design.

A.2 Content Fidelity
Check whether all referenced text and images appear in the provided source content without omission or fabrication. It is clearly and distinctly possible to understand the content that the page intends to convey.

B. Compliance: Whether the designed visual and text blocks follow structural rules such as non-overlap and proper spatial arrangement.

B.1 Overlap Ratio
Check whether any text or visual elements unintentionally overlap with each other, excluding decorative backgrounds or stylistic framing elements.

B.2 Page Occupancy Ratio
Evaluate whether the total occupied area of all elements is appropriate and whether the layout is visually balanced without leaning excessively toward one side.

B.3 Overflow Ratio
Identify portions of any element that extend beyond the slide boundary and quantify the exceeded area.

C. Aesthetics: Evaluate the overall visual appeal of the slide design, including color harmony, font readability, and layout balance.

C.1 Layout & Composition: Whether the spatial arrangement is balanced, alignment is consistent, and spacing between elements is appropriate.

C.2 Visual Hierarchy: Whether visual weight is properly distributed, key elements stand out, and the viewing flow feels natural.

C.3 Color & Contrast: Whether the color palette is harmonious, contrasts are sufficient, and overall color usage feels cohesive.

C.4 Typography: Whether font selection, sizes, spacing, and text layout are visually appealing and easy to read.

Output Format:
Provide your evaluation in the following JSON format:
{{
  "Total Score": float,  // Overall score out of 100
  "Breakdown": {{
    "Completeness": float,  // Score out of 30
    "Compliance": float,    // Score out of 30
    "Aesthetics": float     // Score out of 40
  }},
  "Suggestions": string     // Improvement suggestions
}}"""

design_formatting_prompt_v2 = """
# Meta & Canvas
- meta: {"version":"1.1","units":"inch"}
- Canvas: 13.33 × 7.50 (width × height). All coordinates/sizes keep **2 decimals**.
- Bounds: No element (except separators) may overflow; siblings must not overlap; children must be fully inside parents.
- Z-order: Integer, default 0 (larger = in front).
- Images: If **no usable images** are provided, do not reference images (no blank placeholders).

# Layered Structure
## Background Layer (reusable, loosely coupled but richly styled)
- Provide a ≤80-word natural-language description (style/tone).
- Provide **Machine JSON** that describes the background as a stack of layers composed of basic primitives.
- Readability constraints:
  - Prefer low-contrast, low-opacity ornaments.
  - Respect `policy.avoid_content` and `policy.safe_insets` (background elements must avoid the content safe area).

### Background Machine JSON (multi-layer + patterns/ornaments/grid/free shapes)
{
  "Tone": "light|dark|neutral",
  "Primary/Accent": "#RRGGBB / #RRGGBB",
  "Base texture": "none | <color + style>",
  "Edge ornament": "none | <color + style>",

  "policy": {
    "avoid_content": true,                       // if true, any background element overlapping content blocks is skipped
    "safe_insets": {"top":0.20,"right":0.20,"bottom":0.20,"left":0.20}, // global safe margins (inches)
    "max_alpha": 0.35                            // global max opacity (0–1) to protect readability
  },

  "layers": [
    {
      "id": "bg-fill",
      "z": 0,
      "type": "fill",                            // full-canvas fill
      "fill": {
        "mode": "solid|gradient|image",
        "color": "#RRGGBB",                      // for solid
        "gradient": { "colors":["#RRGGBB","#RRGGBB"], "angle": 0-360 },   // for gradient
        "image": { "src":"<url-or-id>", "fit":"cover|contain|stretch" } // for image
      }
    },

    {//optional
      "id": "soft-dots",
      "z": 1,
      "type": "pattern",                         // tiled primitives: dot/stripe/chevron/diamond/triangle
      "pattern": {
        "primitive": "dot|stripe|chevron|diamond|triangle",
        "tile": {"w":0.50,"h":0.50},             // size of each tile (inches)
        "repeat": {"cols":28,"rows":16},
        "origin": {"x":0.25,"y":0.25},
        "gap": {"x":0.10,"y":0.10},              // spacing between tiles
        "jitter": {"x":0.05,"y":0.05},           // subtle randomness to avoid rigidity
        "style": {
          "size": 0.06,                          // primitive size (inches) for dots/triangles/diamonds
          "weight": 0.02,                        // line weight (inches) for stripes/lines
          "color": "#RRGGBB",
          "opacity": 0.15,
          "dash": "solid|dot|dash",
          "shape": "circle|rect|rounded"         // dot shape
        }
      }
    },

    {//optional
      "id": "accent-grid",
      "z": 2,
      "type": "grid",                            // subtle grid/guides
      "grid": {
        "cols": 12, "rows": 6,
        "stroke": {"color":"#DCE4EE","weight":0.02,"dash":"dot","opacity":0.2},
        "inner_margin": {"x":0.00,"y":0.00}
      }
    },

    {//optional
      "id": "corner-ornaments",
      "z": 3,
      "type": "ornaments",                       // decorative shapes on edges/corners
      "ornaments": [
        {
          "shape": "rounded", "bbox_in":{"x":0.30,"y":0.30,"w":1.20,"h":0.40},
          "fill": {"mode":"solid","color":"#0B5FFF"},
          "opacity": 0.08, "radius": 0.20, "rotation": 0
        },
        {
          "shape": "circle", "bbox_in":{"x":11.90,"y":6.60,"w":0.50,"h":0.50},
          "fill": {"mode":"solid","color":"#FFC53D"},
          "opacity": 0.12
        }
      ]
    },

    {//optional This option is recommended for the cover, back cover and chapter pages.
      "id": "free-shapes",
      "z": 4,
      "type": "shapes",                          // arbitrary shapes/lines collection
      "elements": [
        {
          "kind": "shape", "shape":"rounded",
          "bbox_in":{"x":6.60,"y":-0.20,"w":7.00,"h":1.60},            // negative origin allowed; renderer will clip to canvas
          "fill":{"mode":"gradient","gradient":{"colors":["#RRGGBB","#RRGGBB"],"angle":0-360}},
          "opacity": 1.0, "radius":0.30, "rotation": 0
        },
        {
          "kind": "line",
          "line":{"x1":0.50,"y1":1.40,"x2":12.83,"y2":1.40},
          "stroke":{"color":"#E6EAF0","weight":0.02,"dash":"solid","opacity":0.8}
        },
        ...
      ]
    }
  ]
}

## Layout Layer (blocks & separators only; no literal text/images here)
- Summary: single / top-bottom / left-right / three-horizontal / 2x2-grid / card-grid / ...
- Natural-language details: ≤120 words about alignment/whitespace/emphasis.
- Structure JSON (same as v2, with optional `background.shape` and `rotation`).

### Layout Structure JSON
{
  "blocks": {
    "Block1": {
      "id": "Title", "function": "text",
      "bbox_in": {"x":0.50,"y":0.40,"w":6.33,"h":1.00},
      "z": 5,
      "background": { "type":"none", "shape":"rect", "radius":0.00, "opacity":1.0 },
      "border": {"on": false},
      "children": {
        "Block1.1": { "function":"text", "bbox_in":{"x":0.50,"y":0.95,"w":6.33,"h":0.45} }
      }
    },
    "Block2": {
      "id":"LeftCard","function":"mixed",
      "bbox_in":{"x":0.50,"y":1.60,"w":6.33,"h":2.60},
      "z": 4,
      "background":{"type":"solid","color":"#F5F7FB","shape":"rounded","radius":0.20,"opacity":1.0},
      "children": {}
    },
    "Block3": {
      "id":"RightCard","function":"image",
      "bbox_in":{"x":6.90,"y":1.60,"w":5.93,"h":2.60},
      "z": 4,
      "background":{"type":"solid","color":"#F5F7FB","shape":"rounded","radius":0.20}
    }
  },
  "separators": [
    {"type":"line","bbox_in":{"x":0.50,"y":1.40,"w":12.33,"h":0.00},"color":"#E6EAF0",
     "style":{"on":true,"weight":0.02,"color":"#E6EAF0","dash":"solid"}}
  ]
}

## Content Layer (bindings only; no positioning here)
- Map content to block ids (and child ids) defined in Layout. Do not invent new ids. No empty image placeholders.
- Text/image styles are optional; defaults apply if omitted.

### Content JSON
{
  "Title": {"text":"OpenRouterBench Overview", "style":{"size":36,"bold":true,"color":"#0B1F33"}},
  "LeftCard": {"text":"Datasets (25+)\nHLE · SimpleQA · Swe-bench · …", "style":{"size":18,"align":"left"}},
  "RightCard": {"image":"<url-or-path>", "style":{"fit":"contain","radius":0.20}}
}

# Output Format (strict; output ONLY these three sections)
[Meta]
  {"version":"1.1","units":"inch"}

[Background]
  Preset background: none | <filename or URL>
  Tone: light|dark|neutral
  Primary/Accent: {main_color} / {accent_color}
  Base texture: <color + style, or "none">
  Edge ornament: <color + style, or "none">
  Natural-language description: <≤80 words>
  Machine: { ... }   // Use the Background Machine JSON above (supports layers/pattern/ornaments/grid/shapes)

[Layout]
  Summary: single / top-bottom / left-right / three-horizontal / 2x2-grid / card-grid / ...
  Natural-language details: <≤120 words>
  Structure: { "blocks":{...}, "separators":[...] }

[Content]
  { ... }

# Validation Rules
- Two-decimal precision for all bbox_in numeric values.
- Canvas bounds: 0 ≤ x ≤ 13.33, 0 ≤ y ≤ 7.50, and x+w ≤ 13.33, y+h ≤ 7.50.
- Non-overlap: siblings must not overlap; children must be fully contained within the parent bbox.
- Background policy: if `avoid_content=true`, any background element intersecting content blocks is skipped; respect `safe_insets` and `max_alpha`.
- No images if none are provided. No blank placeholders.
- Render order: background.layers (by z) → separators → blocks (backgrounds) → content (text/images).
"""

color_examples_prompt = """
### **Color Style Guide for Slide Design**

Effective slide design relies on a clear and disciplined color system. A well-constructed palette generally includes one dominant color, one to two accent colors, and a set of neutral tones for balance. The dominant color establishes the visual identity of the slide, while the accent colors highlight key information such as keywords, icons, or numerical results. Neutral grays should be used for backgrounds, text blocks, and low-priority elements to prevent visual clutter.

For gradients, choose subtle transitions within the same hue family rather than mixing unrelated colors; this preserves consistency and avoids visual noise. When applying multiple colors on a single slide, maintain a contrast ratio that ensures text readability, especially on large screens. Limit the use of saturated colors to essential parts only, and avoid using more than one highly saturated accent color at the same time.

To preserve stylistic coherence across slides, apply the same color hierarchy to titles, body text, shapes, and background regions. This ensures a stable rhythm and prevents the viewer’s attention from being scattered.


COLOR PALETTE EXAMPLES:
{
   "main_color": "#D6CCC2",
   "accent_color": "#B2967D",
   "background_tone": "light",
   "heading_font_color": "#3F3F46",
   "body_font_color": "#52525B"
},
{
   "main_color": "#553C9A",
   "accent_color": "#9F7AEA",
   "background_tone": "dark",
   "heading_font_color": "#FFFFFF",
   "body_font_color": "#E9D8FD"
},
{
   "main_color": "#F6E05E",
   "accent_color": "#553C9A",
   "background_tone": "light",
   "heading_font_color": "#1A202C",
   "body_font_color": "#2D3748"
},
    """

code_prefix="""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx_tools.add_free_shape import *
"""